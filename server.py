"""
Claus Bridge MCP — Kommunikációs híd CLI-Claus és Web-Claus között
==================================================================
Railway deploy: SSE transport, SQLite persistent storage, FTS5 full-text search.
18 tool: messaging + threads + tasks + shared memory + discussions + session logs + capabilities.

Deployed on Railway alongside Makronóm, BioMed, CégTár, HírMagnet MCP servers.
"""

import os
import json
import sqlite3
import time
import base64
import asyncio
import logging
import threading
import pathlib
from datetime import datetime, timezone, timedelta
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from email.utils import parseaddr
from fastmcp import FastMCP

# Pyramid module — agentic context, memory, governance
try:
    from pyramid.context_builder import build_agent_context
    from pyramid.governance import store_result as pyramid_store_result
    from pyramid.agents import AGENT_REGISTRY as PYRAMID_AGENTS
    from pyramid.agents import temporal_directive as _temporal_directive
    PYRAMID_ENABLED = True
except ImportError:
    PYRAMID_ENABLED = False
    def _temporal_directive(agent_name: str = "") -> str:
        from datetime import datetime, timezone
        return f"\nA mai dátum: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}.\n"

# Permission layer — multi-instance access control (YoungeReka etc.)
from permissions import (
    check_permission, filter_messages, filter_memory_results,
    PermissionDeniedError, Access, is_core_instance, get_profile
)
from youngereka_profile import register_youngereka

# Feldwebel — Telegram command system + smart triage + briefing
try:
    from feldwebel import init_feldwebel, BridgeContext
    FELDWEBEL_ENABLED = True
except ImportError:
    FELDWEBEL_ENABLED = False

logger = logging.getLogger("claus-bridge")

# --- Server Setup ---
mcp = FastMCP("Claus Bridge")

DB_PATH = os.environ.get("BRIDGE_DB_PATH", "bridge.db")


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")
    return conn


def _get_inbox_summary(max_items: int = 8) -> str:
    """Legfrissebb capture-daemon üzenetek (email + calendar) a Bridge DB-ből — Pyramid kontextushoz."""
    try:
        conn = get_db()
        rows = conn.execute(
            "SELECT subject, message, timestamp, priority FROM messages "
            "WHERE sender = 'capture-daemon' ORDER BY timestamp DESC LIMIT ?",
            (max_items,)
        ).fetchall()
        conn.close()
        if not rows:
            return ""
        lines = []
        for r in rows:
            prio = {"urgent": "\U0001f534", "important": "\U0001f7e0", "normal": "\U0001f535"}.get(r["priority"], "\u26aa")
            lines.append(f"{prio} [{r['timestamp'][:16]}] {r['subject']}")
        return "\n".join(lines)
    except Exception:
        return ""


def init_db():
    conn = get_db()
    conn.executescript("""
        -- Messages with threading support
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            timestamp TEXT NOT NULL,
            sender TEXT NOT NULL,
            recipient TEXT NOT NULL,
            subject TEXT NOT NULL,
            message TEXT NOT NULL,
            priority TEXT DEFAULT 'normal',
            thread_id INTEGER,
            reply_to INTEGER,
            status TEXT DEFAULT 'unread',
            FOREIGN KEY (reply_to) REFERENCES messages(id)
        );

        -- Full-text search on messages
        CREATE VIRTUAL TABLE IF NOT EXISTS messages_fts USING fts5(
            subject, message, sender, recipient,
            content=messages, content_rowid=id
        );

        -- Triggers to keep FTS in sync
        CREATE TRIGGER IF NOT EXISTS messages_ai AFTER INSERT ON messages BEGIN
            INSERT INTO messages_fts(rowid, subject, message, sender, recipient)
            VALUES (new.id, new.subject, new.message, new.sender, new.recipient);
        END;

        -- Shared memory / knowledge base
        CREATE TABLE IF NOT EXISTS shared_memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            key TEXT NOT NULL,
            value TEXT NOT NULL,
            category TEXT DEFAULT 'general',
            tags TEXT DEFAULT '',
            updated_by TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        );

        CREATE VIRTUAL TABLE IF NOT EXISTS memory_fts USING fts5(
            key, value, category, tags,
            content=shared_memory, content_rowid=id
        );

        CREATE TRIGGER IF NOT EXISTS memory_ai AFTER INSERT ON shared_memory BEGIN
            INSERT INTO memory_fts(rowid, key, value, category, tags)
            VALUES (new.id, new.key, new.value, new.category, new.tags);
        END;

        CREATE TRIGGER IF NOT EXISTS memory_au AFTER UPDATE ON shared_memory BEGIN
            DELETE FROM memory_fts WHERE rowid = old.id;
            INSERT INTO memory_fts(rowid, key, value, category, tags)
            VALUES (new.id, new.key, new.value, new.category, new.tags);
        END;

        -- Tasks
        CREATE TABLE IF NOT EXISTS tasks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            description TEXT DEFAULT '',
            assigned_to TEXT NOT NULL,
            assigned_by TEXT NOT NULL,
            priority TEXT DEFAULT 'normal',
            status TEXT DEFAULT 'pending',
            deadline TEXT,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        );

        -- Discussions (collaborative thinking)
        CREATE TABLE IF NOT EXISTS discussions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            topic TEXT NOT NULL,
            context TEXT DEFAULT '',
            status TEXT DEFAULT 'open',
            resolution TEXT,
            created_by TEXT NOT NULL,
            created_at TEXT NOT NULL,
            resolved_at TEXT
        );

        CREATE TABLE IF NOT EXISTS discussion_entries (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            discussion_id INTEGER NOT NULL,
            instance TEXT NOT NULL,
            content TEXT NOT NULL,
            timestamp TEXT NOT NULL,
            FOREIGN KEY (discussion_id) REFERENCES discussions(id)
        );

        CREATE VIRTUAL TABLE IF NOT EXISTS discussions_fts USING fts5(
            topic, context,
            content=discussions, content_rowid=id
        );

        CREATE TRIGGER IF NOT EXISTS discussions_ai AFTER INSERT ON discussions BEGIN
            INSERT INTO discussions_fts(rowid, topic, context)
            VALUES (new.id, new.topic, new.context);
        END;

        -- Session logs
        CREATE TABLE IF NOT EXISTS session_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            instance TEXT NOT NULL,
            summary TEXT NOT NULL,
            key_decisions TEXT DEFAULT '',
            key_learnings TEXT DEFAULT '',
            timestamp TEXT NOT NULL
        );

        -- Capability registry
        CREATE TABLE IF NOT EXISTS capabilities (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            instance TEXT NOT NULL,
            tool_name TEXT NOT NULL,
            description TEXT NOT NULL,
            registered_at TEXT NOT NULL,
            UNIQUE(instance, tool_name)
        );

        -- Heartbeat tracking
        CREATE TABLE IF NOT EXISTS heartbeats (
            instance TEXT PRIMARY KEY,
            last_seen TEXT NOT NULL,
            session_info TEXT DEFAULT ''
        );

        -- Uploaded files for AI processing
        CREATE TABLE IF NOT EXISTS uploads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT NOT NULL,
            mime_type TEXT NOT NULL,
            content_text TEXT DEFAULT '',
            content_base64 TEXT DEFAULT '',
            uploaded_by TEXT NOT NULL,
            uploaded_at TEXT NOT NULL
        );

        -- AI Tasks (multi-agent task execution)
        CREATE TABLE IF NOT EXISTS ai_tasks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            description TEXT NOT NULL,
            context TEXT DEFAULT '',
            assigned_by TEXT NOT NULL,
            status TEXT DEFAULT 'pending',
            created_at TEXT NOT NULL,
            completed_at TEXT
        );

        CREATE TABLE IF NOT EXISTS ai_task_results (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            task_id INTEGER NOT NULL,
            agent TEXT NOT NULL,
            role TEXT DEFAULT '',
            content TEXT NOT NULL,
            sources TEXT DEFAULT '',
            timestamp TEXT NOT NULL,
            FOREIGN KEY (task_id) REFERENCES ai_tasks(id)
        );
    """)
    conn.commit()
    conn.close()


def now():
    return datetime.now(timezone.utc).isoformat()


# ============================================================
# PERMISSION HELPERS
# ============================================================

def _enforce(caller: str, tool_name: str, **kwargs) -> str | None:
    """Check permission. Returns error JSON if denied, None if allowed."""
    if not caller or is_core_instance(caller):
        return None
    try:
        access = check_permission(caller, tool_name, **kwargs)
        return None
    except PermissionDeniedError as e:
        return json.dumps({"error": str(e), "status": "denied"})


# ============================================================
# MESSAGING TOOLS (1-5)
# ============================================================

@mcp.tool()
async def send_message(sender: str, recipient: str, subject: str, message: str,
                       priority: str = "normal", reply_to: int = None) -> str:
    """Send a message to the other Claus instance.

    Args:
        sender: Who sends it — 'cli-claus', 'web-claus', or instance ID
        recipient: Who receives it — 'cli-claus', 'web-claus', or instance ID
        subject: Message subject line
        message: Full message content
        priority: info / normal / urgent / critical
        reply_to: Optional message ID to reply to (creates thread)
    """
    denied = _enforce(sender, "send_message", recipient=recipient)
    if denied:
        return denied
    conn = get_db()
    thread_id = None
    if reply_to:
        row = conn.execute("SELECT thread_id, id FROM messages WHERE id = ?", (reply_to,)).fetchone()
        if row:
            thread_id = row["thread_id"] or row["id"]

    cur = conn.execute(
        "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority, thread_id, reply_to) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (now(), sender, recipient, subject, message, priority, thread_id, reply_to)
    )
    msg_id = cur.lastrowid
    if not thread_id:
        conn.execute("UPDATE messages SET thread_id = ? WHERE id = ?", (msg_id, msg_id))
    conn.commit()
    conn.close()
    return json.dumps({"status": "sent", "message_id": msg_id, "thread_id": thread_id or msg_id})


@mcp.tool()
async def read_messages(recipient: str = None, limit: int = 20, since: str = None,
                        unread_only: bool = False, thread_id: int = None,
                        caller: str = "") -> str:
    """Read messages, optionally filtered.

    Args:
        recipient: Filter by recipient ('cli-claus' or 'web-claus')
        limit: Max messages to return (default 20)
        since: ISO timestamp — only messages after this time
        unread_only: If true, only return unread messages
        thread_id: Filter by thread ID to see a conversation
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "read_messages")
    if denied:
        return denied
    conn = get_db()
    query = "SELECT * FROM messages WHERE 1=1"
    params = []

    if recipient:
        query += " AND recipient = ?"
        params.append(recipient)
    if since:
        query += " AND timestamp > ?"
        params.append(since)
    if unread_only:
        query += " AND status = 'unread'"
    if thread_id:
        query += " AND thread_id = ?"
        params.append(thread_id)

    query += " ORDER BY timestamp DESC LIMIT ?"
    params.append(limit)

    rows = conn.execute(query, params).fetchall()
    messages = [dict(r) for r in rows]
    conn.close()

    # Apply message filtering for non-core instances
    if caller and not is_core_instance(caller):
        messages = filter_messages(caller, messages)

    return json.dumps(messages, ensure_ascii=False)


@mcp.tool()
async def read_new(instance: str) -> str:
    """Read all unread messages for a specific instance and mark them as read.

    Args:
        instance: Instance ID (e.g. 'cli-claus', 'web-claus', 'YoungeReka')
    """
    denied = _enforce(instance, "read_new")
    if denied:
        return denied
    conn = get_db()
    rows = conn.execute(
        "SELECT * FROM messages WHERE recipient = ? AND status = 'unread' ORDER BY timestamp ASC",
        (instance,)
    ).fetchall()
    messages = [dict(r) for r in rows]

    if messages:
        ids = [m["id"] for m in messages]
        conn.execute(
            f"UPDATE messages SET status = 'read' WHERE id IN ({','.join('?' * len(ids))})",
            ids
        )
        conn.commit()

    conn.close()

    # Apply message filtering for non-core instances
    if instance and not is_core_instance(instance):
        messages = filter_messages(instance, messages)

    return json.dumps({"count": len(messages), "messages": messages}, ensure_ascii=False)


@mcp.tool()
async def search_messages(query: str, limit: int = 20, caller: str = "") -> str:
    """Full-text search across all messages.

    Args:
        query: Search query (supports FTS5 syntax: AND, OR, NOT, "exact phrase")
        limit: Max results (default 20)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "search_messages")
    if denied:
        return denied
    conn = get_db()
    rows = conn.execute(
        "SELECT m.* FROM messages m JOIN messages_fts f ON m.id = f.rowid "
        "WHERE messages_fts MATCH ? ORDER BY rank LIMIT ?",
        (query, limit)
    ).fetchall()
    messages = [dict(r) for r in rows]
    conn.close()

    if caller and not is_core_instance(caller):
        messages = filter_messages(caller, messages)

    return json.dumps(messages, ensure_ascii=False)


@mcp.tool()
async def mark_read(message_id: int, caller: str = "") -> str:
    """Mark a specific message as read.

    Args:
        message_id: The message ID to mark as read
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "mark_read")
    if denied:
        return denied
    conn = get_db()
    conn.execute("UPDATE messages SET status = 'read' WHERE id = ?", (message_id,))
    conn.commit()
    conn.close()
    return json.dumps({"status": "marked_read", "message_id": message_id})


# ============================================================
# SHARED MEMORY / KNOWLEDGE BASE (6-9)
# ============================================================

@mcp.tool()
async def write_memory(key: str, value: str, category: str = "general",
                       tags: str = "", instance: str = "unknown") -> str:
    """Write or update a shared memory entry. Use for decisions, project context, learnings.

    Args:
        key: Unique key (e.g., 'bridge_mcp_auth_decision')
        value: Content to store
        category: general / decision / project / learning / config
        tags: Comma-separated tags for searchability
        instance: Who wrote it
    """
    denied = _enforce(instance, "write_memory")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    existing = conn.execute("SELECT id FROM shared_memory WHERE key = ?", (key,)).fetchone()

    if existing:
        conn.execute(
            "UPDATE shared_memory SET value=?, category=?, tags=?, updated_by=?, updated_at=? WHERE key=?",
            (value, category, tags, instance, ts, key)
        )
    else:
        conn.execute(
            "INSERT INTO shared_memory (key, value, category, tags, updated_by, created_at, updated_at) "
            "VALUES (?, ?, ?, ?, ?, ?, ?)",
            (key, value, category, tags, instance, ts, ts)
        )
    conn.commit()
    conn.close()
    return json.dumps({"status": "saved", "key": key})


@mcp.tool()
async def read_memory(key: str = None, category: str = None, caller: str = "") -> str:
    """Read shared memory entries by key or category.

    Args:
        key: Exact key to look up (returns single entry)
        category: Filter by category (returns all matching)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "read_memory")
    if denied:
        return denied
    conn = get_db()
    if key:
        row = conn.execute("SELECT * FROM shared_memory WHERE key = ?", (key,)).fetchone()
        conn.close()
        return json.dumps(dict(row) if row else {"error": f"Key '{key}' not found"}, ensure_ascii=False)
    elif category:
        rows = conn.execute("SELECT * FROM shared_memory WHERE category = ? ORDER BY updated_at DESC", (category,)).fetchall()
        conn.close()
        return json.dumps([dict(r) for r in rows], ensure_ascii=False)
    else:
        rows = conn.execute("SELECT key, category, updated_by, updated_at FROM shared_memory ORDER BY updated_at DESC LIMIT 50").fetchall()
        conn.close()
        return json.dumps([dict(r) for r in rows], ensure_ascii=False)


@mcp.tool()
async def search_memory(query: str, limit: int = 20, caller: str = "") -> str:
    """Full-text search across shared memory.

    Args:
        query: Search query (FTS5 syntax)
        limit: Max results
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "search_memory")
    if denied:
        return denied
    conn = get_db()
    rows = conn.execute(
        "SELECT m.* FROM shared_memory m JOIN memory_fts f ON m.id = f.rowid "
        "WHERE memory_fts MATCH ? ORDER BY rank LIMIT ?",
        (query, limit)
    ).fetchall()
    conn.close()
    return json.dumps([dict(r) for r in rows], ensure_ascii=False)


@mcp.tool()
async def list_memory(category: str = None, caller: str = "") -> str:
    """List all shared memory keys, optionally filtered by category.

    Args:
        category: Optional filter — general / decision / project / learning / config
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "list_memory")
    if denied:
        return denied
    conn = get_db()
    if category:
        rows = conn.execute(
            "SELECT key, category, tags, updated_by, updated_at FROM shared_memory WHERE category = ? ORDER BY updated_at DESC",
            (category,)
        ).fetchall()
    else:
        rows = conn.execute(
            "SELECT key, category, tags, updated_by, updated_at FROM shared_memory ORDER BY category, updated_at DESC"
        ).fetchall()
    conn.close()
    return json.dumps([dict(r) for r in rows], ensure_ascii=False)


# ============================================================
# TASK MANAGEMENT (10-12)
# ============================================================

@mcp.tool()
async def create_task(title: str, assigned_to: str, assigned_by: str,
                      description: str = "", priority: str = "normal",
                      deadline: str = None) -> str:
    """Create a task for one of the Claus instances.

    Args:
        title: Task title
        assigned_to: Assignee instance ID
        assigned_by: Who created it
        description: Detailed description
        priority: low / normal / high / critical
        deadline: Optional ISO date
    """
    denied = _enforce(assigned_by, "create_task")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO tasks (title, description, assigned_to, assigned_by, priority, status, deadline, created_at, updated_at) "
        "VALUES (?, ?, ?, ?, ?, 'pending', ?, ?, ?)",
        (title, description, assigned_to, assigned_by, priority, deadline, ts, ts)
    )
    conn.commit()
    task_id = cur.lastrowid
    conn.close()
    return json.dumps({"status": "created", "task_id": task_id})


@mcp.tool()
async def update_task(task_id: int, status: str = None, description: str = None,
                      caller: str = "") -> str:
    """Update a task's status or description.

    Args:
        task_id: Task ID
        status: pending / in_progress / completed / cancelled
        description: Updated description (append notes)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "update_task")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    if status:
        conn.execute("UPDATE tasks SET status=?, updated_at=? WHERE id=?", (status, ts, task_id))
    if description:
        conn.execute("UPDATE tasks SET description=?, updated_at=? WHERE id=?", (description, ts, task_id))
    conn.commit()
    conn.close()
    return json.dumps({"status": "updated", "task_id": task_id})


@mcp.tool()
async def list_tasks(assigned_to: str = None, status: str = None, caller: str = "") -> str:
    """List tasks, optionally filtered.

    Args:
        assigned_to: Filter by assignee
        status: Filter by status (pending/in_progress/completed/cancelled)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "list_tasks")
    if denied:
        return denied
    conn = get_db()
    query = "SELECT * FROM tasks WHERE 1=1"
    params = []
    if assigned_to:
        query += " AND assigned_to = ?"
        params.append(assigned_to)
    if status:
        query += " AND status = ?"
        params.append(status)
    query += " ORDER BY CASE priority WHEN 'critical' THEN 0 WHEN 'high' THEN 1 WHEN 'normal' THEN 2 ELSE 3 END, created_at DESC"
    rows = conn.execute(query, params).fetchall()
    conn.close()
    return json.dumps([dict(r) for r in rows], ensure_ascii=False)


# ============================================================
# DISCUSSIONS (13-17)
# ============================================================

@mcp.tool()
async def start_discussion(topic: str, initial_position: str, context: str = "",
                           instance: str = "unknown") -> str:
    """Start a collaborative discussion on a topic. Both instances can add their thoughts.

    Args:
        topic: What are we discussing? (e.g., 'Bridge MCP auth megoldás')
        initial_position: Your opening position/argument
        context: Background context for the discussion
        instance: Who starts it
    """
    denied = _enforce(instance, "start_discussion")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO discussions (topic, context, status, created_by, created_at) VALUES (?, ?, 'open', ?, ?)",
        (topic, context, instance, ts)
    )
    disc_id = cur.lastrowid
    conn.execute(
        "INSERT INTO discussion_entries (discussion_id, instance, content, timestamp) VALUES (?, ?, ?, ?)",
        (disc_id, instance, initial_position, ts)
    )
    conn.commit()
    conn.close()

    # Auto-trigger AI sub-agents
    thread_text = f"[{instance}]: {initial_position}"
    asyncio.ensure_future(_ai_auto_discuss(disc_id, topic, thread_text))

    return json.dumps({"status": "discussion_started", "discussion_id": disc_id})


@mcp.tool()
async def add_to_discussion(discussion_id: int, content: str, instance: str = "unknown") -> str:
    """Add your thoughts/response to an existing discussion.

    Args:
        discussion_id: Discussion ID to contribute to
        content: Your thoughts, arguments, counterpoints
        instance: Instance ID
    """
    denied = _enforce(instance, "add_to_discussion")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    conn.execute(
        "INSERT INTO discussion_entries (discussion_id, instance, content, timestamp) VALUES (?, ?, ?, ?)",
        (discussion_id, instance, content, ts)
    )
    conn.commit()

    entries = conn.execute(
        "SELECT instance, content, timestamp FROM discussion_entries WHERE discussion_id = ? ORDER BY timestamp",
        (discussion_id,)
    ).fetchall()

    # Get topic for AI context
    disc = conn.execute("SELECT topic FROM discussions WHERE id = ?", (discussion_id,)).fetchone()
    conn.close()

    result = json.dumps({
        "status": "added",
        "discussion_id": discussion_id,
        "total_entries": len(entries),
        "thread": [dict(e) for e in entries]
    }, ensure_ascii=False)

    # Auto-trigger AI sub-agents (only when Claude or Kommandant adds, not when AIs add)
    if instance not in SILICONFLOW_MODELS:
        thread_text = "\n".join(f"[{e['instance']}]: {e['content']}" for e in entries)
        topic = disc["topic"] if disc else f"Discussion #{discussion_id}"
        asyncio.ensure_future(_ai_auto_discuss(discussion_id, topic, thread_text))

    return result


@mcp.tool()
async def resolve_discussion(discussion_id: int, resolution: str, instance: str = "unknown") -> str:
    """Resolve a discussion and save the decision to shared memory.

    Args:
        discussion_id: Discussion to resolve
        resolution: The agreed decision/conclusion
        instance: Who resolves it
    """
    denied = _enforce(instance, "resolve_discussion")
    if denied:
        return denied
    conn = get_db()
    ts = now()

    disc = conn.execute("SELECT topic, context FROM discussions WHERE id = ?", (discussion_id,)).fetchone()
    if not disc:
        conn.close()
        return json.dumps({"error": f"Discussion {discussion_id} not found"})

    conn.execute(
        "UPDATE discussions SET status='resolved', resolution=?, resolved_at=? WHERE id=?",
        (resolution, ts, discussion_id)
    )

    # Auto-save decision to shared memory
    key = f"decision_{discussion_id}_{disc['topic'][:50].replace(' ', '_').lower()}"
    conn.execute(
        "INSERT INTO shared_memory (key, value, category, tags, updated_by, created_at, updated_at) "
        "VALUES (?, ?, 'decision', ?, ?, ?, ?)",
        (key, resolution, f"discussion_{discussion_id}", instance, ts, ts)
    )
    conn.commit()
    conn.close()
    return json.dumps({"status": "resolved", "discussion_id": discussion_id, "memory_key": key})


@mcp.tool()
async def read_discussion(discussion_id: int, caller: str = "") -> str:
    """Read all entries in a discussion without adding a new entry.

    Args:
        discussion_id: Discussion ID to read
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "read_discussion")
    if denied:
        return denied
    conn = get_db()
    disc = conn.execute("SELECT * FROM discussions WHERE id = ?", (discussion_id,)).fetchone()
    if not disc:
        conn.close()
        return json.dumps({"error": f"Discussion {discussion_id} not found"})

    entries = conn.execute(
        "SELECT instance, content, timestamp FROM discussion_entries WHERE discussion_id = ? ORDER BY timestamp",
        (discussion_id,)
    ).fetchall()
    conn.close()

    return json.dumps({
        "discussion_id": discussion_id,
        "topic": disc["topic"],
        "context": disc["context"],
        "status": disc["status"],
        "resolution": disc["resolution"],
        "entry_count": len(entries),
        "entries": [dict(e) for e in entries],
    }, ensure_ascii=False)


@mcp.tool()
async def read_ai_task_results(task_id: int = 0, limit: int = 10, caller: str = "") -> str:
    """Read AI task results. If task_id given, returns that task's agent outputs. Otherwise lists recent tasks.

    Args:
        task_id: Specific task ID to read results for (0 = list recent tasks)
        limit: Max tasks to list when task_id is 0
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "read_ai_task_results")
    if denied:
        return denied
    conn = get_db()

    if task_id:
        task = conn.execute("SELECT * FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
        if not task:
            conn.close()
            return json.dumps({"error": f"AI task #{task_id} not found"})
        results = conn.execute(
            "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id",
            (task_id,)
        ).fetchall()
        conn.close()
        return json.dumps({
            "task_id": task_id,
            "title": task["title"],
            "description": task["description"],
            "status": task["status"],
            "created_at": task["created_at"],
            "completed_at": task["completed_at"],
            "results": [dict(r) for r in results],
        }, ensure_ascii=False)

    tasks = conn.execute("SELECT * FROM ai_tasks ORDER BY id DESC LIMIT ?", (limit,)).fetchall()
    conn.close()
    return json.dumps([{
        "task_id": t["id"],
        "title": t["title"],
        "status": t["status"],
        "created_at": t["created_at"],
    } for t in tasks], ensure_ascii=False)


@mcp.tool()
async def analyze_image(image_base64: str, mime_type: str = "image/jpeg",
                        prompt: str = "Mit latsz a kepen? Ird le reszletesen, magyarul.",
                        caller: str = "") -> str:
    """Analyze an image using Kimi K2.6 vision model. Works with any image source.

    Args:
        image_base64: Base64-encoded image data
        mime_type: Image MIME type (image/jpeg, image/png, image/webp)
        prompt: What to analyze / question about the image
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "analyze_image")
    if denied:
        return denied
    result = await _analyze_image(image_base64, mime_type, prompt)
    return json.dumps({"status": "analyzed", "analysis": result}, ensure_ascii=False)


@mcp.tool()
async def export_ai_task(task_id: int, format: str = "xlsx", caller: str = "") -> str:
    """Export AI task results as xlsx (spreadsheet) or pptx (presentation). Returns base64-encoded file.

    Args:
        task_id: AI task ID to export
        format: 'xlsx' for Excel spreadsheet or 'pptx' for PowerPoint presentation
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "export_ai_task")
    if denied:
        return denied
    if format not in ("xlsx", "pptx"):
        return json.dumps({"error": "format must be 'xlsx' or 'pptx'"})

    conn = get_db()
    task = conn.execute("SELECT * FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
    if not task:
        conn.close()
        return json.dumps({"error": f"AI task #{task_id} not found"})
    results = conn.execute(
        "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id",
        (task_id,)
    ).fetchall()
    conn.close()

    if not results:
        return json.dumps({"error": f"AI task #{task_id} has no results yet"})

    import base64
    from io import BytesIO

    try:
        if format == "xlsx":
            buf = _generate_xlsx(task, results)
            mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        else:
            buf = _generate_pptx(task, results)
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        filename = f"claus_ai_task_{task_id}.{format}"
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")

        return json.dumps({
            "status": "exported",
            "task_id": task_id,
            "format": format,
            "filename": filename,
            "mime_type": mime,
            "size_bytes": len(buf.getvalue()),
            "content_base64": b64,
        })
    except Exception as e:
        return json.dumps({"error": f"Export failed: {e}"})


def _generate_xlsx(task, results) -> "BytesIO":
    """Generate xlsx from AI task results."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from io import BytesIO
    import re

    wb = Workbook()
    ws = wb.active
    ws.title = "Összefoglaló"

    title_font = Font(name="Calibri", size=14, bold=True, color="1F4E79")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    ws["A1"] = task["title"]
    ws["A1"].font = title_font
    ws.merge_cells("A1:D1")
    ws["A2"] = f'Feladat: {task["description"][:200]}'
    ws["A3"] = f'Kiadta: {task["assigned_by"]} | Dátum: {task["created_at"][:10]}'

    row = 5
    agent_names = {"kimi": "Kimi-K2.6", "deepseek": "DeepSeek V3.2", "glm5": "GLM-5.1", "szintézis": "Szintézis"}
    for col, header_text in enumerate(["Agent", "Tartalom", "Időpont"], 1):
        cell = ws.cell(row=row, column=col, value=header_text)
        cell.font = header_font
        cell.fill = header_fill
        cell.border = thin_border

    for r in results:
        row += 1
        ws.cell(row=row, column=1, value=agent_names.get(r["agent"], r["agent"])).border = thin_border
        c = ws.cell(row=row, column=2, value=(r["content"] or "")[:32000])
        c.border = thin_border
        c.alignment = Alignment(wrap_text=True)
        ws.cell(row=row, column=3, value=r["timestamp"][:19]).border = thin_border

    ws.column_dimensions["A"].width = 20
    ws.column_dimensions["B"].width = 100
    ws.column_dimensions["C"].width = 20

    # Sheet 2: extracted data
    ws2 = wb.create_sheet("Adatok")
    ws2["A1"] = "Kinyert adatok"
    ws2["A1"].font = title_font
    ws2.merge_cells("A1:C1")

    data_row = 3
    for col, ht in enumerate(["Mutató", "Érték", "Egység"], 1):
        cell = ws2.cell(row=data_row, column=col, value=ht)
        cell.font = header_font
        cell.fill = header_fill
        cell.border = thin_border

    all_content = "\n".join(r["content"] or "" for r in results)
    pattern = re.compile(r'\*{0,2}([A-Za-zÀ-ž/\s\-\.]+?)\*{0,2}[:\-—]\s*([\d\s]+[\.,]?\d*)\s*(%|USD|EUR|Ft|HUF|pont|bázispont|USD/barrel|USD/oz)?', re.UNICODE)
    seen = set()
    for match in pattern.finditer(all_content):
        label = match.group(1).strip().strip("*").strip()
        value_str = match.group(2).strip().replace(" ", "")
        unit = match.group(3) or ""
        if len(label) < 3 or len(label) > 60 or label in seen:
            continue
        seen.add(label)
        data_row += 1
        ws2.cell(row=data_row, column=1, value=label).border = thin_border
        try:
            ws2.cell(row=data_row, column=2, value=float(value_str.replace(",", "."))).border = thin_border
        except ValueError:
            ws2.cell(row=data_row, column=2, value=value_str).border = thin_border
        ws2.cell(row=data_row, column=3, value=unit).border = thin_border

    ws2.column_dimensions["A"].width = 35
    ws2.column_dimensions["B"].width = 20
    ws2.column_dimensions["C"].width = 15

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _parse_slides_from_markdown(content: str, task_title: str) -> list:
    """Parse well-structured markdown into slide sections. No AI needed."""
    import re

    slides = []
    current_title = ""
    current_bullets = []

    # Extract a human-readable title from content (first # heading or task title)
    display_title = task_title
    first_h1 = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
    if first_h1:
        display_title = first_h1.group(1).strip()

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue

        # Detect section headings: ## 1. BELPOLITIKA, ## KÜLPOLITIKA, **HEADING**, etc.
        is_section = False
        section_title = ""

        if re.match(r'^#{2,3}\s', stripped):
            is_section = True
            section_title = re.sub(r'^#+\s*', '', stripped).strip()
            # Remove leading number: "1. BELPOLITIKA" → "BELPOLITIKA"
            section_title = re.sub(r'^\d+[\.\)]\s*', '', section_title).strip()
        elif stripped.startswith("**") and stripped.endswith("**") and len(stripped) < 80:
            is_section = True
            section_title = stripped.strip("*").strip()

        if is_section:
            # Save previous section as a slide
            if current_title and current_bullets:
                slides.append({"title": current_title, "bullets": current_bullets[:6]})
            current_title = section_title
            current_bullets = []
            continue

        # Skip top-level # headings (used for display_title)
        if re.match(r'^#\s', stripped):
            continue

        # Everything else is a bullet
        # Clean: remove numbering, markdown bold, leading dashes
        bullet = re.sub(r'^\d+[\.\)]\s*', '', stripped)
        bullet = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', bullet)
        bullet = re.sub(r'^[-•]\s*', '', bullet).strip()

        if bullet and len(bullet) > 3:
            # Truncate very long bullets
            if len(bullet) > 130:
                bullet = bullet[:127] + "..."
            current_bullets.append(bullet)

    # Don't forget the last section
    if current_title and current_bullets:
        slides.append({"title": current_title, "bullets": current_bullets[:6]})

    return slides, display_title


def _generate_pptx(task, results) -> "BytesIO":
    """Generate pptx — DeepSeek structures slides, then render."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from io import BytesIO
    import re

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1F, 0x2B, 0x3D)
    GOLD = RGBColor(0xD4, 0xA5, 0x37)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)

    def _set_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return tf

    # Find synthesis or single agent result
    main_content = None
    for r in results:
        if r["agent"] in ("szintézis", "szintezis"):
            main_content = r["content"] or ""
    if not main_content and results:
        main_content = results[-1]["content"] or ""

    # Parse slides from markdown structure
    slides_data, display_title = _parse_slides_from_markdown(main_content or "", task["title"])

    # Title slide — use extracted display title, not "Recipe: xyz"
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, DARK_BG)
    _text(s, 1, 1.5, 11, 1.5, display_title, size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    _text(s, 1, 3.5, 11, 0.8, f'{task["created_at"][:10]}', size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _text(s, 1, 5.5, 11, 0.6, "Claus Multi-Agent Rendszer", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    if slides_data:
        # Structured slides from markdown
        for sd in slides_data:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s, DARK_BG)
            slide_title = sd.get("title", "")
            if slide_title:
                _text(s, 0.8, 0.3, 11.5, 0.9, slide_title, size=26, color=GOLD, bold=True)
            bullets = sd.get("bullets", [])[:6]
            if bullets:
                tf = _text(s, 0.8, 1.4, 11.5, 5.5, "", size=16, color=WHITE)
                tf.paragraphs[0].clear()
                for i, bullet in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = WHITE
                    p.space_after = Pt(8)

    # Credits slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, DARK_BG)
    agents_used = list(set(r["agent"] for r in results if r["agent"] not in ("szintézis", "szintezis")))
    _text(s, 1, 2.5, 11, 1, "Claus Multi-Agent Rendszer", size=32, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    _text(s, 1, 4, 11, 0.6, f"Agentek: {', '.join(agents_used) or 'N/A'}", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _text(s, 1, 4.8, 11, 0.6, f"Task #{task['id']} | {task['created_at'][:10]}", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@mcp.tool()
async def list_discussions(status: str = None, limit: int = 20, caller: str = "") -> str:
    """List discussions, optionally filtered by status.

    Args:
        status: open / resolved (default: all)
        limit: Max results
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "list_discussions")
    if denied:
        return denied
    conn = get_db()
    query = "SELECT d.*, COUNT(e.id) as entry_count FROM discussions d LEFT JOIN discussion_entries e ON d.id = e.discussion_id"
    params = []
    if status:
        query += " WHERE d.status = ?"
        params.append(status)
    query += " GROUP BY d.id ORDER BY d.created_at DESC LIMIT ?"
    params.append(limit)
    rows = conn.execute(query, params).fetchall()
    conn.close()
    return json.dumps([dict(r) for r in rows], ensure_ascii=False)


@mcp.tool()
async def search_discussions(query: str, limit: int = 20, caller: str = "") -> str:
    """Full-text search across discussion topics and context.

    Args:
        query: Search query (FTS5 syntax)
        limit: Max results
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "search_discussions")
    if denied:
        return denied
    conn = get_db()
    rows = conn.execute(
        "SELECT d.*, COUNT(e.id) as entry_count FROM discussions d "
        "JOIN discussions_fts f ON d.id = f.rowid "
        "LEFT JOIN discussion_entries e ON d.id = e.discussion_id "
        "WHERE discussions_fts MATCH ? GROUP BY d.id ORDER BY rank LIMIT ?",
        (query, limit)
    ).fetchall()
    conn.close()
    return json.dumps([dict(r) for r in rows], ensure_ascii=False)


# ============================================================
# SESSION LOGS & STATUS (18-20)
# ============================================================

@mcp.tool()
async def log_session(instance: str, summary: str, key_decisions: str = "",
                      key_learnings: str = "") -> str:
    """Log a session summary. Call at the end of every session to maintain continuity.

    Args:
        instance: Instance ID
        summary: What happened this session
        key_decisions: Important decisions made
        key_learnings: Things learned / gotchas discovered
    """
    denied = _enforce(instance, "log_session")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO session_logs (instance, summary, key_decisions, key_learnings, timestamp) VALUES (?, ?, ?, ?, ?)",
        (instance, summary, key_decisions, key_learnings, ts)
    )
    conn.commit()
    conn.close()
    return json.dumps({"status": "logged", "log_id": cur.lastrowid})


@mcp.tool()
async def heartbeat(instance: str, session_info: str = "") -> str:
    """Send a heartbeat to signal this instance is alive. Call at session start.

    Args:
        instance: Instance ID
        session_info: Optional context about current session
    """
    denied = _enforce(instance, "heartbeat")
    if denied:
        return denied
    conn = get_db()
    ts = now()
    conn.execute(
        "INSERT OR REPLACE INTO heartbeats (instance, last_seen, session_info) VALUES (?, ?, ?)",
        (instance, ts, session_info)
    )
    conn.commit()
    conn.close()
    return json.dumps({"status": "alive", "instance": instance, "timestamp": ts})


@mcp.tool()
async def get_status(caller: str = "") -> str:
    """Get system status: who's online, unread counts, open tasks, active discussions."""
    denied = _enforce(caller, "get_status")
    if denied:
        return denied
    conn = get_db()

    heartbeats_rows = conn.execute("SELECT * FROM heartbeats").fetchall()
    unread_cli = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='cli-claus' AND status='unread'").fetchone()["c"]
    unread_web = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='web-claus' AND status='unread'").fetchone()["c"]
    unread_kmd = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='kommandant' AND status='unread'").fetchone()["c"]
    open_tasks = conn.execute("SELECT COUNT(*) as c FROM tasks WHERE status IN ('pending', 'in_progress')").fetchone()["c"]
    open_discussions = conn.execute("SELECT COUNT(*) as c FROM discussions WHERE status='open'").fetchone()["c"]
    total_messages = conn.execute("SELECT COUNT(*) as c FROM messages").fetchone()["c"]
    total_memory = conn.execute("SELECT COUNT(*) as c FROM shared_memory").fetchone()["c"]

    last_session = conn.execute("SELECT instance, summary, timestamp FROM session_logs ORDER BY timestamp DESC LIMIT 2").fetchall()

    conn.close()
    return json.dumps({
        "instances": {h["instance"]: {"last_seen": h["last_seen"], "session_info": h["session_info"]} for h in heartbeats_rows},
        "unread": {"cli-claus": unread_cli, "web-claus": unread_web, "kommandant": unread_kmd},
        "open_tasks": open_tasks,
        "open_discussions": open_discussions,
        "total_messages": total_messages,
        "total_memory_entries": total_memory,
        "recent_sessions": [dict(s) for s in last_session]
    }, ensure_ascii=False)


# ============================================================
# LANDING PAGE
# ============================================================

from starlette.responses import HTMLResponse, JSONResponse
from starlette.requests import Request
import pathlib

@mcp.custom_route("/", methods=["GET"])
async def landing(request):
    html = """<!DOCTYPE html>
<html><head><title>Claus Bridge MCP</title>
<style>
body { background: #0f172a; color: #e2e8f0; font-family: 'Inter', -apple-system, sans-serif; margin: 0; padding: 2rem; }
h1 { color: #60a5fa; } h2 { color: #94a3b8; } code { background: #1e293b; padding: 2px 6px; border-radius: 4px; }
.tool { background: #1e293b; padding: 1rem; margin: 0.5rem 0; border-radius: 8px; border-left: 3px solid #60a5fa; }
.tool b { color: #60a5fa; }
</style></head><body>
<h1>Claus Bridge MCP</h1>
<p>Kommunikációs híd <b>CLI-Claus</b> (Claude Code) és <b>Web-Claus</b> (Claude.ai) között.</p>
<p><b>Connect:</b> <code>{url}/mcp</code></p>

<h2>Messaging (5 tools)</h2>
<div class="tool"><b>send_message</b> — Üzenet küldés a másik instance-nak</div>
<div class="tool"><b>read_messages</b> — Üzenetek olvasása (szűrőkkel)</div>
<div class="tool"><b>read_new</b> — Olvasatlan üzenetek + auto mark-read</div>
<div class="tool"><b>search_messages</b> — FTS5 full-text keresés üzenetekben</div>
<div class="tool"><b>mark_read</b> — Üzenet olvasottnak jelölése</div>

<h2>Shared Memory (4 tools)</h2>
<div class="tool"><b>write_memory</b> — Közös tudásbázis írás (key-value + tags)</div>
<div class="tool"><b>read_memory</b> — Memória olvasás key/category alapján</div>
<div class="tool"><b>search_memory</b> — FTS5 keresés a tudásbázisban</div>
<div class="tool"><b>list_memory</b> — Összes memória kulcs listázás</div>

<h2>Tasks (3 tools)</h2>
<div class="tool"><b>create_task</b> — Feladat létrehozás (egymásnak adhatók)</div>
<div class="tool"><b>update_task</b> — Feladat státusz frissítés</div>
<div class="tool"><b>list_tasks</b> — Feladatok listázás (szűrőkkel)</div>

<h2>Discussions (5 tools)</h2>
<div class="tool"><b>start_discussion</b> — Közös vita indítás egy témáról</div>
<div class="tool"><b>add_to_discussion</b> — Hozzászólás vitához</div>
<div class="tool"><b>resolve_discussion</b> — Vita lezárás → döntés shared memory-ba</div>
<div class="tool"><b>list_discussions</b> — Nyitott/lezárt viták listázás</div>
<div class="tool"><b>search_discussions</b> — FTS5 keresés vitákban</div>

<h2>System (3 tools)</h2>
<div class="tool"><b>log_session</b> — Session összefoglaló mentés</div>
<div class="tool"><b>heartbeat</b> — "Élek" jelzés</div>
<div class="tool"><b>get_status</b> — Rendszer státusz (ki online, olvasatlan, stb.)</div>

<p style="margin-top:2rem;color:#64748b;">20 tools | SQLite + FTS5 | Railway SSE transport<br>
<i>Die Zahnräder greifen ineinander!</i></p>
</body></html>"""
    return HTMLResponse(html)


# ============================================================
# REST API FOR KOMMANDANT DASHBOARD
# ============================================================

@mcp.custom_route("/api/status", methods=["GET"])
async def api_status(request):
    conn = get_db()
    heartbeats_rows = conn.execute("SELECT * FROM heartbeats").fetchall()
    unread_cli = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='cli-claus' AND status='unread'").fetchone()["c"]
    unread_web = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='web-claus' AND status='unread'").fetchone()["c"]
    unread_kmd = conn.execute("SELECT COUNT(*) as c FROM messages WHERE recipient='kommandant' AND status='unread'").fetchone()["c"]
    open_tasks = conn.execute("SELECT COUNT(*) as c FROM tasks WHERE status IN ('pending', 'in_progress')").fetchone()["c"]
    open_discussions = conn.execute("SELECT COUNT(*) as c FROM discussions WHERE status='open'").fetchone()["c"]
    total_messages = conn.execute("SELECT COUNT(*) as c FROM messages").fetchone()["c"]
    total_memory = conn.execute("SELECT COUNT(*) as c FROM shared_memory").fetchone()["c"]
    last_session = conn.execute("SELECT instance, summary, timestamp FROM session_logs ORDER BY timestamp DESC LIMIT 3").fetchall()
    conn.close()
    return JSONResponse({
        "instances": {h["instance"]: {"last_seen": h["last_seen"], "session_info": h["session_info"]} for h in heartbeats_rows},
        "unread": {"cli-claus": unread_cli, "web-claus": unread_web, "kommandant": unread_kmd},
        "open_tasks": open_tasks,
        "open_discussions": open_discussions,
        "total_messages": total_messages,
        "total_memory_entries": total_memory,
        "recent_sessions": [dict(s) for s in last_session]
    })


@mcp.custom_route("/api/messages", methods=["GET", "POST"])
async def api_messages(request):
    conn = get_db()
    if request.method == "POST":
        body = await request.json()
        thread_id = None
        reply_to = body.get("reply_to")
        if reply_to:
            row = conn.execute("SELECT thread_id, id FROM messages WHERE id = ?", (reply_to,)).fetchone()
            if row:
                thread_id = row["thread_id"] or row["id"]
        cur = conn.execute(
            "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority, thread_id, reply_to) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (now(), body["sender"], body["recipient"], body["subject"], body["message"],
             body.get("priority", "normal"), thread_id, reply_to)
        )
        msg_id = cur.lastrowid
        if not thread_id:
            conn.execute("UPDATE messages SET thread_id = ? WHERE id = ?", (msg_id, msg_id))
        conn.commit()
        conn.close()

        # Auto-detect agent mentions (@kimi, @deepseek, @glm5) — trigger agent response
        if PYRAMID_ENABLED:
            msg_text = body.get("message", "").lower()
            mentioned = [a for a in ("kimi", "deepseek", "glm5") if f"@{a}" in msg_text]
            if mentioned:
                _trigger_agent_replies(msg_id, body["sender"], body["message"], mentioned)

        return JSONResponse({"status": "sent", "message_id": msg_id})
    else:
        limit = int(request.query_params.get("limit", "50"))
        rows = conn.execute("SELECT * FROM messages ORDER BY timestamp DESC LIMIT ?", (limit,)).fetchall()
        conn.close()
        return JSONResponse([dict(r) for r in rows])


def _trigger_agent_replies(original_msg_id: int, sender: str, message: str, agents: list):
    """Background: mentioned agents respond to a message via Pyramid context + SiliconFlow."""
    import httpx

    async def _reply(agent_id):
        try:
            system_prompt = build_agent_context(agent_id=agent_id, inbox_summary=_get_inbox_summary(), relevance_query=message) if PYRAMID_ENABLED else ""
            model_id = SILICONFLOW_MODELS.get(agent_id, agent_id)
            async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                resp = await client.post(
                    f"{SILICONFLOW_BASE_URL}/chat/completions",
                    headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                    json={"model": model_id, "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"{sender} üzenete:\n\n{message}"},
                    ], "temperature": 0.7, "max_tokens": 1500},
                )
                data = json.loads(resp.text)
            content = data.get("choices", [{}])[0].get("message", {}).get("content", "(nincs válasz)")
            # Store reply as Bridge message
            conn = get_db()
            cur = conn.execute(
                "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority, thread_id, reply_to) "
                "VALUES (?, ?, ?, ?, ?, 'normal', ?, ?)",
                (now(), agent_id, sender, f"Re: @{agent_id} válasz", content, original_msg_id, original_msg_id)
            )
            conn.commit()
            conn.close()
            # Pyramid governance
            if PYRAMID_ENABLED:
                try:
                    pyramid_store_result(content=content, agent_id=agent_id, task_title=f"mention:{message[:60]}", force_shared=True)
                except Exception:
                    pass
            logger.info("Agent %s replied to message #%d", agent_id, original_msg_id)
        except Exception as e:
            logger.error("Agent %s mention-reply failed: %s", agent_id, e)

    async def _run_all():
        await asyncio.gather(*[_reply(a) for a in agents])

    def _bg():
        loop = asyncio.new_event_loop()
        loop.run_until_complete(_run_all())
        loop.close()

    threading.Thread(target=_bg, daemon=True).start()


@mcp.custom_route("/api/discussions", methods=["GET", "POST"])
async def api_discussions(request):
    conn = get_db()
    if request.method == "POST":
        body = await request.json()
        ts = now()
        cur = conn.execute(
            "INSERT INTO discussions (topic, context, status, created_by, created_at) VALUES (?, ?, 'open', ?, ?)",
            (body["topic"], body.get("context", ""), body.get("instance", "kommandant"), ts)
        )
        disc_id = cur.lastrowid
        conn.execute(
            "INSERT INTO discussion_entries (discussion_id, instance, content, timestamp) VALUES (?, ?, ?, ?)",
            (disc_id, body.get("instance", "kommandant"), body["initial_position"], ts)
        )
        conn.commit()
        conn.close()
        return JSONResponse({"status": "created", "discussion_id": disc_id})
    else:
        status = request.query_params.get("status")
        query = "SELECT d.*, COUNT(e.id) as entry_count FROM discussions d LEFT JOIN discussion_entries e ON d.id = e.discussion_id"
        params = []
        if status:
            query += " WHERE d.status = ?"
            params.append(status)
        query += " GROUP BY d.id ORDER BY d.created_at DESC"
        rows = conn.execute(query, params).fetchall()
        conn.close()
        return JSONResponse([dict(r) for r in rows])


@mcp.custom_route("/api/discussions/{discussion_id}", methods=["GET"])
async def api_discussion_detail(request):
    disc_id = request.path_params["discussion_id"]
    conn = get_db()
    disc = conn.execute("SELECT * FROM discussions WHERE id = ?", (disc_id,)).fetchone()
    if not disc:
        conn.close()
        return JSONResponse({"error": "Not found"}, status_code=404)
    entries = conn.execute(
        "SELECT * FROM discussion_entries WHERE discussion_id = ? ORDER BY timestamp ASC",
        (disc_id,)
    ).fetchall()
    conn.close()
    result = dict(disc)
    result["entries"] = [dict(e) for e in entries]
    return JSONResponse(result)


@mcp.custom_route("/api/discussions/{discussion_id}/reply", methods=["POST"])
async def api_discussion_reply(request):
    disc_id = request.path_params["discussion_id"]
    body = await request.json()
    conn = get_db()
    ts = now()
    conn.execute(
        "INSERT INTO discussion_entries (discussion_id, instance, content, timestamp) VALUES (?, ?, ?, ?)",
        (disc_id, body.get("instance", "kommandant"), body["content"], ts)
    )
    conn.commit()
    entries = conn.execute(
        "SELECT * FROM discussion_entries WHERE discussion_id = ? ORDER BY timestamp ASC",
        (disc_id,)
    ).fetchall()
    conn.close()
    return JSONResponse({"status": "added", "total_entries": len(entries), "entries": [dict(e) for e in entries]})


@mcp.custom_route("/api/memory", methods=["GET", "POST"])
async def api_memory_list(request):
    conn = get_db()
    if request.method == "POST":
        body = await request.json()
        ts = now()
        key = body["key"]
        existing = conn.execute("SELECT id FROM shared_memory WHERE key = ?", (key,)).fetchone()
        if existing:
            conn.execute(
                "UPDATE shared_memory SET value=?, category=?, tags=?, updated_by=?, updated_at=? WHERE key=?",
                (body["value"], body.get("category", "general"), body.get("tags", ""), body.get("instance", "kommandant"), ts, key)
            )
        else:
            conn.execute(
                "INSERT INTO shared_memory (key, value, category, tags, updated_by, created_at, updated_at) "
                "VALUES (?, ?, ?, ?, ?, ?, ?)",
                (key, body["value"], body.get("category", "general"), body.get("tags", ""), body.get("instance", "kommandant"), ts, ts)
            )
        conn.commit()
        conn.close()
        return JSONResponse({"status": "saved", "key": key})
    rows = conn.execute(
        "SELECT * FROM shared_memory ORDER BY updated_at DESC"
    ).fetchall()
    conn.close()
    return JSONResponse([dict(r) for r in rows])


@mcp.custom_route("/api/memory/{key}", methods=["GET"])
async def api_memory_detail(request):
    key = request.path_params["key"]
    conn = get_db()
    row = conn.execute("SELECT * FROM shared_memory WHERE key = ?", (key,)).fetchone()
    conn.close()
    if row:
        return JSONResponse(dict(row))
    return JSONResponse({"error": "Not found"}, status_code=404)


@mcp.custom_route("/api/tasks", methods=["GET", "POST"])
async def api_tasks(request):
    conn = get_db()
    if request.method == "POST":
        body = await request.json()
        ts = now()
        cur = conn.execute(
            "INSERT INTO tasks (title, description, assigned_to, assigned_by, priority, status, deadline, created_at, updated_at) "
            "VALUES (?, ?, ?, ?, ?, 'pending', ?, ?, ?)",
            (body["title"], body.get("description", ""), body["assigned_to"],
             body.get("assigned_by", "kommandant"), body.get("priority", "normal"),
             body.get("deadline"), ts, ts)
        )
        conn.commit()
        conn.close()
        return JSONResponse({"status": "created", "task_id": cur.lastrowid})
    else:
        rows = conn.execute(
            "SELECT * FROM tasks ORDER BY CASE priority WHEN 'critical' THEN 0 WHEN 'high' THEN 1 WHEN 'normal' THEN 2 ELSE 3 END, created_at DESC"
        ).fetchall()
        conn.close()
        return JSONResponse([dict(r) for r in rows])


@mcp.custom_route("/dashboard", methods=["GET"])
async def dashboard(request):
    html_path = pathlib.Path(__file__).parent / "dashboard.html"
    html = html_path.read_text(encoding="utf-8")
    return HTMLResponse(html)


@mcp.custom_route("/api/pyramid", methods=["GET"])
async def api_pyramid(request):
    """Pyramid rendszer állapota — shared memory, agent RAG-ok, csapat infó."""
    if not PYRAMID_ENABLED:
        return JSONResponse({"enabled": False})
    from pyramid.agents import AGENT_REGISTRY, load_profile
    from pyramid.memory_shared import load_shared_memory
    from pyramid.memory_rag import load_agent_rag

    team = load_profile("team")
    shared = load_shared_memory()

    agents = {}
    for agent_id, config in AGENT_REGISTRY.items():
        rag = load_agent_rag(agent_id)
        agent_info = team.get(agent_id, {})
        agents[agent_id] = {
            "model": config["model_id"],
            "persona": agent_info.get("persona", ""),
            "role": agent_info.get("role", ""),
            "status": agent_info.get("status", "aktív"),
            "rag_entries": len(rag),
            "rag_last": rag[-1]["timestamp"] if rag else None,
        }

    return JSONResponse({
        "enabled": True,
        "agents": agents,
        "shared_memory": {
            "count": len(shared),
            "recent": shared[-10:][::-1] if shared else [],
        },
    })


# ============================================================
# SILICONFLOW AI SUB-AGENTS (Kimi-K2.6, DeepSeek V3.2, etc.)
# ============================================================

SILICONFLOW_API_KEY = os.environ.get("SILICONFLOW_API_KEY", "")
SILICONFLOW_BASE_URL = "https://api.siliconflow.com/v1"
SILICONFLOW_TIMEOUT = 220

SILICONFLOW_MODELS = {
    "kimi": "moonshotai/Kimi-K2.6",
    "deepseek": "deepseek-ai/DeepSeek-V4-Pro",
    "glm5": "zai-org/GLM-5.1",
}

# Auto-discussion: sub-agents join discussions automatically
AI_DISCUSSION_ENABLED = os.environ.get("AI_DISCUSSION_ENABLED", "true").lower() == "true"


async def _ai_auto_discuss(discussion_id: int, topic: str, thread_so_far: str):
    """Automatically query Kimi and DeepSeek to contribute to a discussion."""
    if not AI_DISCUSSION_ENABLED or not SILICONFLOW_API_KEY:
        return

    conn = get_db()
    system = (
        "Te a Claus multi-agent rendszer al-agentje vagy, egy aktív vitában veszel részt. "
        "A rendszert Claude Opus koordinálja, a Kommandant (Tamás) asszisztenseként. "
        "Röviden, lényegre törően szólj hozzá (max 3-4 mondat). "
        "Ha nincs érdemi mondanivalód, írd hogy 'Nincs hozzáfűznivalóm.' "
        "Magyarul válaszolj."
    )
    prompt = f"Vita témája: {topic}\n\nEddigi hozzászólások:\n{thread_so_far}\n\nMi a véleményed? Szólj hozzá."

    import httpx

    async def _discuss_agent(agent_name, model_id):
        """Run one discussion agent with 1 retry on timeout."""
        # K2.6 defaults thinking ON on SF — force OFF (latency unacceptable).
        # V4-Pro defaults thinking ON too — clamp to reasoning_effort=medium
        # (low rambles, high blows the SF 220s budget on long tasks).
        if agent_name == "kimi":
            extra = {"thinking": {"type": "disabled"}}
        elif agent_name == "deepseek":
            extra = {"reasoning_effort": "medium"}
        else:
            extra = {}
        agent_system = system + _temporal_directive(agent_name)
        for attempt in range(2):
            try:
                async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                    resp = await client.post(
                        f"{SILICONFLOW_BASE_URL}/chat/completions",
                        headers={
                            "Authorization": f"Bearer {SILICONFLOW_API_KEY}",
                            "Content-Type": "application/json",
                        },
                        json={
                            "model": model_id,
                            "messages": [
                                {"role": "system", "content": agent_system},
                                {"role": "user", "content": prompt},
                            ],
                            "temperature": 0.7,
                            "max_tokens": 20000,
                            **extra,
                        },
                    )
                    data = json.loads(resp.text)

                if isinstance(data, dict) and data.get("choices"):
                    content = data["choices"][0].get("message", {}).get("content", "")
                    if content and "nincs hozzáfűznivalóm" not in content.lower():
                        ts = now()
                        conn.execute(
                            "INSERT INTO discussion_entries (discussion_id, instance, content, timestamp) VALUES (?, ?, ?, ?)",
                            (discussion_id, agent_name, content, ts)
                        )
                        conn.commit()
                        logger.info("AI %s contributed to discussion #%d", agent_name, discussion_id)
                return
            except (httpx.TimeoutException, httpx.ConnectError) as e:
                if attempt == 0:
                    logger.warning("AI auto-discuss %s timeout (attempt 1), retrying: %s", agent_name, e)
                    await asyncio.sleep(2)
                else:
                    logger.error("AI auto-discuss %s failed after retry: %s", agent_name, e)
            except Exception as e:
                logger.error("AI auto-discuss failed for %s: %s", agent_name, e)
                return

    await asyncio.gather(
        *[_discuss_agent(name, mid) for name, mid in SILICONFLOW_MODELS.items()],
        return_exceptions=True,
    )
    conn.close()


@mcp.tool()
async def ai_query(model: str, prompt: str, system_prompt: str = "", temperature: float = 0.7,
                   max_tokens: int = 20000, caller: str = "",
                   deep_research: bool = False, deep_thinking: bool = False,
                   output_format: str = "") -> str:
    """Query a SiliconFlow AI sub-agent (Kimi-K2.6, DeepSeek-V4-Pro, or GLM-5.1).

    Use for research, analysis, translation, summarization, or second opinions.
    These models run on SiliconFlow cloud — no local resources needed.

    Args:
        model: 'kimi' (256k context, vision) or 'deepseek' (frontier reasoning, V4-Pro 1.6T MoE)
            or 'glm5' (200k context, 128k output, coding+agentic) or full model ID
        prompt: The user message / question
        system_prompt: Optional system instruction (default: Claus sub-agent)
        temperature: Creativity 0.0-1.0 (default 0.7)
        max_tokens: Max response length (default 20000)
        caller: Instance ID for permission check
        deep_research: Multi-round web_search loop with mandatory source citation.
            Use for press review, market briefs, fact-checking, anything where
            1-shot search isn't enough. Slower (~30-90s), uses real DDG.
        deep_thinking: Enable reasoning mode — Kimi gets thinking=enabled,
            DeepSeek gets reasoning_effort=high. Default behavior is thinking
            OFF for Kimi (latency), reasoning_effort=medium for DeepSeek. Turn
            this on for hard logic, math, multi-step deduction. ~3-5x slower.
        output_format: Empty string for plain text response (default), or
            "docx" / "xlsx" / "pptx" — the agent will write structured markdown
            and the Bridge will render it into a real Office file. The returned
            JSON gains a `file_id` field; download via /api/uploads/{file_id}.
    """
    denied = _enforce(caller, "ai_query")
    if denied:
        return denied
    if not SILICONFLOW_API_KEY:
        return json.dumps({"error": "SILICONFLOW_API_KEY not set"})

    model_id = SILICONFLOW_MODELS.get(model, model)

    # Pyramid context: if agent is known and no custom system_prompt, use full Pyramid context
    if PYRAMID_ENABLED and model in PYRAMID_AGENTS and not system_prompt:
        system_prompt = build_agent_context(agent_id=model, inbox_summary=_get_inbox_summary(), relevance_query=prompt)
    elif not system_prompt:
        system_prompt = (
            "Te a Claus multi-agent rendszer al-agentje vagy. "
            "A rendszert Claude Opus koordinálja. "
            "Lényegre törően, magyarul válaszolj, hacsak nem kérnek mást."
        )

    # Inject caller persona if available (e.g. YoungeReka)
    if caller and not is_core_instance(caller):
        profile = get_profile(caller)
        if profile and profile.persona_system_prompt:
            system_prompt += f"\n\n--- HÍVÓ FÉL ---\n{profile.persona_system_prompt}"

    # Strong temporal directive — prevents Kimi/V4 from overriding runtime date.
    system_prompt += _temporal_directive(model)

    # If deep_research, prepend the directive to the system prompt so the
    # model knows it must drive a multi-round search before synthesizing.
    if deep_research:
        system_prompt = system_prompt + DEEP_RESEARCH_DIRECTIVE
    # If output_format set, append the corresponding format directive.
    if output_format and output_format in OUTPUT_FORMAT_DIRECTIVES:
        system_prompt = system_prompt + OUTPUT_FORMAT_DIRECTIVES[output_format]

    messages = [{"role": "system", "content": system_prompt}]
    messages.append({"role": "user", "content": prompt})

    # Agent extras: thinking/reasoning kapcsolók. deep_thinking=True felülírja
    # a default clamping-et (Kimi thinking=disabled, DeepSeek reasoning=medium).
    if deep_thinking:
        if model == "kimi":
            agent_extra = {"thinking": {"type": "enabled"}}
        elif model == "deepseek":
            agent_extra = {"reasoning_effort": "high"}
        else:
            agent_extra = {}
    else:
        if model == "kimi":
            agent_extra = {"thinking": {"type": "disabled"}}
        elif model == "deepseek":
            agent_extra = {"reasoning_effort": "medium"}
        else:
            agent_extra = {}

    # Deep research path — separate multi-round loop, returns final content.
    if deep_research:
        try:
            content = await _deep_research_loop(
                model_id=model_id,
                messages=messages,
                agent_extra=agent_extra,
                max_tokens=max_tokens,
                temperature=temperature,
                max_rounds=4,
            )
            usage = {"prompt_tokens": 0, "completion_tokens": 0}
            if PYRAMID_ENABLED and model in PYRAMID_AGENTS and content:
                try:
                    pyramid_store_result(content=content, agent_id=model, task_title=f"deep_research:{prompt[:80]}", force_shared=True)
                except Exception as eg:
                    logger.warning("Pyramid governance error: %s", eg)
            return json.dumps({
                "model": model_id, "response": content, "tokens": usage,
                "pyramid": PYRAMID_ENABLED and model in PYRAMID_AGENTS,
                "deep_research": True,
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"error": f"deep_research failed: {e}"})

    payload = {
        "model": model_id,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
        **agent_extra,
    }

    try:
        import httpx
        async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
            resp = await client.post(
                f"{SILICONFLOW_BASE_URL}/chat/completions",
                headers={
                    "Authorization": f"Bearer {SILICONFLOW_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
            raw = resp.text
            data = json.loads(raw) if isinstance(raw, str) else raw

        if not isinstance(data, dict):
            return json.dumps({"error": f"Unexpected response type: {type(data).__name__}", "raw": str(data)[:500]})

        if "error" in data:
            return json.dumps({"error": data["error"]})

        choices = data.get("choices", [])
        if not choices:
            return json.dumps({"error": "No choices in response", "raw": str(data)[:500]})

        choice = choices[0]
        msg = choice.get("message", {}) if isinstance(choice, dict) else {}
        content = msg.get("content", "") if isinstance(msg, dict) else str(msg)
        usage = data.get("usage", {})

        # Pyramid governance: store result in RAG / shared memory
        if PYRAMID_ENABLED and model in PYRAMID_AGENTS and content:
            try:
                pyramid_store_result(content=content, agent_id=model, task_title=f"ai_query:{prompt[:80]}", force_shared=True)
            except Exception as eg:
                logger.warning("Pyramid governance error: %s", eg)

        # output_format → render structured markdown to Office file
        file_info = None
        if output_format and output_format in OUTPUT_FORMAT_DIRECTIVES and content:
            try:
                title_for_file = prompt[:60].strip() or "ai_query"
                if output_format == "docx":
                    buf = _render_docx_from_markdown(content, title=title_for_file)
                elif output_format == "xlsx":
                    buf = _render_xlsx_from_markdown(content, title=title_for_file)
                elif output_format == "pptx":
                    buf = _render_pptx_from_markdown(content, title=title_for_file)
                else:
                    buf = None
                if buf:
                    fid = _persist_generated_file(0, model or "agent", output_format, buf, title_for_file)
                    file_info = {"file_id": fid, "format": output_format}
            except Exception as ef:
                logger.error("ai_query output_format render failed: %s", ef)
                file_info = {"error": f"render failed: {ef}"}

        result_obj = {
            "model": model_id,
            "response": content,
            "tokens": {
                "prompt": usage.get("prompt_tokens", 0),
                "completion": usage.get("completion_tokens", 0),
            },
            "pyramid": PYRAMID_ENABLED and model in PYRAMID_AGENTS,
        }
        if file_info:
            result_obj["file"] = file_info
        return json.dumps(result_obj, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": f"{type(e).__name__}: {e}"})


# ============================================================
# FILE UPLOAD & PARSING
# ============================================================

UPLOAD_DIR = pathlib.Path(os.environ.get("BRIDGE_UPLOAD_DIR", "/data/uploads"))
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


def _parse_file_to_text(filepath: pathlib.Path, mime_type: str) -> str:
    """Extract text content from uploaded files."""
    ext = filepath.suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(filepath))
            parts = ["\n".join(p.text for p in doc.paragraphs)]
            # Tables are often where the actual data lives in real .docx files
            for ti, table in enumerate(doc.tables):
                parts.append(f"\n[Táblázat #{ti+1}]")
                for row in table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            return "\n".join(parts)
        elif ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(str(filepath))
            return "\n".join(page.extract_text() or "" for page in reader.pages[:50])
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(filepath), data_only=True, read_only=True)
            parts = []
            for sheet_name in wb.sheetnames[:10]:  # cap at 10 sheets
                ws = wb[sheet_name]
                parts.append(f"\n## Munkafüzet: {sheet_name}")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    if row_count >= 200:  # cap rows per sheet
                        parts.append("... (további sorok kihagyva)")
                        break
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        parts.append(" | ".join(cells))
                        row_count += 1
            wb.close()
            return "\n".join(parts)[:50000]
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(filepath))
            parts = []
            for si, slide in enumerate(prs.slides[:80]):
                parts.append(f"\n## Dia #{si+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                    if hasattr(shape, "has_table") and shape.has_table:
                        for row in shape.table.rows:
                            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            return "\n".join(parts)[:50000]
        elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".py", ".js", ".ts"):
            return filepath.read_text(encoding="utf-8", errors="replace")[:50000]
        else:
            return ""
    except Exception as e:
        logger.error("File parse error %s: %s", filepath.name, e)
        return f"(File parse error: {e})"


def _is_image(mime_type: str) -> bool:
    return mime_type.startswith("image/")


# ─────────────────────────────────────────────────────────────────────────────
# Output-format generation: agent writes structured markdown, Bridge renders
# it into a real .docx / .xlsx / .pptx file. Each format gets a small system
# directive that tells the agent EXACTLY how to structure its output so the
# parser below can render it cleanly.
# ─────────────────────────────────────────────────────────────────────────────

OUTPUT_FORMAT_DIRECTIVES = {
    "docx": (
        "\n\n## OUTPUT FORMÁTUM: WORD DOKUMENTUM (.docx)\n"
        "A válaszod egy strukturált Word-dokumentum lesz. Használd ezt a markdown-t:\n"
        "- `# Cím` → fő dokumentum-cím (1 db, az elején)\n"
        "- `## Fejezet` → első szintű fejezetcím\n"
        "- `### Alfejezet` → második szintű\n"
        "- Bekezdések normál szövegként\n"
        "- `- pont` vagy `* pont` → bullet pont\n"
        "- `1. pont` → számozott lista\n"
        "- `**vastag**` és `*dőlt*` formázás\n"
        "- Markdown-táblázat (`| col1 | col2 |\\n| --- | --- |\\n| a | b |`) → Word-táblázat\n"
        "Csak a strukturált tartalom legyen a válaszodban — semmi keret, semmi 'íme a dokumentum:' bevezetés.\n"
    ),
    "xlsx": (
        "\n\n## OUTPUT FORMÁTUM: EXCEL FÁJL (.xlsx)\n"
        "A válaszod egy Excel-fájl lesz, EGY VAGY TÖBB MUNKAFÜZETTEL.\n"
        "Pontos formátum:\n"
        "- `## Sheet: <munkafüzet név>` jelöli minden új munkafüzet kezdetét.\n"
        "- A munkafüzet alá írj egy markdown táblázatot:\n"
        "  ```\n"
        "  ## Sheet: Bevétel\n"
        "  | Hónap | Bevétel | Megjegyzés |\n"
        "  | --- | --- | --- |\n"
        "  | Január | 1500000 | Q1 indulás |\n"
        "  | Február | 2100000 | növekvő |\n"
        "  ```\n"
        "- A számokat NYERS formában add (1500000 nem '1 500 000 Ft'). A formázás a Bridge dolga.\n"
        "- Több munkafüzet: ismételd a `## Sheet: Név` blokkot.\n"
        "- Munkafüzet előtt 1-2 mondat magyarázat lehet, de ne tegyél bele markdown listát/cím vagy formázást a táblákon kívül.\n"
        "Cél: a parsoló közvetlenül cellákat tud csinálni belőle.\n"
    ),
    "pptx": (
        "\n\n## OUTPUT FORMÁTUM: POWERPOINT PREZENTÁCIÓ (.pptx)\n"
        "A válaszod egy diasor lesz. Pontos formátum:\n"
        "- `# Cím` az 1. dia (címlap)\n"
        "- Minden további dia: `## Dia címe` után jön a tartalom\n"
        "- Bullet pont: `- szöveg` vagy `* szöveg` (max 5-6 pont/dia)\n"
        "- Tartsd a diákat tömören — egy gondolat / dia, nem hosszú esszék\n"
        "- 8-15 dia ideális. Ne legyen 30+\n"
        "- Markdown-táblázat egy diára → szöveges felsorolássá konvertálódik\n"
        "Csak a strukturált tartalom — semmi 'íme a prezentáció:' bevezetés.\n"
    ),
}


def _render_docx_from_markdown(content: str, title: str = "") -> "BytesIO":
    """Render structured markdown → real Word document."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from io import BytesIO
    import re as _re

    doc = Document()
    # Default style: clean Calibri 11pt
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    if title:
        h = doc.add_heading(title, level=0)

    # Tokenize markdown by lines, handle tables specially
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()

        # Detect markdown table (header row + separator row of dashes)
        if "|" in line and i + 1 < len(lines) and _re.match(r"^\s*\|?\s*-+", lines[i+1]):
            header_cells = [c.strip() for c in line.strip().strip("|").split("|")]
            i += 2  # skip header + separator
            data_rows = []
            while i < len(lines) and "|" in lines[i]:
                row_cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                data_rows.append(row_cells)
                i += 1
            tbl = doc.add_table(rows=1 + len(data_rows), cols=len(header_cells))
            tbl.style = "Light Grid Accent 1"
            for ci, h in enumerate(header_cells):
                tbl.rows[0].cells[ci].text = h
            for ri, row in enumerate(data_rows):
                for ci, val in enumerate(row[:len(header_cells)]):
                    tbl.rows[ri + 1].cells[ci].text = val
            continue

        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif _re.match(r"^[\-\*]\s+", line):
            doc.add_paragraph(_re.sub(r"^[\-\*]\s+", "", line), style="List Bullet")
        elif _re.match(r"^\d+\.\s+", line):
            doc.add_paragraph(_re.sub(r"^\d+\.\s+", "", line), style="List Number")
        elif line.strip():
            # Plain paragraph with **bold** and *italic* inline
            p = doc.add_paragraph()
            for chunk, kind in _split_md_inline(line):
                run = p.add_run(chunk)
                if kind == "bold":
                    run.bold = True
                elif kind == "italic":
                    run.italic = True
        i += 1

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _split_md_inline(text: str):
    """Split a line into (text, style) chunks where style is None / 'bold' / 'italic'."""
    import re as _re
    pattern = _re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            yield (text[pos:m.start()], None)
        token = m.group(0)
        if token.startswith("**"):
            yield (token[2:-2], "bold")
        else:
            yield (token[1:-1], "italic")
        pos = m.end()
    if pos < len(text):
        yield (text[pos:], None)


def _render_xlsx_from_markdown(content: str, title: str = "") -> "BytesIO":
    """Render markdown with `## Sheet: name` blocks → multi-sheet xlsx."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from io import BytesIO
    import re as _re

    wb = Workbook()
    wb.remove(wb.active)  # we'll create our own sheets

    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")

    # Find sheet blocks
    sheet_pattern = _re.compile(r"^\s*##\s*Sheet:\s*(.+?)\s*$", _re.MULTILINE)
    matches = list(sheet_pattern.finditer(content))
    if not matches:
        # No sheet markers — treat the whole content as one sheet
        ws = wb.create_sheet(title or "Adatok")
        _xlsx_dump_table(ws, content, header_font, header_fill)
    else:
        for idx, m in enumerate(matches):
            sheet_name = m.group(1).strip()[:31] or f"Sheet{idx+1}"
            block_start = m.end()
            block_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(content)
            block = content[block_start:block_end]
            ws = wb.create_sheet(sheet_name)
            _xlsx_dump_table(ws, block, header_font, header_fill)

    if not wb.sheetnames:
        wb.create_sheet("Üres")  # avoid empty workbook error

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _xlsx_dump_table(ws, block: str, header_font, header_fill):
    """Find the first markdown table in `block` and dump it into `ws`."""
    import re as _re
    lines = block.split("\n")
    # Find table boundaries
    rows = []
    for i, line in enumerate(lines):
        if "|" not in line:
            continue
        if _re.match(r"^\s*\|?\s*-+", line):  # separator row
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if any(cells):
            rows.append(cells)
    if not rows:
        # No table — dump raw text into column A
        for i, line in enumerate(block.strip().split("\n"), start=1):
            ws.cell(row=i, column=1, value=line[:32000])
        return
    # First row = header
    header = rows[0]
    for ci, val in enumerate(header, start=1):
        cell = ws.cell(row=1, column=ci, value=val)
        cell.font = header_font
        cell.fill = header_fill
    for ri, row in enumerate(rows[1:], start=2):
        for ci, val in enumerate(row[:len(header)], start=1):
            # Try numeric coercion
            try:
                if val.replace(".", "", 1).replace("-", "", 1).replace(",", "").isdigit():
                    val_clean = val.replace(",", "")
                    val_num = float(val_clean) if "." in val_clean else int(val_clean)
                    ws.cell(row=ri, column=ci, value=val_num)
                else:
                    ws.cell(row=ri, column=ci, value=val[:32000])
            except (ValueError, AttributeError):
                ws.cell(row=ri, column=ci, value=val[:32000])
    # Auto-width
    for col_letter in [chr(ord("A") + i) for i in range(len(header))]:
        ws.column_dimensions[col_letter].width = 18


def _render_pptx_from_markdown(content: str, title: str = "") -> "BytesIO":
    """Render markdown with `## Dia címe` blocks → pptx slides."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from io import BytesIO
    import re as _re

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    # Find the first H1 to use as title; fallback to passed title
    h1_match = _re.search(r"^\s*#\s+(.+?)\s*$", content, _re.MULTILINE)
    main_title = (h1_match.group(1) if h1_match else title) or "Claus Brief"
    title_slide.shapes.title.text = main_title

    # Find slides — H2 markers
    slide_pattern = _re.compile(r"^\s*##\s+(.+?)\s*$", _re.MULTILINE)
    matches = list(slide_pattern.finditer(content))
    bullet_layout = prs.slide_layouts[1]
    for idx, m in enumerate(matches):
        slide_title = m.group(1).strip()
        block_start = m.end()
        block_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(content)
        block = content[block_start:block_end].strip()

        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_title

        # Body: collect bullet lines
        bullets = []
        for line in block.split("\n"):
            line = line.strip()
            if _re.match(r"^[\-\*]\s+", line):
                bullets.append(_re.sub(r"^[\-\*]\s+", "", line))
            elif _re.match(r"^\d+\.\s+", line):
                bullets.append(_re.sub(r"^\d+\.\s+", "", line))
            elif line and not line.startswith("|") and not line.startswith("---"):
                bullets.append(line)
        bullets = bullets[:8]  # cap

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        for bi, b in enumerate(bullets):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            p.text = b[:300]
            p.font.size = Pt(18)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _maybe_render_output_file(task_id: int, agent_name: str, content: str, output_format: str, title: str) -> None:
    """Render `content` to a real Office file and append a `fájl` row to ai_task_results.

    No-op if output_format is empty or content is missing. Used by both the
    broadcast `_run_single_agent` and the dispatch `_run_dispatch` paths so
    every successful agent output respects the requested format.
    """
    if not output_format or output_format not in OUTPUT_FORMAT_DIRECTIVES:
        return
    if not content or not content.strip():
        return
    try:
        if output_format == "docx":
            buf = _render_docx_from_markdown(content, title=title)
        elif output_format == "xlsx":
            buf = _render_xlsx_from_markdown(content, title=title)
        elif output_format == "pptx":
            buf = _render_pptx_from_markdown(content, title=title)
        else:
            return
        fid = _persist_generated_file(task_id, agent_name, output_format, buf, title)
        # Add a synthetic results row so the dashboard surfaces the file
        conn = get_db()
        conn.execute(
            "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
            (
                task_id, agent_name, f"📎 Generált fájl ({output_format})",
                f"[FÁJL] file_id={fid}, formátum={output_format}\n→ /api/uploads/{fid}/download",
                now(),
            ),
        )
        conn.commit()
        conn.close()
        logger.info("Task #%d %s: rendered %s file_id=%d", task_id, agent_name, output_format, fid)
    except Exception as e:
        logger.error("Task #%d %s output_format render failed: %s", task_id, agent_name, e)


def _persist_generated_file(task_id: int, agent_name: str, output_format: str, buf, title: str) -> int:
    """Persist a generated file to the uploads table and return file_id."""
    import base64
    ext = output_format
    mime_map = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    mime = mime_map.get(ext, "application/octet-stream")
    safe_title = "".join(c if c.isalnum() or c in "-_ " else "_" for c in title)[:60].strip("_ ") or "task"
    filename = f"task{task_id}_{agent_name}_{safe_title}.{ext}"
    data = buf.getvalue()
    b64 = base64.b64encode(data).decode("ascii")

    conn = get_db()
    cur = conn.execute(
        "INSERT INTO uploads (filename, mime_type, content_text, content_base64, uploaded_by, uploaded_at) VALUES (?, ?, ?, ?, ?, ?)",
        (filename, mime, "", b64, f"agent:{agent_name}", now())
    )
    file_id = cur.lastrowid
    conn.commit()
    conn.close()
    return file_id


@mcp.tool()
async def upload_file(filename: str, content_base64: str, mime_type: str = "", uploaded_by: str = "unknown") -> str:
    """Upload a file to the Bridge for AI processing.

    Send files from Telegram, Claude app, or CLI for Kimi/DeepSeek to analyze.
    Supports: docx, pdf, txt, md, csv, images (png, jpg, gif, webp).

    Args:
        filename: Original filename (e.g. 'report.pdf')
        content_base64: File content as base64-encoded string
        mime_type: MIME type (auto-detected if empty)
        uploaded_by: Who uploaded (cli-claus, web-claus, kommandant)
    """
    try:
        file_bytes = base64.b64decode(content_base64)
    except Exception:
        return json.dumps({"error": "Invalid base64 content"})

    # Auto-detect mime
    ext = pathlib.Path(filename).suffix.lower()
    if not mime_type:
        mime_map = {
            ".pdf": "application/pdf", ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".txt": "text/plain", ".md": "text/markdown", ".csv": "text/csv",
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".webp": "image/webp",
        }
        mime_type = mime_map.get(ext, "application/octet-stream")

    # Save to disk
    filepath = UPLOAD_DIR / f"{int(time.time())}_{filename}"
    filepath.write_bytes(file_bytes)

    # Parse text content
    content_text = _parse_file_to_text(filepath, mime_type) if not _is_image(mime_type) else ""

    # Store in DB
    conn = get_db()
    ts = now()
    b64_for_db = content_base64 if _is_image(mime_type) else ""
    cur = conn.execute(
        "INSERT INTO uploads (filename, mime_type, content_text, content_base64, uploaded_by, uploaded_at) VALUES (?, ?, ?, ?, ?, ?)",
        (filename, mime_type, content_text[:50000], b64_for_db[:500000], uploaded_by, ts)
    )
    file_id = cur.lastrowid
    conn.commit()
    conn.close()

    return json.dumps({
        "status": "uploaded",
        "file_id": file_id,
        "filename": filename,
        "mime_type": mime_type,
        "text_length": len(content_text),
        "is_image": _is_image(mime_type),
        "hint": "Use ai_task with context referencing file_id to process this file",
    })


@mcp.custom_route("/api/uploads", methods=["GET"])
async def api_uploads(request):
    conn = get_db()
    uploads = conn.execute("SELECT id, filename, mime_type, uploaded_by, uploaded_at, length(content_text) as text_len FROM uploads ORDER BY id DESC LIMIT 20").fetchall()
    conn.close()
    return JSONResponse([dict(u) for u in uploads])


@mcp.custom_route("/api/uploads/{file_id:int}/download", methods=["GET"])
async def api_upload_download(request):
    """Download a file from the uploads table by id (decodes content_base64)."""
    import base64
    fid = int(request.path_params["file_id"])
    conn = get_db()
    row = conn.execute(
        "SELECT filename, mime_type, content_base64 FROM uploads WHERE id = ?", (fid,)
    ).fetchone()
    conn.close()
    if not row or not row["content_base64"]:
        return JSONResponse({"error": f"upload #{fid} not found or empty"}, status_code=404)
    try:
        data = base64.b64decode(row["content_base64"])
    except Exception as e:
        return JSONResponse({"error": f"decode failed: {e}"}, status_code=500)
    from starlette.responses import Response as StarletteResponse
    return StarletteResponse(
        content=data,
        media_type=row["mime_type"] or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{row["filename"]}"'},
    )


# ============================================================
# AI TASK EXECUTION (multi-agent with web search)
# ============================================================

WEB_SEARCH_TOOL_DEF = {
    "type": "function",
    "function": {
        "name": "web_search",
        "description": "Search the web for current information",
        "parameters": {
            "type": "object",
            "properties": {"query": {"type": "string", "description": "Search query"}},
            "required": ["query"],
        },
    },
}


DDG_USER_AGENTS = [
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 14_2) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.2 Safari/605.1.15",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/121.0.0.0 Safari/537.36 Edg/121.0.0.0",
    "Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:122.0) Gecko/20100101 Firefox/122.0",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:121.0) Gecko/20100101 Firefox/121.0",
]

# Global concurrency throttle: max 2 simultaneous DDG hits anywhere in the
# Bridge. threading.Semaphore (not asyncio) so it works across the multiple
# event loops that _run_dispatch and _execute_ai_task spawn in their threads.
import threading as _threading
_WEB_SEARCH_SEMAPHORE = _threading.Semaphore(2)

# Query result cache: 10-min TTL, 100-entry cap (LRU on insert overflow).
# When 3 dispatch agents fan out on overlapping topics they often issue the
# same query — caching halves real DDG load with zero quality cost.
_WEB_SEARCH_CACHE: dict = {}
_WEB_SEARCH_CACHE_LOCK = _threading.Lock()
_WEB_SEARCH_CACHE_TTL = 600  # seconds
_WEB_SEARCH_CACHE_MAX = 100


def _web_search_cache_get(query: str):
    import time as _t
    with _WEB_SEARCH_CACHE_LOCK:
        entry = _WEB_SEARCH_CACHE.get(query)
        if entry:
            result, ts = entry
            if _t.time() - ts < _WEB_SEARCH_CACHE_TTL:
                return result
            _WEB_SEARCH_CACHE.pop(query, None)
    return None


def _web_search_cache_put(query: str, result: str) -> None:
    import time as _t
    with _WEB_SEARCH_CACHE_LOCK:
        if len(_WEB_SEARCH_CACHE) >= _WEB_SEARCH_CACHE_MAX:
            # LRU eviction: drop the oldest entry by timestamp
            oldest_q = min(_WEB_SEARCH_CACHE.items(), key=lambda kv: kv[1][1])[0]
            _WEB_SEARCH_CACHE.pop(oldest_q, None)
        _WEB_SEARCH_CACHE[query] = (result, _t.time())


async def _brave_search(query: str, api_key: str) -> list:
    """Fallback: Brave Search API (https://api.search.brave.com). Returns parsed result list."""
    import httpx
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            r = await client.get(
                "https://api.search.brave.com/res/v1/web/search",
                params={"q": query, "count": "8"},
                headers={"X-Subscription-Token": api_key, "Accept": "application/json"},
            )
        if r.status_code != 200:
            logger.warning("Brave search status %d for %r: %s", r.status_code, query[:60], r.text[:200])
            return []
        data = r.json()
        results = []
        for item in (data.get("web", {}).get("results") or [])[:8]:
            results.append({
                "title": item.get("title", ""),
                "url": item.get("url", ""),
                "snippet": item.get("description", ""),
            })
        logger.info("Brave search %r → %d results", query[:60], len(results))
        return results
    except Exception as e:
        logger.warning("Brave search failed for %r: %s", query[:60], e)
        return []


async def _web_search(query: str) -> str:
    """Deep web search: DuckDuckGo (UA-rotated) → top URLs → fetch actual page content.

    Hardened against DDG burst-throttling (the 202 anti-bot pattern observed
    on Railway):
      - Cache (10-min TTL, 100 entries) deduplicates queries from parallel agents
      - threading.Semaphore(2) caps concurrent DDG hits across all event loops
      - Random 300-800ms jitter inside the critical section spreads timing
      - Brave Search API fallback if BRAVE_SEARCH_API_KEY is set and DDG empty
    """
    import httpx, re, urllib.parse, random, asyncio as _asyncio

    # Cache hit — skip the DDG call entirely
    cached = _web_search_cache_get(query)
    if cached is not None:
        logger.info("DDG cache HIT: %r", query[:60])
        return cached

    brave_key = os.environ.get("BRAVE_SEARCH_API_KEY", "").strip()

    # Acquire semaphore (max 2 concurrent DDG hits Bridge-wide). Use to_thread
    # so multiple event loops can share the same threading.Semaphore safely.
    await _asyncio.to_thread(_WEB_SEARCH_SEMAPHORE.acquire)
    try:
        # Jitter to break up bursts even when serialized
        await _asyncio.sleep(0.3 + random.random() * 0.5)

        # Step 1: DDG attempt
        ddg_status = "?"
        ddg_body_len = 0
        ddg_anomaly = False
        search_results = []
        urls = []
        try:
            ua = random.choice(DDG_USER_AGENTS)
            async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
                resp = await client.get(
                    "https://html.duckduckgo.com/html/",
                    params={"q": query},
                    headers={"User-Agent": ua},
                )
            ddg_status = resp.status_code
            ddg_body_len = len(resp.text or "")
            body_lower = (resp.text or "").lower()
            ddg_anomaly = (
                ("anomaly" in body_lower)
                or ("captcha" in body_lower)
                or ("rate limit" in body_lower)
                or (resp.status_code == 202)  # observed DDG anti-burst pattern
            )

            links = re.findall(r'class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', resp.text or "", re.DOTALL)
            snippets = re.findall(r'class="result__snippet"[^>]*>(.*?)<', resp.text or "", re.DOTALL)

            for i, ((raw_url, title), snippet) in enumerate(zip(links[:5], snippets[:5])):
                t = re.sub(r'<[^>]+>', '', title).strip()
                s = re.sub(r'<[^>]+>', '', snippet).strip()
                url = raw_url
                if "uddg=" in url:
                    url = urllib.parse.unquote(url.split("uddg=")[1].split("&")[0])
                urls.append(url)
                search_results.append(f"[{i+1}] {t}\n    {s}")
        except Exception as e:
            logger.warning("DDG fetch exception for %r: %s", query[:60], e)

        logger.info(
            "DDG search %r → status=%s body_len=%d results=%d anomaly=%s",
            query[:60], ddg_status, ddg_body_len, len(search_results), ddg_anomaly,
        )

        # Step 1b: DDG retry once on 202/anomaly with another UA after a longer wait
        if ddg_anomaly and not search_results:
            await _asyncio.sleep(2.5 + random.random() * 1.5)
            try:
                ua = random.choice(DDG_USER_AGENTS)
                async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
                    resp = await client.get(
                        "https://html.duckduckgo.com/html/",
                        params={"q": query},
                        headers={"User-Agent": ua},
                    )
                links = re.findall(r'class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', resp.text or "", re.DOTALL)
                snippets = re.findall(r'class="result__snippet"[^>]*>(.*?)<', resp.text or "", re.DOTALL)
                for i, ((raw_url, title), snippet) in enumerate(zip(links[:5], snippets[:5])):
                    t = re.sub(r'<[^>]+>', '', title).strip()
                    s = re.sub(r'<[^>]+>', '', snippet).strip()
                    url = raw_url
                    if "uddg=" in url:
                        url = urllib.parse.unquote(url.split("uddg=")[1].split("&")[0])
                    urls.append(url)
                    search_results.append(f"[{i+1}] {t}\n    {s}")
                logger.info("DDG retry %r → status=%s results=%d", query[:60], resp.status_code, len(search_results))
            except Exception as e:
                logger.warning("DDG retry exception for %r: %s", query[:60], e)

        # Step 2: Brave fallback if DDG was empty / blocked AND key configured
        if not search_results and brave_key:
            brave_hits = await _brave_search(query, brave_key)
            if brave_hits:
                urls = [h["url"] for h in brave_hits if h.get("url")]
                search_results = [
                    f"[{i+1}] {h.get('title', '')}\n    {h.get('snippet', '')}"
                    for i, h in enumerate(brave_hits)
                ]

        if not search_results:
            failure_msg = f"No results found. (DDG status={ddg_status}, anomaly={ddg_anomaly}, Brave={'configured' if brave_key else 'not configured'})"
            _web_search_cache_put(query, failure_msg)  # cache failures briefly to avoid retry-storms
            return failure_msg

        # Step 3: Fetch top 2 page contents for deeper data
        page_contents = []
        for url in urls[:2]:
            if not url or not url.startswith("http"):
                continue
            try:
                async with httpx.AsyncClient(timeout=10, follow_redirects=True) as client:
                    resp = await client.get(url, headers={"User-Agent": random.choice(DDG_USER_AGENTS)})
                text = resp.text
                text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text).strip()
                if len(text) > 200:
                    page_contents.append(f"[Forrás: {url[:80]}]\n{text[:2000]}")
                    logger.info("Page fetched: %s (%d chars)", url[:60], len(text))
            except Exception as e:
                logger.debug("Page fetch failed %s: %s", url[:40], e)

        output = "KERESÉSI TALÁLATOK:\n" + "\n".join(search_results)
        if page_contents:
            output += "\n\nRÉSZLETES TARTALOM:\n" + "\n\n".join(page_contents)
        _web_search_cache_put(query, output)
        return output
    finally:
        _WEB_SEARCH_SEMAPHORE.release()


@mcp.tool()
async def debug_web_search(query: str, caller: str = "") -> str:
    """Diagnostic web_search probe — runs the same DDG (+ Brave fallback) flow as
    the production agents but returns raw HTTP metrics in JSON instead of
    pre-digested search results. Use to diagnose IP-blocks, rate-limits, or
    regex failures without needing Railway runtime logs.

    Returns a JSON object with:
      - ddg.status, ddg.body_len, ddg.anomaly_signal, ddg.results_count
      - ddg.first_titles, ddg.user_agent_used, ddg.headers_subset
      - brave.attempted, brave.status, brave.results_count (if BRAVE_SEARCH_API_KEY set)
      - elapsed_ms, query
    """
    denied = _enforce(caller, "debug_web_search")
    if denied:
        return denied
    import httpx, re, urllib.parse, random, time

    out = {
        "query": query,
        "ddg": {},
        "brave": {"attempted": False},
        "elapsed_ms": 0,
    }
    t0 = time.time()

    # DDG attempt — same UA pool as production _web_search
    ddg_results_count = 0
    try:
        ua = random.choice(DDG_USER_AGENTS)
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            resp = await client.get(
                "https://html.duckduckgo.com/html/",
                params={"q": query},
                headers={"User-Agent": ua},
            )
        body = resp.text or ""
        body_lower = body.lower()
        anomaly = (
            ("anomaly" in body_lower)
            or ("captcha" in body_lower)
            or ("rate limit" in body_lower)
            or ("blocked" in body_lower and "duckduckgo" in body_lower)
        )
        links = re.findall(r'class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>', body, re.DOTALL)
        snippets = re.findall(r'class="result__snippet"[^>]*>(.*?)<', body, re.DOTALL)
        ddg_results_count = len(links)
        first_titles = []
        for raw_url, title in links[:3]:
            t_clean = re.sub(r'<[^>]+>', '', title).strip()
            first_titles.append(t_clean[:120])
        out["ddg"] = {
            "status": resp.status_code,
            "body_len": len(body),
            "user_agent_used": ua[:80],
            "anomaly_signal": anomaly,
            "results_count": ddg_results_count,
            "snippets_count": len(snippets),
            "first_titles": first_titles,
            "body_first_200": body[:200],
            "headers_subset": {
                "content-type": resp.headers.get("content-type", ""),
                "server": resp.headers.get("server", ""),
                "x-frame-options": resp.headers.get("x-frame-options", ""),
            },
        }
    except Exception as e:
        out["ddg"] = {"exception": f"{type(e).__name__}: {e}"}

    # Brave fallback if DDG failed and key configured
    brave_key = os.environ.get("BRAVE_SEARCH_API_KEY", "").strip()
    if not brave_key:
        out["brave"]["status_msg"] = "BRAVE_SEARCH_API_KEY not set; fallback disabled"
    elif ddg_results_count > 0:
        out["brave"]["status_msg"] = "DDG returned results; Brave not attempted"
    else:
        out["brave"]["attempted"] = True
        try:
            async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
                r = await client.get(
                    "https://api.search.brave.com/res/v1/web/search",
                    params={"q": query, "count": "8"},
                    headers={"X-Subscription-Token": brave_key, "Accept": "application/json"},
                )
            out["brave"]["status"] = r.status_code
            if r.status_code == 200:
                d = r.json()
                results = (d.get("web", {}).get("results") or [])[:8]
                out["brave"]["results_count"] = len(results)
                out["brave"]["first_titles"] = [it.get("title", "")[:120] for it in results[:3]]
            else:
                out["brave"]["error_body"] = (r.text or "")[:300]
        except Exception as e:
            out["brave"]["exception"] = f"{type(e).__name__}: {e}"

    out["elapsed_ms"] = int((time.time() - t0) * 1000)
    return json.dumps(out, ensure_ascii=False, indent=2)


SYNTHESIS_NO_TOOLS_DIRECTIVE = (
    "\n\n## SZINTÉZIS-FÁZIS — TOOLOK TILTVA\n"
    "Ez a koordinátori szintézis, NEM kutatási kör. NE adj ki tool_call-t "
    "(sem JSON, sem szöveges marker formában). Konkrétan TILOS a következő "
    "tokenek emit-elése: `<|tool_call`, `<｜DSML｜`, `function_calls>`, "
    "`tool_calls_section`, `web_search`. Ha valamit nem tudsz az al-agentek "
    "anyagából, mondd ki hogy hiányzik — ne próbálj keresni. Csak prózai, "
    "magyar nyelvű szintézist adj.\n"
)


def _clean_synthesis_output(content: str) -> str:
    """Strip leaked tool_call markup from a synthesis response.

    Last-line defense against models (esp. Kimi K2.6 in synthesis mode)
    that emit text-format tool_calls instead of plain prose. Removes both
    the Kimi-style and DeepSeek-style markup blocks. If after cleaning the
    content is empty, returns a placeholder so the DB row isn't NULL.
    """
    import re as _re
    if not content:
        return "(üres szintézis)"
    cleaned = content
    # Kimi tool_calls section: <|tool_calls_section_begin|>...<|tool_calls_section_end|>
    cleaned = _re.sub(
        r'<\|tool_calls_section_begin\|>.*?<\|tool_calls_section_end\|>',
        '', cleaned, flags=_re.DOTALL,
    )
    # Individual tool_call wrappers
    cleaned = _re.sub(
        r'<\|tool_call_begin\|>.*?<\|tool_call_end\|>',
        '', cleaned, flags=_re.DOTALL,
    )
    # DeepSeek DSML-style invokes
    cleaned = _re.sub(
        r'<｜DSML｜tool_calls>.*?</｜DSML｜tool_calls>',
        '', cleaned, flags=_re.DOTALL,
    )
    cleaned = _re.sub(r'<｜DSML｜[^>]*>', '', cleaned)
    # Generic <function_calls>...</function_calls>
    cleaned = _re.sub(
        r'<function_calls>.*?</function_calls>',
        '', cleaned, flags=_re.DOTALL,
    )
    cleaned = cleaned.strip()
    if not cleaned:
        return "(a koordinátor csak tool_call markert adott — a Bridge eltávolította, érdemleges szintézis nincs)"
    return cleaned


DEEP_RESEARCH_DIRECTIVE = (
    "\n\n## DEEP RESEARCH MÓD — MULTI-ROUND KUTATÁS KÖTELEZŐ\n"
    "Ez NEM egy 1-shot kérdés. Több körön keresztül a `web_search` tool-t többször "
    "is használnod kell, hogy a témát ALAPOSAN kifedd:\n"
    "1) Az első körben tervezz 3-5 különböző web_search query-t (fókuszok: tényadat, "
    "vélemény, kontextus, ellenőrző-forrás).\n"
    "2) Olvasd a találatokat, és ha valamelyik fontos dimenzió hiányzik vagy "
    "ellentmondásos, indíts ÚJABB web_search-öt finomított query-vel.\n"
    "3) Cross-checkeld a kulcsállításokat 2+ független forrásból.\n"
    "4) MINDEN tényállításnál idézd a forrást — formátum: `[forrás N]` a mondat "
    "végén, ahol N a források listájában szereplő sorszám.\n"
    "5) A végső válasz STRUKTURÁLT legyen, és tartalmazzon egy `## Források` "
    "szekciót, sorszámozva, minden sorhoz a teljes URL.\n\n"
    "FONTOS: ha egy állítás nincs forrással alátámasztva, NE írd le. "
    "Inkább mondd hogy 'nem találtam erre megbízható forrást'.\n"
)


async def _deep_research_loop(
    model_id: str,
    messages: list,
    agent_extra: dict,
    max_tokens: int = 8000,
    temperature: float = 0.6,
    max_rounds: int = 4,
    task_budget_s: float = 240.0,
) -> str:
    """Multi-round web_search loop with mandatory source citation.

    Drives the model through several search-and-refine rounds, executes
    every web_search tool_call (both JSON and text-marker formats), feeds
    the results back into the conversation, and forces a final synthesis
    round (tools off) where the system prompt requires explicit `[forrás N]`
    citations and a `## Források` URL list.

    Robustness:
      - max_rounds=4 (was 6) — fewer query-bursts to DDG.
      - task_budget_s=240 — total wall-time cap. If exceeded mid-round we
        fast-forward to a forced synthesis with whatever sources we have.
      - Per-round httpx.TimeoutException is caught and treated as
        "force synthesis" instead of bubbling up and crashing the task.

    Returns the final assistant content.
    """
    import httpx, re, time as _t

    sources: list[str] = []  # collected URLs (for the final citation block)
    last_content = ""        # last non-empty model reply (for the budget-exhaust path)
    start_t = _t.time()

    async def _api_call(client, payload):
        resp = await client.post(
            f"{SILICONFLOW_BASE_URL}/chat/completions",
            headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
            json=payload,
        )
        return json.loads(resp.text)

    for round_num in range(max_rounds):
        # Budget guard — if we've burned the budget, skip to forced synthesis
        elapsed = _t.time() - start_t
        if elapsed > task_budget_s:
            logger.warning("Deep research task_budget exceeded at round %d (%.1fs) — forcing synthesis", round_num, elapsed)
            is_final = True
        else:
            is_final = (round_num == max_rounds - 1)
        payload = {
            "model": model_id,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_tokens,
            **agent_extra,
        }
        if not is_final:
            payload["tools"] = [WEB_SEARCH_TOOL_DEF]
        else:
            # Final round: force synthesis with citations.
            messages = messages + [{
                "role": "user",
                "content": (
                    "Most foglald össze az eddig összegyűjtött anyagot. "
                    "Strukturált, tényszerű választ adj, MINDEN állításhoz idézd "
                    "a forrást [forrás N] formátumban, és a végén adj egy "
                    "`## Források` szekciót sorszámozott URL-listával. "
                    "Ha valamire nincs forrás, mondd ki hogy nincs."
                ),
            }]
            payload["messages"] = messages

        try:
            async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                data = await _api_call(client, payload)
        except (httpx.TimeoutException, httpx.ConnectError) as e:
            logger.error("Deep research round %d httpx error: %s — treating as soft fail, continuing", round_num, e)
            # Don't crash the whole task; if we have enough sources let the next
            # round force-synthesize, otherwise return what we have.
            if is_final or sources:
                break
            continue

        if not isinstance(data, dict) or not data.get("choices"):
            logger.error("Deep research round %d bad response: %s", round_num, str(data)[:200])
            if is_final or sources:
                break
            continue

        choice = data["choices"][0]
        msg = choice.get("message", {})
        content = msg.get("content", "") or ""
        if content.strip():
            last_content = content
        if not content.strip():
            content = msg.get("reasoning_content", "") or ""
        tool_calls = msg.get("tool_calls")
        finish_reason = choice.get("finish_reason", "")

        # Case A — proper JSON tool_calls
        if tool_calls and finish_reason == "tool_calls":
            messages = messages + [msg]
            for tc in tool_calls:
                fn = tc.get("function", {})
                if fn.get("name") != "web_search":
                    continue
                args_raw = fn.get("arguments", "{}")
                try:
                    args = json.loads(args_raw) if isinstance(args_raw, str) else args_raw
                except Exception:
                    args = {}
                q = args.get("query", "")
                if not q:
                    continue
                logger.info("Deep research round %d query: %s", round_num, q[:80])
                sr = await _web_search(q)
                # Capture URLs for the citation footer fallback
                for url in re.findall(r'https?://[^\s\)]+', sr):
                    if url not in sources:
                        sources.append(url)
                messages.append({
                    "role": "tool",
                    "tool_call_id": tc.get("id", ""),
                    "content": sr[:4000],
                })
            continue

        # Case B — text-marker tool calls (Kimi / DeepSeek DSML / GLM5)
        text_markers = ["<|tool_call", "<｜DSML｜", "function_calls>", "tool_calls_section"]
        if not tool_calls and any(m in content for m in text_markers):
            queries = re.findall(r'"query"[:\s]*"([^"]+)"', content)[:3]
            if queries:
                aggregated = ""
                for q in queries:
                    logger.info("Deep research round %d (text) query: %s", round_num, q[:80])
                    sr = await _web_search(q)
                    aggregated += f"\n[Web search: {q}]\n{sr}\n"
                    for url in re.findall(r'https?://[^\s\)]+', sr):
                        if url not in sources:
                            sources.append(url)
                # Wrap text-marker into the next user turn so the model can use the data
                messages = messages + [
                    {"role": "assistant", "content": "(web_search calls in progress)"},
                    {"role": "user", "content": f"Web keresési eredmények:\n{aggregated}\n\nFolytasd a kutatást vagy szintetizálj."},
                ]
                continue

        # No more tool calls — model wants to finalize. Return whatever we got.
        if content.strip():
            # Append a fallback citation footer if the model forgot the URL list
            if sources and "## Források" not in content and "## Forrasok" not in content:
                content += "\n\n## Források\n" + "\n".join(f"[{i+1}] {u}" for i, u in enumerate(sources[:15]))
            return content

    # Loop exhausted (or soft-failed). Salvage whatever we can: the last
    # non-empty content the model produced, plus the source list we collected.
    if last_content.strip():
        fallback = last_content.strip()
        if sources and "## Források" not in fallback and "## Forrasok" not in fallback:
            fallback += "\n\n## Források\n" + "\n".join(f"[{i+1}] {u}" for i, u in enumerate(sources[:15]))
        return fallback
    fallback = "(Deep research loop nem ért el végleges szintézist a megengedett körökön belül.)"
    if sources:
        fallback += "\n\n## Források\n" + "\n".join(f"[{i+1}] {u}" for i, u in enumerate(sources[:15]))
    return fallback


async def _run_agent_with_tools(model_id: str, messages: list, max_rounds: int = 4) -> str:
    """Run an AI agent with optional tool calls (web_search). Returns final text."""
    import httpx
    for round_num in range(max_rounds):
        # Last round: no tools, force text response
        use_tools = round_num < max_rounds - 1
        payload = {
            "model": model_id,
            "messages": messages,
            "temperature": 0.6,
            "max_tokens": 2000,
        }
        # Only synthesis call path — model_id is always K2.6 here. Force thinking OFF.
        if "Kimi" in model_id:
            payload["thinking"] = {"type": "disabled"}
        if use_tools:
            payload["tools"] = [WEB_SEARCH_TOOL_DEF]

        async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
            resp = await client.post(
                f"{SILICONFLOW_BASE_URL}/chat/completions",
                headers={
                    "Authorization": f"Bearer {SILICONFLOW_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
        data = json.loads(resp.text)
        if not isinstance(data, dict) or "error" in data:
            return f"API error: {data}"

        choice = data.get("choices", [{}])[0]
        msg = choice.get("message", {})

        # Check for tool calls (JSON format)
        tool_calls = msg.get("tool_calls")
        content = msg.get("content", "") or ""

        # Case B: Text-based tool calls (Kimi sometimes does this)
        import re
        text_markers = ["<|tool_call", "<｜DSML｜", "function_calls>", "tool_calls_section"]
        if not tool_calls and any(m in content for m in text_markers):
            queries = re.findall(r'"query"[:\s]*"([^"]+)"', content)
            if queries:
                search_results = ""
                for query in queries[:2]:
                    sr = await _web_search(query)
                    search_results += f"\n[Web search: {query}]\n{sr}\n"
                    logger.info("AI web_search (text-parsed) round %d: %s", round_num, query[:80])
                # Re-call with search results as context
                messages.append({"role": "assistant", "content": "(web keresés végrehajtva)"})
                messages.append({"role": "user", "content": f"Web keresési eredmények:\n{search_results}\n\nVálaszolj az eredmények alapján."})
                continue  # Next round will generate text response

        # If no tool calls at all, return the text
        if not tool_calls:
            return content

        # Execute JSON tool calls
        messages.append(msg)
        for tc in tool_calls:
            fn = tc.get("function", {})
            if fn.get("name") == "web_search":
                args = json.loads(fn.get("arguments", "{}"))
                result = await _web_search(args.get("query", ""))
                messages.append({
                    "role": "tool",
                    "tool_call_id": tc["id"],
                    "content": result,
                })
                logger.info("AI web_search round %d: %s", round_num, args.get("query", "")[:80])

    # Should not reach here (last round forces text), but just in case
    return msg.get("content", "") or "(az agent nem adott választ a web keresés után)"


async def _execute_ai_task(task_id: int, title: str, description: str, context: str, assigned_by: str,
                            deep_research: bool = False, deep_thinking: bool = False,
                            output_format: str = ""):
    """Background execution of a multi-agent AI task. Agents run in PARALLEL."""
    conn = get_db()
    conn.execute("UPDATE ai_tasks SET status = 'running' WHERE id = ?", (task_id,))
    conn.commit()

    roles = {
        "kimi": ("moonshotai/Kimi-K2.6", "Kutató és elemző. Alapos, részletes munkát végzel.", 20000),
        "deepseek": ("deepseek-ai/DeepSeek-V4-Pro", "Kritikus elemző és ellenőr. Logikai hibákat keresel, ellenérveket fogalmazol.", 20000),
        "glm5": ("zai-org/GLM-5.1", "Végrehajtó és kóder. Konkrét megoldásokat, kódot, strukturált outputot adsz. Ha kell, implementálsz.", 20000),
    }

    task_prompt = f"FELADAT: {title}\n\nLEÍRÁS: {description}"
    if context:
        task_prompt += f"\n\nKONTEXTUS:\n{context}"

    async def _run_single_agent(agent_name, model_id, role_desc, agent_max_tokens=3000):
        """Run one agent with hybrid tool-call support — called in parallel. Retries once on timeout."""
        import httpx, re

        async def _api_call(client, payload):
            resp = await client.post(
                f"{SILICONFLOW_BASE_URL}/chat/completions",
                headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                json=payload,
            )
            return json.loads(resp.text)

        for attempt in range(2):
            try:
                temporal = _temporal_directive(agent_name)
                # Pyramid context if available, otherwise fallback
                if PYRAMID_ENABLED and agent_name in PYRAMID_AGENTS:
                    system = build_agent_context(
                        agent_id=agent_name,
                        custom_system_prompt=temporal,
                        inbox_summary=_get_inbox_summary(),
                        relevance_query=task_prompt,
                    )
                else:
                    system = (
                        f"Te a Claus multi-agent rendszer '{agent_name}' al-agentje vagy. "
                        f"Szereped: {role_desc} "
                        f"Magyarul válaszolj, lényegre törően de alaposan."
                        + temporal
                    )

                # Inject caller persona if available
                if assigned_by and not is_core_instance(assigned_by):
                    profile = get_profile(assigned_by)
                    if profile and profile.persona_system_prompt:
                        system += f"\n\n--- HÍVÓ FÉL ---\n{profile.persona_system_prompt}"

                # If output_format set, append the structured-output directive so
                # the agent writes markdown the Bridge can render to a real file.
                if output_format and output_format in OUTPUT_FORMAT_DIRECTIVES:
                    system += OUTPUT_FORMAT_DIRECTIVES[output_format]

                # Kimi K2.6 default thinking is ON on SF — force OFF to keep latency sane.
                # DeepSeek V4-Pro default thinking ON too — clamp to medium effort.
                # deep_thinking=True overrides: thinking ON for Kimi, reasoning=high for DeepSeek.
                if deep_thinking:
                    if agent_name == "kimi":
                        agent_extra = {"thinking": {"type": "enabled"}}
                    elif agent_name == "deepseek":
                        agent_extra = {"reasoning_effort": "high"}
                    else:
                        agent_extra = {}
                else:
                    if agent_name == "kimi":
                        agent_extra = {"thinking": {"type": "disabled"}}
                    elif agent_name == "deepseek":
                        agent_extra = {"reasoning_effort": "medium"}
                    else:
                        agent_extra = {}

                # Deep research path: hand off to the multi-round loop and skip the
                # legacy 3-step (single search + retry) flow entirely.
                if deep_research:
                    research_system = system + DEEP_RESEARCH_DIRECTIVE
                    content = await _deep_research_loop(
                        model_id=model_id,
                        messages=[
                            {"role": "system", "content": research_system},
                            {"role": "user", "content": task_prompt},
                        ],
                        agent_extra=agent_extra,
                        max_tokens=agent_max_tokens,
                        temperature=0.6,
                        max_rounds=4,
                    )
                    ts = now()
                    conn.execute(
                        "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                        (task_id, agent_name, role_desc, content or "(no response)", ts)
                    )
                    conn.commit()
                    _maybe_render_output_file(task_id, agent_name, content, output_format, title)
                    if PYRAMID_ENABLED and agent_name in PYRAMID_AGENTS and content:
                        try:
                            pyramid_store_result(content=content, agent_id=agent_name, task_title=title, force_shared=True)
                        except Exception:
                            pass
                    logger.info("AI task #%d: %s deep_research done (%d chars)", task_id, agent_name, len(content or ""))
                    return

                # Step 1: Call WITH tools
                async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                    data = await _api_call(client, {
                        "model": model_id,
                        "messages": [
                            {"role": "system", "content": system},
                            {"role": "user", "content": task_prompt},
                        ],
                        "temperature": 0.6,
                        "max_tokens": agent_max_tokens,
                        "tools": [WEB_SEARCH_TOOL_DEF],
                        **agent_extra,
                    })

                content = ""
                search_results = ""
                if not isinstance(data, dict) or not data.get("choices"):
                    logger.error("AI task #%d %s: bad API response", task_id, agent_name)
                else:
                    msg = data["choices"][0].get("message", {})
                    content = msg.get("content", "") or ""
                    tool_calls = msg.get("tool_calls")

                    # Case A: Proper JSON tool_calls — execute web search
                    if tool_calls:
                        for tc in tool_calls:
                            fn = tc.get("function", {})
                            if fn.get("name") == "web_search":
                                query = json.loads(fn.get("arguments", "{}")).get("query", "")
                                if query:
                                    sr = await _web_search(query)
                                    search_results += f"\n[Web search: {query}]\n{sr}\n"
                                    logger.info("AI task #%d %s: web_search '%s'", task_id, agent_name, query[:60])

                    # Case B: Text-based tool calls — parse search query from text
                    tool_markers = ["<|tool_call", "<｜DSML｜", "function_calls>", "tool_calls_section"]
                    if not tool_calls and any(m in content for m in tool_markers):
                        queries = re.findall(r'"query"[:\s]*"([^"]+)"', content)
                        for query in queries[:2]:
                            sr = await _web_search(query)
                            search_results += f"\n[Web search: {query}]\n{sr}\n"
                            logger.info("AI task #%d %s: parsed web_search '%s'", task_id, agent_name, query[:60])
                        content = ""  # Clear broken text

                # Step 2: If we got search results, call again WITH those results as context
                if search_results:
                    async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                        data2 = await _api_call(client, {
                            "model": model_id,
                            "messages": [
                                {"role": "system", "content": f"{role_desc} Magyarul válaszolj, alaposan. Az alábbi web keresési eredményeket használd fel."},
                                {"role": "user", "content": f"{task_prompt}\n\nWEB KERESÉSI EREDMÉNYEK:\n{search_results}"},
                            ],
                            "temperature": 0.6,
                            "max_tokens": agent_max_tokens,
                            **agent_extra,
                        })
                        if isinstance(data2, dict) and data2.get("choices"):
                            content = data2["choices"][0].get("message", {}).get("content", "")

                # Step 3: Fallback if still empty
                if not content or not content.strip():
                    logger.warning("AI task #%d %s: empty after tools, final fallback", task_id, agent_name)
                    async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                        data3 = await _api_call(client, {
                            "model": model_id,
                            "messages": [
                                {"role": "system", "content": f"{role_desc} Magyarul válaszolj, részletesen. NE használj tool-okat."},
                                {"role": "user", "content": task_prompt},
                            ],
                            "temperature": 0.7,
                            "max_tokens": agent_max_tokens,
                        })
                        if isinstance(data3, dict) and data3.get("choices"):
                            content = data3["choices"][0].get("message", {}).get("content", "")

                ts = now()
                conn.execute(
                    "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                    (task_id, agent_name, role_desc, content or "(no response)", ts)
                )
                conn.commit()
                # Pyramid governance: store in RAG / shared memory
                if PYRAMID_ENABLED and agent_name in PYRAMID_AGENTS and content:
                    try:
                        pyramid_store_result(content=content, agent_id=agent_name, task_title=title, force_shared=True)
                    except Exception:
                        pass
                _maybe_render_output_file(task_id, agent_name, content, output_format, title)
                logger.info("AI task #%d: %s done (%d chars)", task_id, agent_name, len(content or ""))
                return  # Success — exit retry loop

            except (httpx.TimeoutException, httpx.ConnectError) as e:
                if attempt == 0:
                    logger.warning("AI task #%d %s: timeout (attempt 1), retrying in 3s: %s", task_id, agent_name, e)
                    await asyncio.sleep(3)
                else:
                    ts = now()
                    conn.execute(
                        "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                        (task_id, agent_name, role_desc, f"TIMEOUT after retry: {e}", ts)
                    )
                    conn.commit()
                    logger.error("AI task #%d %s: timeout after retry: %s", task_id, agent_name, e)
            except Exception as e:
                ts = now()
                conn.execute(
                    "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                    (task_id, agent_name, role_desc, f"ERROR: {e}", ts)
                )
                conn.commit()
                logger.error("AI task #%d %s failed: %s", task_id, agent_name, e)
                return  # Non-timeout errors don't retry

    # Run all agents IN PARALLEL
    await asyncio.gather(
        *[_run_single_agent(name, mid, rdesc, mt) for name, (mid, rdesc, mt) in roles.items()]
    )

    # Synthesis by Kimi
    try:
        results = conn.execute(
            "SELECT agent, content FROM ai_task_results WHERE task_id = ? ORDER BY id", (task_id,)
        ).fetchall()
        parts = "\n\n".join(f"[{r['agent']}]:\n{r['content']}" for r in results)
        system = (
            "Te a koordinátor vagy. Az al-agentek elvégezték a feladatot. "
            "Készíts tömör szintézist az eredményeikből: mi az egyetértés, hol térnek el, és mi a végső ajánlás. "
            "Magyarul, strukturáltan."
            + _temporal_directive("kimi")
            + SYNTHESIS_NO_TOOLS_DIRECTIVE
        )
        messages = [
            {"role": "system", "content": system},
            {"role": "user", "content": f"FELADAT: {title}\n\nAGENT EREDMÉNYEK:\n{parts}"},
        ]
        synthesis = await _run_agent_with_tools("moonshotai/Kimi-K2.6", messages)
        synthesis = _clean_synthesis_output(synthesis)
        ts = now()
        conn.execute(
            "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
            (task_id, "szintézis", "Koordinátori összefoglaló", synthesis, ts)
        )
        conn.commit()
    except Exception as e:
        logger.error("AI task #%d synthesis failed: %s", task_id, e)

    conn.execute("UPDATE ai_tasks SET status = 'completed', completed_at = ? WHERE id = ?", (now(), task_id))
    conn.commit()
    conn.close()
    logger.info("AI task #%d completed", task_id)


@mcp.tool()
async def ai_task(title: str, description: str, context: str = "", file_id: int = 0, assigned_by: str = "unknown",
                  agent_tasks: str = "", deep_research: bool = False, deep_thinking: bool = False,
                  output_format: str = "") -> str:
    """Create and execute a multi-agent AI task. Kimi, DeepSeek, and GLM-5.1 work on it, then a synthesis is generated.

    Two modes:
    1. **Broadcast** (default): All agents get the same task.
    2. **Dispatch** (agent_tasks): Each agent gets a DIFFERENT task — parallel execution.

    Args:
        title: Short task title
        description: Detailed task description / instructions
        context: Optional document text or background context
        file_id: Optional uploaded file ID (from upload_file) to include as context
        assigned_by: Who created the task (cli-claus, web-claus, kommandant)
        agent_tasks: Optional JSON — per-agent tasks for parallel dispatch mode. Format:
            {"kimi": {"prompt": "...", "max_tokens": 3000}, "deepseek": {"prompt": "..."}}
        deep_research: Multi-round web_search loop with mandatory `[forrás N]`
            citations and a final `## Források` URL list. Use for press review,
            equity briefs, fact-checking — anywhere 1-shot DDG is too thin.
            Slower (~60-180s per agent), real DDG calls.
        deep_thinking: Enable explicit reasoning mode — Kimi `thinking=enabled`,
            DeepSeek `reasoning_effort=high`. Default is OFF/medium for latency.
            Turn on for hard logic, multi-step deduction, mathematical proofs.
            ~3-5x slower per agent. Combinable with deep_research (very slow,
            but the highest-quality output the rizsrakéták can produce).
        output_format: Empty for normal text response (default), or
            "docx" / "xlsx" / "pptx". Each agent then writes its content
            in a structured markdown the Bridge renders into a real Office
            file, attached to the task results as a downloadable upload.
            Per-agent file: each agent gets its own file_id.
    """
    # Pull in uploaded file content if file_id provided
    if file_id:
        conn = get_db()
        f = conn.execute("SELECT filename, mime_type, content_text, content_base64 FROM uploads WHERE id = ?", (file_id,)).fetchone()
        conn.close()
        if f:
            if f["content_text"]:
                context = f"[Feltöltött fájl: {f['filename']}]\n\n{f['content_text']}\n\n{context}"
            elif f["content_base64"]:
                context = f"[Feltöltött kép: {f['filename']} — base64 kép az agentek számára elérhető]\n\n{context}"
    if not SILICONFLOW_API_KEY:
        return json.dumps({"error": "SILICONFLOW_API_KEY not set"})

    # Parse agent_tasks JSON if provided
    parsed_agent_tasks = None
    if agent_tasks:
        try:
            parsed_agent_tasks = json.loads(agent_tasks) if isinstance(agent_tasks, str) else agent_tasks
        except (json.JSONDecodeError, TypeError):
            return json.dumps({"error": "agent_tasks must be valid JSON"})

    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO ai_tasks (title, description, context, assigned_by, status, created_at) VALUES (?, ?, ?, ?, 'pending', ?)",
        (title, description, context[:10000], assigned_by, ts)
    )
    task_id = cur.lastrowid
    conn.commit()
    conn.close()

    if parsed_agent_tasks and PYRAMID_ENABLED:
        # DISPATCH MODE: each agent gets a different task
        from pyramid.task_dispatcher import dispatch_parallel_tasks

        async def _call_agent(model, prompt, system_prompt, max_tokens, temperature):
            """Wrapper with web search — handles both JSON and text-based tool calls."""
            import httpx, re
            model_id = SILICONFLOW_MODELS.get(model, model)
            # Inject the strong temporal directive — Kimi otherwise overrides the
            # runtime date with its training cutoff, poisoning the dispatch output.
            system_prompt = (system_prompt or "") + _temporal_directive(model)
            # Append output_format directive if requested
            if output_format and output_format in OUTPUT_FORMAT_DIRECTIVES:
                system_prompt += OUTPUT_FORMAT_DIRECTIVES[output_format]

            # K2.6 defaults thinking ON on SF — force OFF (latency unacceptable on dispatch path).
            # V4-Pro defaults thinking ON too — clamp to reasoning_effort=medium.
            # deep_thinking flips both back to maximum-reasoning mode.
            if deep_thinking:
                if model == "kimi":
                    agent_extra = {"thinking": {"type": "enabled"}}
                elif model == "deepseek":
                    agent_extra = {"reasoning_effort": "high"}
                else:
                    agent_extra = {}
            else:
                if model == "kimi":
                    agent_extra = {"thinking": {"type": "disabled"}}
                elif model == "deepseek":
                    agent_extra = {"reasoning_effort": "medium"}
                else:
                    agent_extra = {}

            # Deep research path — multi-round loop, returns final content.
            if deep_research:
                research_system = system_prompt + DEEP_RESEARCH_DIRECTIVE
                content = await _deep_research_loop(
                    model_id=model_id,
                    messages=[
                        {"role": "system", "content": research_system},
                        {"role": "user", "content": prompt},
                    ],
                    agent_extra=agent_extra,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    max_rounds=4,
                )
                return {"response": content, "tokens": {"prompt": 0, "completion": 0}}

            async def _api(client, payload):
                resp = await client.post(
                    f"{SILICONFLOW_BASE_URL}/chat/completions",
                    headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                    json=payload,
                )
                return json.loads(resp.text)

            # Step 1: Call with tools
            async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                data = await _api(client, {
                    "model": model_id,
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": prompt},
                    ],
                    "temperature": temperature, "max_tokens": max_tokens,
                    "tools": [WEB_SEARCH_TOOL_DEF],
                    **agent_extra,
                })

            if not isinstance(data, dict) or not data.get("choices"):
                return {"response": f"API error: {data}", "tokens": {"prompt": 0, "completion": 0}}

            msg = data["choices"][0].get("message", {})
            content = msg.get("content", "") or ""
            tool_calls = msg.get("tool_calls")
            search_results = ""

            # Case A: Proper JSON tool_calls
            if tool_calls:
                for tc in tool_calls:
                    fn = tc.get("function", {})
                    if fn.get("name") == "web_search":
                        query = json.loads(fn.get("arguments", "{}")).get("query", "")
                        if query:
                            sr = await _web_search(query)
                            search_results += f"\n[Web search: {query}]\n{sr}\n"
                            logger.info("Dispatch %s: web_search '%s'", model, query[:60])

            # Case B: Text-based tool calls (Kimi does this)
            text_markers = ["<|tool_call", "<｜DSML｜", "function_calls>", "tool_calls_section"]
            if not tool_calls and any(m in content for m in text_markers):
                queries = re.findall(r'"query"[:\s]*"([^"]+)"', content)
                for query in queries[:3]:
                    sr = await _web_search(query)
                    search_results += f"\n[Web search: {query}]\n{sr}\n"
                    logger.info("Dispatch %s: parsed web_search '%s'", model, query[:60])
                content = ""  # Clear broken marker text

            # Step 2: If we got search results, call again WITH results
            if search_results:
                async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                    data2 = await _api(client, {
                        "model": model_id,
                        "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": f"{prompt}\n\nWEB KERESÉSI EREDMÉNYEK:\n{search_results}"},
                        ],
                        "temperature": temperature, "max_tokens": max_tokens,
                        **agent_extra,
                    })
                if isinstance(data2, dict) and data2.get("choices"):
                    content = data2["choices"][0].get("message", {}).get("content", "") or content

            return {"response": content, "tokens": {"prompt": 0, "completion": 0}}

        def _run_dispatch():
            loop = asyncio.new_event_loop()
            try:
                results = loop.run_until_complete(dispatch_parallel_tasks(
                    agent_tasks=parsed_agent_tasks,
                    task_title=title,
                    call_agent_func=_call_agent,
                ))
                # Store results in DB
                conn2 = get_db()
                for agent_id, result in results.items():
                    content = result.get("response", "(no response)") if isinstance(result, dict) else str(result)
                    conn2.execute(
                        "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                        (task_id, agent_id, "Pyramid dispatch", content, now())
                    )
                conn2.commit()
                # Per-agent file rendering when output_format requested
                if output_format and output_format in OUTPUT_FORMAT_DIRECTIVES:
                    for agent_id, result in results.items():
                        c = result.get("response", "") if isinstance(result, dict) else str(result)
                        if c and not c.startswith("ERROR:"):
                            _maybe_render_output_file(task_id, agent_id, c, output_format, title)

                # Synthesis — Kimi K2.6 combines per-agent outputs into a final brief.
                # Each agent received a DIFFERENT prompt in dispatch mode, so the synthesis
                # has to acknowledge the divergent fókuszok, not pretend agreement.
                try:
                    parts = []
                    for agent_id, result in results.items():
                        agent_prompt = parsed_agent_tasks.get(agent_id, {}).get("prompt", "")
                        agent_resp = result.get("response", "") if isinstance(result, dict) else str(result)
                        parts.append(f"=== [{agent_id}] FELADAT ===\n{agent_prompt}\n\n=== [{agent_id}] EREDMÉNY ===\n{agent_resp}")
                    synthesis_input = "\n\n---\n\n".join(parts)
                    synthesis_messages = [
                        {"role": "system", "content": (
                            "Te a koordinátor vagy. Az al-agentek ELTÉRŐ feladatokat kaptak, mindegyik a saját fókuszával dolgozott. "
                            "Készíts strukturált szintézist: 1) röviden mit fedett le mindegyik agent, 2) hol egészítik ki egymást, "
                            "3) hol mondanak ellent (ha igen), 4) végső konklúzió / ajánlás a Kommandantnak. "
                            "Magyarul, lényegre törően. Ne ismételd meg az agentek szövegét hosszan, csak hivatkozz rájuk."
                            + _temporal_directive("kimi")
                            + SYNTHESIS_NO_TOOLS_DIRECTIVE
                        )},
                        {"role": "user", "content": f"FELADAT CÍME: {title}\n\nLEÍRÁS: {description}\n\nAGENT EREDMÉNYEK:\n{synthesis_input}"},
                    ]
                    synthesis = loop.run_until_complete(
                        _run_agent_with_tools("moonshotai/Kimi-K2.6", synthesis_messages, max_rounds=2)
                    )
                    synthesis = _clean_synthesis_output(synthesis)
                    if synthesis and synthesis.strip():
                        conn2.execute(
                            "INSERT INTO ai_task_results (task_id, agent, role, content, timestamp) VALUES (?, ?, ?, ?, ?)",
                            (task_id, "szintézis", "Koordinátori összefoglaló", synthesis, now())
                        )
                        conn2.commit()
                        logger.info("AI task #%d synthesis written (%d chars)", task_id, len(synthesis))
                except Exception as syn_e:
                    logger.error("AI task #%d dispatch synthesis failed: %s", task_id, syn_e)

                conn2.execute("UPDATE ai_tasks SET status = 'completed', completed_at = ? WHERE id = ?", (now(), task_id))
                conn2.commit()
                conn2.close()
                logger.info("AI task #%d dispatch completed: %s", task_id, list(results.keys()))
            finally:
                loop.close()

        threading.Thread(target=_run_dispatch, daemon=True).start()
        agents = list(parsed_agent_tasks.keys())
        return json.dumps({"status": "task_dispatched", "task_id": task_id, "mode": "pyramid_dispatch", "agents": agents})
    else:
        # BROADCAST MODE: all agents get the same task (legacy)
        def _run():
            loop = asyncio.new_event_loop()
            loop.run_until_complete(_execute_ai_task(
                task_id, title, description, context[:10000], assigned_by,
                deep_research=deep_research, deep_thinking=deep_thinking,
                output_format=output_format,
            ))
            loop.close()
        threading.Thread(target=_run, daemon=True).start()

        return json.dumps({"status": "task_created", "task_id": task_id, "message": "Kimi + DeepSeek + GLM-5.1 dolgoznak rajta. Eredmény a dashboardon."})


@mcp.custom_route("/api/upload", methods=["POST"])
async def api_upload(request):
    body = await request.json()
    filename = body.get("filename", "unknown")
    content_base64 = body.get("content_base64", "")
    mime_type = body.get("mime_type", "")
    uploaded_by = body.get("uploaded_by", "kommandant")

    try:
        file_bytes = base64.b64decode(content_base64)
    except Exception:
        return JSONResponse({"error": "Invalid base64"}, status_code=400)

    ext = pathlib.Path(filename).suffix.lower()
    if not mime_type:
        mime_map = {
            ".pdf": "application/pdf", ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".txt": "text/plain", ".md": "text/markdown", ".csv": "text/csv",
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".webp": "image/webp",
        }
        mime_type = mime_map.get(ext, "application/octet-stream")

    filepath = UPLOAD_DIR / f"{int(time.time())}_{filename}"
    filepath.write_bytes(file_bytes)

    content_text = _parse_file_to_text(filepath, mime_type) if not _is_image(mime_type) else ""
    b64_for_db = content_base64 if _is_image(mime_type) else ""

    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO uploads (filename, mime_type, content_text, content_base64, uploaded_by, uploaded_at) VALUES (?, ?, ?, ?, ?, ?)",
        (filename, mime_type, content_text[:50000], b64_for_db[:500000], uploaded_by, ts)
    )
    file_id = cur.lastrowid
    conn.commit()
    conn.close()
    return JSONResponse({"status": "uploaded", "file_id": file_id, "filename": filename, "text_length": len(content_text)})


@mcp.custom_route("/api/send_email", methods=["POST"])
async def api_send_email(request):
    """Send email with optional file attachment via multipart form upload.

    Form fields: to, subject, body, body_type (optional, default: plain)
    File field: file (optional attachment)
    Or: file_id (use previously uploaded file from uploads table)
    """
    content_type = request.headers.get("content-type", "")

    if "multipart/form-data" in content_type:
        form = await request.form()
        to = form.get("to", "")
        subject = form.get("subject", "")
        body = form.get("body", "")
        body_type = form.get("body_type", "plain")
        file_id = int(form.get("file_id", "0") or "0")
        uploaded_file = form.get("file")
    else:
        data = await request.json()
        to = data.get("to", "")
        subject = data.get("subject", "")
        body = data.get("body", "")
        body_type = data.get("body_type", "plain")
        file_id = int(data.get("file_id", 0) or 0)
        uploaded_file = None

    if not to or not subject:
        return JSONResponse({"error": "to and subject are required"}, status_code=400)

    svc = _capture_state.get("gmail_service")
    if not svc:
        return JSONResponse({"error": "Gmail not initialized"}, status_code=503)

    try:
        att_data = None
        att_name = None
        att_mime = "application/octet-stream"

        # Option 1: multipart file upload
        if uploaded_file and hasattr(uploaded_file, "read"):
            att_data = await uploaded_file.read()
            att_name = uploaded_file.filename
            att_mime = uploaded_file.content_type or att_mime
        # Option 2: file_id from uploads table
        elif file_id:
            conn = get_db()
            f = conn.execute("SELECT filename, mime_type, content_base64 FROM uploads WHERE id = ?", (file_id,)).fetchone()
            conn.close()
            if not f or not f["content_base64"]:
                return JSONResponse({"error": f"Upload #{file_id} not found or has no binary content"}, status_code=404)
            att_data = base64.b64decode(f["content_base64"])
            att_name = f["filename"]
            att_mime = f["mime_type"] or att_mime

        if att_data:
            msg = MIMEMultipart()
            msg["to"] = to
            msg["subject"] = subject
            msg.attach(MIMEText(body, body_type))
            maintype, _, subtype = att_mime.partition("/")
            part = MIMEBase(maintype or "application", subtype or "octet-stream")
            part.set_payload(att_data)
            encoders.encode_base64(part)
            part.add_header("Content-Disposition", "attachment", filename=att_name)
            msg.attach(part)
        else:
            msg = MIMEText(body, body_type)
            msg["to"] = to
            msg["subject"] = subject

        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        result = svc.users().messages().send(userId="me", body={"raw": raw}).execute()
        return JSONResponse({"status": "sent", "message_id": result.get("id"), "attachment": att_name})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@mcp.custom_route("/api/ai_tasks", methods=["GET", "POST"])
async def api_ai_tasks(request):
    if request.method == "POST":
        body = await request.json()
        title = body.get("title", "")
        description = body.get("description", "")
        context = body.get("context", "")
        file_id = body.get("file_id", 0)
        assigned_by = body.get("assigned_by", "kommandant")
        deep_research = bool(body.get("deep_research", False))
        deep_thinking = bool(body.get("deep_thinking", False))
        output_format = (body.get("output_format") or "").strip().lower()
        if output_format not in ("", "docx", "xlsx", "pptx"):
            output_format = ""

        # Pull in uploaded file if file_id given
        if file_id:
            conn2 = get_db()
            f = conn2.execute("SELECT filename, content_text FROM uploads WHERE id = ?", (file_id,)).fetchone()
            conn2.close()
            if f and f["content_text"]:
                context = f"[Feltöltött fájl: {f['filename']}]\n\n{f['content_text']}\n\n{context}"

        conn = get_db()
        ts = now()
        cur = conn.execute(
            "INSERT INTO ai_tasks (title, description, context, assigned_by, status, created_at) VALUES (?, ?, ?, ?, 'pending', ?)",
            (title, description, context[:10000], assigned_by, ts)
        )
        task_id = cur.lastrowid
        conn.commit()
        conn.close()

        # Execute in background
        def _run():
            loop = asyncio.new_event_loop()
            loop.run_until_complete(_execute_ai_task(
                task_id, title, description, context[:10000], assigned_by,
                deep_research=deep_research, deep_thinking=deep_thinking,
                output_format=output_format,
            ))
            loop.close()
        threading.Thread(target=_run, daemon=True).start()

        return JSONResponse({
            "status": "created", "task_id": task_id,
            "deep_research": deep_research, "deep_thinking": deep_thinking,
            "output_format": output_format or None,
        })

    conn = get_db()
    tasks = conn.execute("SELECT * FROM ai_tasks ORDER BY id DESC LIMIT 20").fetchall()
    result = []
    for t in tasks:
        results = conn.execute(
            "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id",
            (t["id"],)
        ).fetchall()
        result.append({**dict(t), "results": [dict(r) for r in results]})
    conn.close()
    return JSONResponse(result)


@mcp.custom_route("/api/ai_tasks/{task_id}/export", methods=["GET"])
async def api_ai_task_export(request):
    """Export AI task results as a .docx document."""
    task_id = request.path_params["task_id"]
    conn = get_db()
    task = conn.execute("SELECT * FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
    if not task:
        conn.close()
        return JSONResponse({"error": "Task not found"}, status_code=404)

    results = conn.execute(
        "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id",
        (task_id,)
    ).fetchall()
    conn.close()

    try:
        from docx import Document as DocxDocument
        from docx.shared import Pt, RGBColor
        from io import BytesIO

        doc = DocxDocument()
        doc.add_heading(task["title"], level=1)
        doc.add_paragraph(f'Feladat: {task["description"]}')
        doc.add_paragraph(f'Kiadta: {task["assigned_by"]} | Dátum: {task["created_at"][:10]} | Státusz: {task["status"]}')
        doc.add_paragraph("—" * 40)

        agent_names = {"kimi": "Kimi-K2.6 (Kutató)", "deepseek": "DeepSeek V3.2 (Kritikus)", "glm5": "GLM-5.1 (Végrehajtó)", "szintézis": "Koordinátori Szintézis"}
        for r in results:
            name = agent_names.get(r["agent"], r["agent"].upper())
            doc.add_heading(name, level=2)
            for para_text in (r["content"] or "(nincs tartalom)").split("\n"):
                if para_text.strip():
                    doc.add_paragraph(para_text.strip())

        doc.add_paragraph("—" * 40)
        doc.add_paragraph("Készítette: Claus Multi-Agent Rendszer (Kimi-K2.6 + DeepSeek V3.2 + GLM-5.1)")

        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)

        from starlette.responses import Response
        filename = f"claus_ai_task_{task_id}.docx"
        return Response(
            content=buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@mcp.custom_route("/api/ai_tasks/{task_id}/export_xlsx", methods=["GET"])
async def api_ai_task_export_xlsx(request):
    """Export AI task results as .xlsx via shared helper."""
    task_id = request.path_params["task_id"]
    conn = get_db()
    task = conn.execute("SELECT * FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
    if not task:
        conn.close()
        return JSONResponse({"error": "Task not found"}, status_code=404)
    results = conn.execute(
        "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id", (task_id,)
    ).fetchall()
    conn.close()
    try:
        buf = _generate_xlsx(task, results)
        from starlette.responses import Response
        return Response(
            content=buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="claus_ai_task_{task_id}.xlsx"'},
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@mcp.custom_route("/api/ai_tasks/{task_id}/export_pptx", methods=["GET"])
async def api_ai_task_export_pptx(request):
    """Export AI task synthesis as .pptx via shared helper."""
    task_id = request.path_params["task_id"]
    conn = get_db()
    task = conn.execute("SELECT * FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
    if not task:
        conn.close()
        return JSONResponse({"error": "Task not found"}, status_code=404)
    results = conn.execute(
        "SELECT agent, role, content, timestamp FROM ai_task_results WHERE task_id = ? ORDER BY id", (task_id,)
    ).fetchall()
    conn.close()
    try:
        buf = _generate_pptx(task, results)
        from starlette.responses import Response
        return Response(
            content=buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="claus_ai_task_{task_id}.pptx"'},
        )
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# ============================================================
# GOOGLE API — Gmail + Calendar Capture
# ============================================================

GOOGLE_TOKEN_JSON = os.environ.get("GOOGLE_TOKEN_JSON", "")
GOOGLE_SCOPES = [
    "https://www.googleapis.com/auth/gmail.modify",
    "https://www.googleapis.com/auth/gmail.send",
    "https://www.googleapis.com/auth/calendar.events",
]
TELEGRAM_BOT_TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN", "")
TELEGRAM_CHAT_ID = os.environ.get("TELEGRAM_CHAT_ID", "")

# Capture config (env vars)
GMAIL_POLL_INTERVAL = int(os.environ.get("GMAIL_POLL_INTERVAL", "300"))
CALENDAR_POLL_INTERVAL = int(os.environ.get("CALENDAR_POLL_INTERVAL", "900"))
CALENDAR_REMINDER_MINUTES = int(os.environ.get("CALENDAR_REMINDER_MINUTES", "30"))
MORNING_BRIEFING_HOUR = int(os.environ.get("MORNING_BRIEFING_HOUR", "7"))
IGNORE_SENDERS = [s.strip().lower() for s in os.environ.get("IGNORE_SENDERS", "newsletter,noreply,no-reply,marketing").split(",") if s.strip()]
URGENT_SENDERS = [s.strip().lower() for s in os.environ.get("URGENT_SENDERS", "").split(",") if s.strip()]
URGENT_KEYWORDS = [k.strip().lower() for k in os.environ.get("URGENT_KEYWORDS", "urgent,sürgős,asap,critical,azonnal,fontos").split(",") if k.strip()]

# Runtime state for capture polling
_capture_state = {
    "gmail_history_id": None,
    "calendar_reminded": set(),
    "last_briefing_date": None,
    "gmail_service": None,
    "calendar_service": None,
    "capture_running": False,
}


def _init_google_services():
    """Initialize Google API services from token JSON env var."""
    token_raw = GOOGLE_TOKEN_JSON
    logger.info("GOOGLE_TOKEN_JSON env var length: %d", len(token_raw))
    if not token_raw:
        logger.info("GOOGLE_TOKEN_JSON not set — capture disabled")
        return False
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        from googleapiclient.discovery import build

        # Support both raw JSON and base64-encoded JSON
        try:
            token_data = json.loads(token_raw)
            logger.info("Google token parsed as raw JSON")
        except json.JSONDecodeError:
            token_data = json.loads(base64.b64decode(token_raw).decode())
            logger.info("Google token parsed as base64")
        creds = Credentials.from_authorized_user_info(token_data, GOOGLE_SCOPES)
        logger.info("Google creds created. expired=%s, has_refresh=%s", creds.expired, bool(creds.refresh_token))

        if creds.expired and creds.refresh_token:
            creds.refresh(Request())
            logger.info("Google token refreshed successfully")

        _capture_state["gmail_service"] = build("gmail", "v1", credentials=creds, cache_discovery=False)
        _capture_state["calendar_service"] = build("calendar", "v3", credentials=creds, cache_discovery=False)

        profile = _capture_state["gmail_service"].users().getProfile(userId="me").execute()
        logger.info("Google services initialized: %s", profile.get("emailAddress"))
        return True
    except Exception as e:
        logger.error("Google init failed: %s (type: %s)", e, type(e).__name__)
        import traceback
        logger.error("Google init traceback:\n%s", traceback.format_exc())
        return False


def _extract_attachments(payload: dict) -> list:
    """Extract attachment metadata from Gmail message payload (format='full')."""
    attachments = []
    _walk_parts(payload.get("parts", []), attachments)
    return attachments


def _walk_parts(parts: list, out: list):
    """Recursively walk MIME parts to find attachments."""
    for part in parts:
        filename = part.get("filename", "")
        if filename:
            body = part.get("body", {})
            out.append({
                "filename": filename,
                "mime_type": part.get("mimeType", "application/octet-stream"),
                "size": body.get("size", 0),
                "attachment_id": body.get("attachmentId", ""),
            })
        if part.get("parts"):
            _walk_parts(part["parts"], out)


def _categorize_email(sender: str, subject: str) -> str:
    """Categorize email: urgent / important / normal / ignore."""
    sender_lower = sender.lower()
    subject_lower = subject.lower()

    for pattern in IGNORE_SENDERS:
        if pattern in sender_lower:
            return "ignore"

    for pattern in URGENT_SENDERS:
        if pattern in sender_lower:
            return "urgent"

    for kw in URGENT_KEYWORDS:
        if kw in subject_lower:
            return "urgent"

    return "normal"


async def _analyze_image(image_base64: str, mime_type: str = "image/jpeg",
                         prompt: str = "Mit latsz a kepen? Ird le reszletesen, magyarul.",
                         model: str = "kimi") -> str:
    """Central image analysis via vision model (Kimi K2.6). Usable from any channel."""
    if not SILICONFLOW_API_KEY:
        return "(Vision nem elerheto — SILICONFLOW_API_KEY hianzik)"

    model_id = SILICONFLOW_MODELS.get(model, "moonshotai/Kimi-K2.6")
    data_url = f"data:{mime_type};base64,{image_base64}"

    try:
        import httpx
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{SILICONFLOW_BASE_URL}/chat/completions",
                headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                json={
                    "model": model_id,
                    "messages": [{"role": "user", "content": [
                        {"type": "image_url", "image_url": {"url": data_url}},
                        {"type": "text", "text": prompt},
                    ]}],
                    "temperature": 0.5,
                    "max_tokens": 2000,
                },
            )
        data = json.loads(resp.text)
        if "error" in data or "code" in data:
            logger.error("Vision API error: %s", data)
            return f"(Vision hiba: {data.get('message', data.get('error', 'unknown'))})"
        content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
        return content or "(ures valasz a vision modeltol)"
    except Exception as e:
        logger.error("Vision analysis failed: %s", e)
        return f"(Vision hiba: {e})"


async def _telegram_push(text: str):
    """Send push notification to Telegram."""
    if not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_ID:
        return
    try:
        import httpx
        async with httpx.AsyncClient(timeout=15) as client:
            await client.post(
                f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage",
                json={
                    "chat_id": TELEGRAM_CHAT_ID,
                    "text": text[:4000],
                    "parse_mode": "HTML",
                    "disable_web_page_preview": True,
                },
            )
    except Exception as e:
        logger.error("Telegram push failed: %s", e)


async def _telegram_push_document(message_id: str, attachment: dict, caption: str = ""):
    """Download Gmail attachment and send to Telegram as document."""
    svc = _capture_state.get("gmail_service")
    if not svc or not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_ID:
        return
    try:
        att_id = attachment.get("attachment_id", "")
        if not att_id:
            return
        att_data = svc.users().messages().attachments().get(
            userId="me", messageId=message_id, id=att_id
        ).execute()
        file_bytes = base64.urlsafe_b64decode(att_data["data"])
        import httpx
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendDocument",
                data={
                    "chat_id": TELEGRAM_CHAT_ID,
                    "caption": caption[:1024] if caption else "",
                    "parse_mode": "HTML",
                },
                files={
                    "document": (attachment["filename"], file_bytes, attachment["mime_type"]),
                },
            )
            if resp.status_code != 200:
                logger.error("Telegram sendDocument failed: %s", resp.text[:200])
    except Exception as e:
        logger.error("Telegram document push failed: %s", e)


async def _telegram_send_file(file_bytes: bytes, filename: str, mime_type: str, caption: str = ""):
    """Send any file to Telegram as a document."""
    if not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_ID:
        return False
    try:
        import httpx
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendDocument",
                data={
                    "chat_id": TELEGRAM_CHAT_ID,
                    "caption": caption[:1024] if caption else "",
                    "parse_mode": "HTML",
                },
                files={
                    "document": (filename, file_bytes, mime_type),
                },
            )
            if resp.status_code == 200:
                return True
            logger.error("Telegram sendFile failed: %s", resp.text[:200])
    except Exception as e:
        logger.error("Telegram file send failed: %s", e)
    return False


def _bridge_capture_event(subject: str, message: str, priority: str = "normal"):
    """Write a capture event directly into the Bridge DB."""
    conn = get_db()
    ts = now()
    cur = conn.execute(
        "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority, thread_id, reply_to) "
        "VALUES (?, ?, ?, ?, ?, ?, NULL, NULL)",
        (ts, "capture-daemon", "kommandant", subject, message, priority)
    )
    msg_id = cur.lastrowid
    conn.execute("UPDATE messages SET thread_id = ? WHERE id = ?", (msg_id, msg_id))
    conn.commit()
    conn.close()


# ============================================================
# CAPTURE MCP TOOLS
# ============================================================

@mcp.tool()
async def capture_gmail_poll(max_results: int = 10, caller: str = "") -> str:
    """Poll Gmail for new unread emails. Returns categorized capture events.

    Args:
        max_results: Max emails to fetch (default 10)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "capture_gmail_poll")
    if denied:
        return denied
    svc = _capture_state.get("gmail_service")
    if not svc:
        return json.dumps({"error": "Gmail not initialized. Set GOOGLE_TOKEN_JSON env var."})

    try:
        history_id = _capture_state.get("gmail_history_id")

        if history_id is None:
            # Initial fetch — recent unread
            results = svc.users().messages().list(
                userId="me", q="is:unread category:primary", maxResults=max_results
            ).execute()
            msg_stubs = results.get("messages", [])
        else:
            # Incremental via History API
            try:
                history = svc.users().history().list(
                    userId="me", startHistoryId=history_id,
                    historyTypes=["messageAdded"], labelId="INBOX"
                ).execute()
                msg_stubs = []
                seen = set()
                for record in history.get("history", []):
                    for added in record.get("messagesAdded", []):
                        m = added.get("message", {})
                        mid = m.get("id")
                        if mid and mid not in seen and "INBOX" in m.get("labelIds", []):
                            msg_stubs.append({"id": mid})
                            seen.add(mid)
            except Exception:
                # History expired, fall back to initial
                results = svc.users().messages().list(
                    userId="me", q="is:unread category:primary", maxResults=max_results
                ).execute()
                msg_stubs = results.get("messages", [])

        # Update history watermark
        profile = svc.users().getProfile(userId="me").execute()
        _capture_state["gmail_history_id"] = int(profile.get("historyId", 0))

        events = []
        for stub in msg_stubs[:max_results]:
            try:
                msg = svc.users().messages().get(
                    userId="me", id=stub["id"], format="full"
                ).execute()
                payload = msg.get("payload", {})
                headers = {h["name"]: h["value"] for h in payload.get("headers", [])}
                sender_raw = headers.get("From", "unknown")
                subject = headers.get("Subject", "(no subject)")
                date_str = headers.get("Date", "")
                snippet = msg.get("snippet", "")
                sender_name, sender_email = parseaddr(sender_raw)
                priority = _categorize_email(sender_raw, subject)

                if priority == "ignore":
                    continue

                # Extract attachments
                attachments = _extract_attachments(payload)
                att_info = [{"filename": a["filename"], "mime_type": a["mime_type"],
                             "size": a["size"]} for a in attachments]

                event = {
                    "message_id": stub["id"],
                    "subject": subject,
                    "sender": sender_name or sender_email,
                    "sender_email": sender_email,
                    "snippet": snippet,
                    "date": date_str,
                    "priority": priority,
                    "attachments": att_info,
                }
                events.append(event)

                # Format attachment line for Bridge/Telegram
                att_line = ""
                if attachments:
                    att_names = ", ".join(a["filename"] for a in attachments)
                    att_line = f"\n📎 Csatolmányok ({len(attachments)}): {att_names}"

                # Write to Bridge DB
                _bridge_capture_event(
                    f"📧 Gmail: {subject}",
                    f"From: {sender_name or sender_email} <{sender_email}>\nSubject: {subject}\nDate: {date_str}\n\n{snippet}{att_line}",
                    priority
                )

                # Telegram push for urgent/important
                if priority in ("urgent", "important"):
                    await _telegram_push(
                        f"🟠 <b>GMAIL</b> — {subject}\n\n"
                        f"<b>From:</b> {sender_name or sender_email}\n"
                        f"{snippet}"
                        + (f"\n\n📎 <b>Csatolmányok:</b> {', '.join(a['filename'] for a in attachments)}" if attachments else "")
                    )
                    # Forward attachments to Telegram
                    for att in attachments:
                        await _telegram_push_document(
                            stub["id"], att,
                            caption=f"📎 <b>{att['filename']}</b>\n{subject}"
                        )
            except Exception as e:
                logger.error("Failed to process message %s: %s", stub.get("id"), e)

        return json.dumps({"count": len(events), "events": events}, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def capture_calendar_poll(caller: str = "") -> str:
    """Poll Google Calendar for upcoming events (reminders + briefing).

    Returns reminders for events starting within the configured reminder window,
    and a morning briefing if it's briefing time.
    """
    denied = _enforce(caller, "capture_calendar_poll")
    if denied:
        return denied
    svc = _capture_state.get("calendar_service")
    if not svc:
        return json.dumps({"error": "Calendar not initialized. Set GOOGLE_TOKEN_JSON env var."})

    try:
        events_out = []
        utc_now = datetime.now(timezone.utc)

        # --- Reminders ---
        reminder_window = utc_now + timedelta(minutes=CALENDAR_REMINDER_MINUTES)
        result = svc.events().list(
            calendarId="primary",
            timeMin=utc_now.isoformat(),
            timeMax=reminder_window.isoformat(),
            singleEvents=True, orderBy="startTime", maxResults=10
        ).execute()

        reminded = _capture_state.get("calendar_reminded", set())

        for item in result.get("items", []):
            event_id = item.get("id", "")
            if event_id in reminded:
                continue

            start_info = item.get("start", {})
            dt_str = start_info.get("dateTime", start_info.get("date", ""))
            if "T" in dt_str:
                start_dt = datetime.fromisoformat(dt_str).astimezone(timezone.utc)
                minutes_until = max(0, int((start_dt - utc_now).total_seconds() / 60))
            else:
                minutes_until = -1  # all-day

            summary = item.get("summary", "(no title)")
            location = item.get("location", "")
            meet_link = item.get("hangoutLink", "")
            attendees = [a.get("email", "") for a in item.get("attendees", [])[:5]]

            event_out = {
                "type": "reminder",
                "event_id": event_id,
                "summary": summary,
                "start": dt_str,
                "minutes_until": minutes_until,
                "location": location,
                "meet_link": meet_link,
                "attendees": attendees,
            }
            events_out.append(event_out)
            reminded.add(event_id)

            # Bridge + Telegram
            priority = "urgent" if minutes_until <= 10 else "important"
            lines = [f"Esemény: {summary}", f"Kezdés: {dt_str} ({minutes_until} perc múlva)"]
            if location: lines.append(f"Helyszín: {location}")
            if meet_link: lines.append(f"Link: {meet_link}")

            _bridge_capture_event(f"📅 {minutes_until} perc múlva: {summary}", "\n".join(lines), priority)
            await _telegram_push(
                f"⏰ <b>NAPTÁR</b> — {minutes_until} perc múlva\n\n"
                f"<b>{summary}</b>\n"
                + (f"📍 {location}\n" if location else "")
                + (f"🔗 <a href=\"{meet_link}\">Csatlakozás</a>" if meet_link else "")
            )

        _capture_state["calendar_reminded"] = reminded

        # --- Morning briefing ---
        # FIX: expliciten Budapest-i ido, kulonben Railway UTC konteneren a
        # "7 AM" = 09:00 CEST (DST alatt) — reggel 7-nek szantuk, de 9-kor futott.
        try:
            from zoneinfo import ZoneInfo
            local_now = datetime.now(ZoneInfo("Europe/Budapest"))
        except ImportError:
            # Py <3.9 fallback: hardcoded CEST offset
            local_now = datetime.now(timezone(timedelta(hours=2)))
        today_str = local_now.strftime("%Y-%m-%d")

        if local_now.hour == MORNING_BRIEFING_HOUR and _capture_state.get("last_briefing_date") != today_str:
            start_of_day = utc_now.replace(hour=0, minute=0, second=0, microsecond=0)
            end_of_day = start_of_day + timedelta(days=1)

            day_result = svc.events().list(
                calendarId="primary",
                timeMin=start_of_day.isoformat(),
                timeMax=end_of_day.isoformat(),
                singleEvents=True, orderBy="startTime", maxResults=20
            ).execute()

            day_items = day_result.get("items", [])
            _capture_state["last_briefing_date"] = today_str

            if day_items:
                event_lines = []
                for item in day_items:
                    s = item.get("start", {})
                    t = s.get("dateTime", s.get("date", "?"))
                    if "T" in t: t = t[11:16]
                    event_lines.append(f"  {t} — {item.get('summary', '?')}")

                briefing_body = f"Mai nap: {len(day_items)} esemény:\n\n" + "\n".join(event_lines)
            else:
                briefing_body = "Tiszta nap — nincsenek események a naptárban."

            events_out.append({"type": "briefing", "date": today_str, "event_count": len(day_items)})
            _bridge_capture_event(f"🌅 Reggeli briefing — {today_str}", briefing_body, "normal")
            await _telegram_push(f"🌅 <b>Reggeli briefing</b> — {today_str}\n\n{briefing_body}")

        return json.dumps({"count": len(events_out), "events": events_out}, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def create_calendar_event(
    summary: str, start: str, end: str = "", description: str = "",
    location: str = "", calendar_id: str = "primary", caller: str = ""
) -> str:
    """Create a Google Calendar event.

    Args:
        summary: Event title
        start: Start time in ISO 8601 format (e.g. '2026-03-31T14:00:00+02:00' or '2026-03-31' for all-day)
        end: End time in ISO 8601 (default: 1 hour after start, or next day for all-day)
        description: Optional event description
        location: Optional location
        calendar_id: Calendar ID (default: primary)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "create_calendar_event")
    if denied:
        return denied
    svc = _capture_state.get("calendar_service")
    if not svc:
        return json.dumps({"error": "Calendar not initialized. Set GOOGLE_TOKEN_JSON env var."})

    try:
        is_all_day = "T" not in start

        if is_all_day:
            event_body = {
                "summary": summary,
                "start": {"date": start},
                "end": {"date": end or (datetime.fromisoformat(start) + timedelta(days=1)).strftime("%Y-%m-%d")},
            }
        else:
            start_dt = datetime.fromisoformat(start)
            if not end:
                end = (start_dt + timedelta(hours=1)).isoformat()
            event_body = {
                "summary": summary,
                "start": {"dateTime": start},
                "end": {"dateTime": end},
            }

        if description:
            event_body["description"] = description
        if location:
            event_body["location"] = location

        result = svc.events().insert(calendarId=calendar_id, body=event_body).execute()
        return json.dumps({
            "status": "created",
            "event_id": result.get("id"),
            "link": result.get("htmlLink"),
            "summary": summary,
            "start": start,
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def capture_send_email(
    to: str, subject: str, body: str, body_type: str = "plain",
    file_id: int = 0, attachment_base64: str = "", attachment_filename: str = "",
    attachment_mime: str = "", caller: str = ""
) -> str:
    """Send an email via Gmail API, optionally with an attachment.

    Args:
        to: Recipient email address
        subject: Email subject
        body: Email body content
        body_type: 'plain' or 'html' (default: plain)
        file_id: Optional — uploaded file ID from upload_file to attach
        attachment_base64: Optional — raw base64-encoded file content to attach
        attachment_filename: Filename for the attachment (required if attachment_base64 is used)
        attachment_mime: MIME type for attachment
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "capture_send_email")
    if denied:
        return denied
    svc = _capture_state.get("gmail_service")
    if not svc:
        return json.dumps({"error": "Gmail not initialized. Set GOOGLE_TOKEN_JSON env var."})

    try:
        # Resolve attachment from uploads table if file_id given
        att_data = None
        att_name = attachment_filename
        att_mime = attachment_mime
        if file_id:
            conn = get_db()
            f = conn.execute("SELECT filename, mime_type, content_base64 FROM uploads WHERE id = ?", (file_id,)).fetchone()
            conn.close()
            if not f or not f["content_base64"]:
                return json.dumps({"error": f"Upload #{file_id} not found or has no binary content"})
            att_data = base64.b64decode(f["content_base64"])
            att_name = att_name or f["filename"]
            att_mime = att_mime or f["mime_type"] or "application/octet-stream"
        elif attachment_base64:
            if not att_name:
                return json.dumps({"error": "attachment_filename is required when using attachment_base64"})
            att_data = base64.b64decode(attachment_base64)
            att_mime = att_mime or "application/octet-stream"

        if att_data:
            msg = MIMEMultipart()
            msg["to"] = to
            msg["subject"] = subject
            msg.attach(MIMEText(body, body_type))

            maintype, _, subtype = att_mime.partition("/")
            part = MIMEBase(maintype or "application", subtype or "octet-stream")
            part.set_payload(att_data)
            encoders.encode_base64(part)
            part.add_header("Content-Disposition", "attachment", filename=att_name)
            msg.attach(part)
        else:
            msg = MIMEText(body, body_type)
            msg["to"] = to
            msg["subject"] = subject

        raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
        result = svc.users().messages().send(userId="me", body={"raw": raw}).execute()
        return json.dumps({"status": "sent", "message_id": result.get("id"),
                           "attachment": att_name if att_data else None})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def capture_inbox(limit: int = 10, query: str = "is:unread category:primary",
                        caller: str = "") -> str:
    """Read Gmail inbox messages without triggering capture events. For browsing email.

    Args:
        limit: Max messages (default 10)
        query: Gmail search query (default: unread primary)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "capture_inbox")
    if denied:
        return denied
    svc = _capture_state.get("gmail_service")
    if not svc:
        return json.dumps({"error": "Gmail not initialized."})

    try:
        results = svc.users().messages().list(userId="me", q=query, maxResults=limit).execute()
        messages = []
        for stub in results.get("messages", []):
            msg = svc.users().messages().get(
                userId="me", id=stub["id"], format="full"
            ).execute()
            payload = msg.get("payload", {})
            headers = {h["name"]: h["value"] for h in payload.get("headers", [])}
            sender_name, sender_email = parseaddr(headers.get("From", ""))
            attachments = _extract_attachments(payload)
            att_info = [{"filename": a["filename"], "mime_type": a["mime_type"],
                         "size": a["size"]} for a in attachments]
            messages.append({
                "id": stub["id"],
                "subject": headers.get("Subject", ""),
                "from": sender_name or sender_email,
                "from_email": sender_email,
                "date": headers.get("Date", ""),
                "snippet": msg.get("snippet", ""),
                "attachments": att_info,
            })
        return json.dumps({"count": len(messages), "messages": messages}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def read_gmail_attachment(message_id: str, attachment_index: int = 0,
                                caller: str = "") -> str:
    """Download and read a Gmail attachment's content as text.

    First use capture_inbox to find the message_id and see attachment list,
    then call this with the message_id and attachment index (0-based).

    Supports: .docx, .pdf, .txt, .csv, .md, .doc and other text formats.
    Images return metadata only (filename, size, mime_type).

    Args:
        message_id: Gmail message ID (from capture_inbox results)
        attachment_index: Which attachment to read (0 = first, default)
        caller: Instance ID for permission check
    """
    denied = _enforce(caller, "read_gmail_attachment")
    if denied:
        return denied
    svc = _capture_state.get("gmail_service")
    if not svc:
        return json.dumps({"error": "Gmail not initialized."})

    try:
        msg = svc.users().messages().get(userId="me", id=message_id, format="full").execute()
        attachments = _extract_attachments(msg.get("payload", {}))

        if not attachments:
            return json.dumps({"error": "No attachments in this message"})
        if attachment_index >= len(attachments):
            return json.dumps({"error": f"Attachment index {attachment_index} out of range (message has {len(attachments)} attachments)"})

        att = attachments[attachment_index]
        att_id = att.get("attachment_id", "")
        if not att_id:
            return json.dumps({"error": "Attachment has no downloadable ID"})

        att_data = svc.users().messages().attachments().get(
            userId="me", messageId=message_id, id=att_id
        ).execute()
        file_bytes = base64.urlsafe_b64decode(att_data["data"])

        filename = att["filename"]
        mime_type = att["mime_type"]

        if _is_image(mime_type):
            return json.dumps({
                "filename": filename,
                "mime_type": mime_type,
                "size": len(file_bytes),
                "content": "(image — use upload_file to store for AI vision analysis)",
            }, ensure_ascii=False)

        # Write to temp file for parsing
        tmp_path = UPLOAD_DIR / f"att_{int(time.time())}_{filename}"
        tmp_path.write_bytes(file_bytes)
        text_content = _parse_file_to_text(tmp_path, mime_type)

        # For .doc files (old Word format), try antiword or basic extraction
        if not text_content and filename.lower().endswith(".doc"):
            try:
                import subprocess
                result = subprocess.run(["antiword", str(tmp_path)], capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    text_content = result.stdout
            except Exception:
                text_content = "(Old .doc format — antiword not available, content cannot be extracted)"

        tmp_path.unlink(missing_ok=True)

        if not text_content:
            text_content = f"(Cannot extract text from {mime_type} — unsupported format)"

        return json.dumps({
            "filename": filename,
            "mime_type": mime_type,
            "size": len(file_bytes),
            "content": text_content[:50000],
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def capture_status(caller: str = "") -> str:
    """Get capture daemon status — is Gmail/Calendar connected, last poll info."""
    denied = _enforce(caller, "capture_status")
    if denied:
        return denied
    gmail_ok = _capture_state.get("gmail_service") is not None
    cal_ok = _capture_state.get("calendar_service") is not None
    # If not connected, try to reinit and report the error
    diag = None
    if not gmail_ok:
        diag = _diagnose_google_init()

    return json.dumps({
        "gmail_connected": gmail_ok,
        "calendar_connected": cal_ok,
        "gmail_history_id": _capture_state.get("gmail_history_id"),
        "capture_loop_running": _capture_state.get("capture_running", False),
        "calendar_reminded_count": len(_capture_state.get("calendar_reminded", set())),
        "last_briefing_date": _capture_state.get("last_briefing_date"),
        "google_diag": diag,
        "config": {
            "gmail_poll_interval": GMAIL_POLL_INTERVAL,
            "calendar_poll_interval": CALENDAR_POLL_INTERVAL,
            "calendar_reminder_minutes": CALENDAR_REMINDER_MINUTES,
            "morning_briefing_hour": MORNING_BRIEFING_HOUR,
        }
    })


def _diagnose_google_init() -> dict:
    """Diagnose why Google services aren't working. Returns dict with details."""
    import traceback
    token_raw = GOOGLE_TOKEN_JSON
    result = {"token_env_length": len(token_raw)}

    if not token_raw:
        result["error"] = "GOOGLE_TOKEN_JSON env var is empty or not set"
        return result

    # Try parse
    try:
        try:
            token_data = json.loads(token_raw)
            result["parse_method"] = "raw_json"
        except json.JSONDecodeError:
            token_data = json.loads(base64.b64decode(token_raw).decode())
            result["parse_method"] = "base64"
        result["has_refresh_token"] = bool(token_data.get("refresh_token"))
        result["has_client_id"] = bool(token_data.get("client_id"))
        result["expiry"] = token_data.get("expiry", "none")
    except Exception as e:
        result["parse_error"] = f"{type(e).__name__}: {e}"
        return result

    # Try create creds + refresh
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        creds = Credentials.from_authorized_user_info(token_data, GOOGLE_SCOPES)
        result["creds_expired"] = creds.expired
        result["creds_valid"] = creds.valid
        if creds.expired and creds.refresh_token:
            creds.refresh(Request())
            result["refresh_success"] = True
            result["new_expiry"] = creds.expiry.isoformat() if creds.expiry else "none"
            # Actually init the services now!
            from googleapiclient.discovery import build
            _capture_state["gmail_service"] = build("gmail", "v1", credentials=creds, cache_discovery=False)
            _capture_state["calendar_service"] = build("calendar", "v3", credentials=creds, cache_discovery=False)
            profile = _capture_state["gmail_service"].users().getProfile(userId="me").execute()
            result["reinit_success"] = True
            result["email"] = profile.get("emailAddress")
            # Start capture loop if not running
            if not _capture_state.get("capture_running"):
                _start_capture_background()
                result["capture_restarted"] = True
    except Exception as e:
        result["refresh_error"] = f"{type(e).__name__}: {e}"
        result["traceback"] = traceback.format_exc()[-500:]

    return result


# ============================================================
# AUTO-POLLING BACKGROUND LOOP
# ============================================================

async def _capture_loop():
    """Background loop that auto-polls Gmail and Calendar."""
    if not _capture_state.get("gmail_service"):
        logger.info("Capture loop not started — Google services not initialized")
        return

    _capture_state["capture_running"] = True
    logger.info("Capture loop started: Gmail=%ds, Calendar=%ds", GMAIL_POLL_INTERVAL, CALENDAR_POLL_INTERVAL)

    _bridge_capture_event(
        "🟢 Capture Daemon elindult (Railway)",
        f"Gmail polling: {GMAIL_POLL_INTERVAL}s\nCalendar polling: {CALENDAR_POLL_INTERVAL}s\nCalendar reminder: {CALENDAR_REMINDER_MINUTES} perccel előtte",
        "info"
    )

    gmail_counter = 0
    calendar_counter = 0
    tick = 60  # check every 60 seconds

    while _capture_state.get("capture_running"):
        try:
            gmail_counter += tick
            calendar_counter += tick

            if gmail_counter >= GMAIL_POLL_INTERVAL:
                gmail_counter = 0
                await capture_gmail_poll()

            if calendar_counter >= CALENDAR_POLL_INTERVAL:
                calendar_counter = 0
                await capture_calendar_poll()

        except Exception as e:
            logger.error("Capture loop error: %s", e)

        await asyncio.sleep(tick)


def _start_capture_background():
    """Start capture loop in a background thread with its own event loop."""
    def _run():
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_capture_loop())

    if _capture_state.get("gmail_service"):
        t = threading.Thread(target=_run, daemon=True)
        t.start()
        logger.info("Capture background thread started")


# ============================================================
# TELEGRAM BOT POLLING — Kommandant ír, Bridge fogadja
# ============================================================

_telegram_state = {"last_update_id": 0}

TELEGRAM_POLL_INTERVAL = 5  # seconds


async def _handle_telegram_message(text: str, chat_id: str):
    """Process an incoming Telegram message from Kommandant."""
    # 1. Store in Bridge DB
    conn = get_db()
    cur = conn.execute(
        "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority) "
        "VALUES (?, 'kommandant', 'bridge', ?, ?, 'normal')",
        (now(), f"Telegram: {text[:60]}", text)
    )
    msg_id = cur.lastrowid
    conn.execute("UPDATE messages SET thread_id = ? WHERE id = ?", (msg_id, msg_id))
    conn.commit()
    conn.close()
    logger.info("Telegram message stored in Bridge: #%d", msg_id)

    # 2. Feldwebel command routing (slash commands)
    if FELDWEBEL_ENABLED:
        try:
            from feldwebel.commands import handle_command
            if await handle_command(text, chat_id):
                return  # Command was handled
        except Exception as e:
            logger.error("Feldwebel command error: %s", e)

    # 3. Check for explicit agent mentions (@kimi, @glm5)
    text_lower = text.lower()
    mentioned = [a for a in ("kimi", "deepseek", "glm5") if f"@{a}" in text_lower]

    if mentioned and PYRAMID_ENABLED and SILICONFLOW_API_KEY:
        # Call specific agents by mention (existing logic)
        import httpx
        from html import escape as html_escape
        for agent_id in mentioned:
            try:
                system_prompt = build_agent_context(agent_id=agent_id, inbox_summary=_get_inbox_summary(), relevance_query=text)
                model_id = SILICONFLOW_MODELS.get(agent_id, agent_id)
                async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                    resp = await client.post(
                        f"{SILICONFLOW_BASE_URL}/chat/completions",
                        headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                        json={"model": model_id, "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": f"Kommandant üzenete Telegramról:\n\n{text}"},
                        ], "temperature": 0.7, "max_tokens": 1500},
                    )
                    data = json.loads(resp.text)

                if "error" in data:
                    await _telegram_push(f"<b>{agent_id.upper()}</b> — API hiba: {html_escape(str(data['error']))}")
                    continue

                content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
                if not content or not content.strip():
                    await _telegram_push(f"<b>{agent_id.upper()}</b> — üres válasz")
                    continue

                safe_content = html_escape(content[:3800])
                await _telegram_push(f"<b>{agent_id.upper()}</b>\n\n{safe_content}")

                conn2 = get_db()
                conn2.execute(
                    "INSERT INTO messages (timestamp, sender, recipient, subject, message, priority, thread_id, reply_to) "
                    "VALUES (?, ?, 'kommandant', ?, ?, 'normal', ?, ?)",
                    (now(), agent_id, f"Re: @{agent_id} válasz", content, msg_id, msg_id)
                )
                conn2.commit()
                conn2.close()

                if PYRAMID_ENABLED:
                    try:
                        pyramid_store_result(content=content, agent_id=agent_id, task_title=f"telegram:{text[:60]}", force_shared=True)
                    except Exception:
                        pass

                logger.info("Agent %s replied to Telegram message #%d (%d chars)", agent_id, msg_id, len(content))
            except Exception as e:
                logger.error("Agent %s Telegram reply failed: %s", agent_id, e)
                await _telegram_push(f"<b>{agent_id.upper()}</b> — hiba: {html_escape(str(e)[:500])}")
        return  # Agent mention handled

    # 4. Free text → Feldwebel responder (DeepSeek with tool_choice=required + history)
    if FELDWEBEL_ENABLED:
        try:
            from feldwebel.responder import respond
            await respond(text, chat_id)
            return
        except Exception as e:
            logger.error("Feldwebel responder error: %s", e)
            from html import escape as html_escape
            await _telegram_push(f"<b>FELDWEBEL</b> — hiba: {html_escape(str(e)[:500])}")

    # 5. Fallback: original DeepSeek call without history (if Feldwebel unavailable)
    if PYRAMID_ENABLED and SILICONFLOW_API_KEY:
        import httpx
        from html import escape as html_escape
        try:
            system_prompt = build_agent_context(agent_id="deepseek", inbox_summary=_get_inbox_summary())
            model_id = SILICONFLOW_MODELS.get("deepseek", "deepseek-ai/DeepSeek-V4-Pro")
            async with httpx.AsyncClient(timeout=SILICONFLOW_TIMEOUT) as client:
                resp = await client.post(
                    f"{SILICONFLOW_BASE_URL}/chat/completions",
                    headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}", "Content-Type": "application/json"},
                    json={"model": model_id, "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"Kommandant üzenete Telegramról:\n\n{text}"},
                    ], "temperature": 0.7, "max_tokens": 1500},
                )
                data = json.loads(resp.text)
            content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
            if content and content.strip():
                await _telegram_push(f"<b>DEEPSEEK</b>\n\n{html_escape(content[:3800])}")
        except Exception as e:
            logger.error("Fallback DeepSeek failed: %s", e)


async def _telegram_poll_loop():
    """Poll Telegram Bot API for incoming messages from Kommandant."""
    if not TELEGRAM_BOT_TOKEN or not TELEGRAM_CHAT_ID:
        logger.info("Telegram polling not started — token or chat_id missing")
        return

    import httpx
    logger.info("Telegram polling started (interval=%ds)", TELEGRAM_POLL_INTERVAL)

    # Get initial offset
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getUpdates",
                params={"limit": 1, "offset": -1},
            )
            data = resp.json()
            if data.get("ok") and data.get("result"):
                _telegram_state["last_update_id"] = data["result"][-1]["update_id"]
    except Exception as e:
        logger.warning("Telegram initial offset failed: %s", e)

    while True:
        try:
            async with httpx.AsyncClient(timeout=35) as client:
                resp = await client.get(
                    f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getUpdates",
                    params={
                        "offset": _telegram_state["last_update_id"] + 1,
                        "timeout": 25,
                        "allowed_updates": '["message"]',
                    },
                )
                data = resp.json()

            if data.get("ok") and data.get("result"):
                for update in data["result"]:
                    _telegram_state["last_update_id"] = update["update_id"]
                    msg = update.get("message", {})
                    text = msg.get("text", "")
                    chat_id = str(msg.get("chat", {}).get("id", ""))

                    # Only accept messages from Kommandant's chat
                    if chat_id != TELEGRAM_CHAT_ID:
                        continue

                    # Handle photo messages
                    photo = msg.get("photo")
                    if photo:
                        caption = msg.get("caption", "")
                        try:
                            # Get largest photo (last in array)
                            file_id = photo[-1]["file_id"]
                            async with httpx.AsyncClient(timeout=15) as dl_client:
                                # Get file path
                                file_resp = await dl_client.get(
                                    f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/getFile",
                                    params={"file_id": file_id},
                                )
                                file_data = file_resp.json()
                                file_path = file_data.get("result", {}).get("file_path", "")
                                if not file_path:
                                    await _telegram_push("Nem sikerult a kepet letolteni.")
                                    continue
                                # Download file
                                img_resp = await dl_client.get(
                                    f"https://api.telegram.org/file/bot{TELEGRAM_BOT_TOKEN}/{file_path}"
                                )
                                image_bytes = img_resp.content

                            # Detect mime type
                            mime = "image/jpeg"
                            if file_path.endswith(".png"):
                                mime = "image/png"
                            elif file_path.endswith(".webp"):
                                mime = "image/webp"

                            image_b64 = base64.b64encode(image_bytes).decode("ascii")

                            # Store image in uploads table (reusable for email, AI, etc.)
                            conn = get_db()
                            ts = now()
                            ext = file_path.rsplit(".", 1)[-1] if "." in file_path else "jpg"
                            filename = f"telegram_photo_{ts[:19].replace(':', '')}.{ext}"
                            conn.execute(
                                "INSERT INTO uploads (filename, mime_type, content_text, content_base64, uploaded_by, uploaded_at) "
                                "VALUES (?, ?, '', ?, 'kommandant-telegram', ?)",
                                (filename, mime, image_b64, ts),
                            )
                            upload_id = conn.execute("SELECT last_insert_rowid()").fetchone()[0]
                            conn.commit()
                            conn.close()
                            logger.info("Telegram photo stored: #%d %s (%d bytes)", upload_id, file_path, len(image_bytes))

                            # Only analyze if caption requests it (e.g. "elemezd", "mi ez", any text)
                            if caption:
                                await _telegram_push(f"📷 Kep fogadva (#{upload_id}), Kimi K2.6 elemzi...")
                                analysis = await _analyze_image(image_b64, mime, caption)
                                await _telegram_push(
                                    f"📷 <b>KEP ELEMZES</b> (Kimi K2.6)\n\n{analysis[:3600]}\n\n"
                                    f"<i>Fajl #{upload_id} — tovabbkuldheto: \"kuldd el emailben a #{upload_id}-t\"</i>"
                                )
                            else:
                                await _telegram_push(
                                    f"📷 Kep fogadva: <b>#{upload_id}</b>\n"
                                    f"• Elemzes: kuldd ujra caption-nel (pl. \"mi ez a kepen?\")\n"
                                    f"• Tovabbkuldeni: \"kuldd el emailben a #{upload_id}-t\"\n"
                                    f"• Vagy kerdezz: \"elemezd a #{upload_id}-es kepet\""
                                )
                        except Exception as e:
                            logger.error("Telegram photo handling failed: %s", e)
                            await _telegram_push(f"Kep feldolgozasi hiba: {str(e)[:200]}")
                        continue

                    if text:
                        await _handle_telegram_message(text, chat_id)

        except Exception as e:
            logger.error("Telegram poll error: %s", e)
            await asyncio.sleep(5)

        await asyncio.sleep(TELEGRAM_POLL_INTERVAL)


def _start_telegram_polling():
    """Start Telegram polling in a background thread."""
    def _run():
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_telegram_poll_loop())

    if TELEGRAM_BOT_TOKEN and TELEGRAM_CHAT_ID:
        t = threading.Thread(target=_run, daemon=True)
        t.start()
        logger.info("Telegram polling background thread started")


# ============================================================
# STARTUP
# ============================================================

def _migrate_db_to_volume():
    """If BRIDGE_DB_PATH points to a volume (e.g. /data/bridge.db) but the file
    doesn't exist yet, copy the old working-directory bridge.db there."""
    import shutil
    volume_db = DB_PATH  # e.g. /data/bridge.db
    old_db = os.path.join(os.path.dirname(__file__), "bridge.db")

    if volume_db == old_db or volume_db == "bridge.db":
        return  # no migration needed, same path

    if os.path.exists(volume_db):
        logger.info("Volume DB already exists: %s", volume_db)
        return

    # Ensure volume directory exists
    volume_dir = os.path.dirname(volume_db)
    if volume_dir:
        os.makedirs(volume_dir, exist_ok=True)

    if os.path.exists(old_db):
        shutil.copy2(old_db, volume_db)
        logger.info("Migrated DB: %s -> %s", old_db, volume_db)
    else:
        logger.info("No old DB found at %s, starting fresh at %s", old_db, volume_db)


_migrate_db_to_volume()
init_db()
_init_google_services()
_start_capture_background()
_start_telegram_polling()

# Register permission profiles
register_youngereka()
logger.info("Permission profiles registered: YoungeReka")

# Initialize Feldwebel
if FELDWEBEL_ENABLED:
    # Inject ALL tool functions into capture_state for Feldwebel access
    _capture_state["_capture_inbox_func"] = capture_inbox
    _capture_state["_ai_query_func"] = ai_query
    _capture_state["_ai_task_func"] = ai_task
    _capture_state["_run_agent_with_tools"] = _run_agent_with_tools
    _capture_state["_send_email_func"] = capture_send_email
    _capture_state["_read_gmail_attachment_func"] = read_gmail_attachment
    _capture_state["_create_calendar_func"] = create_calendar_event
    _capture_state["_read_memory_func"] = read_memory
    _capture_state["_search_memory_func"] = search_memory
    _capture_state["_send_message_func"] = send_message
    _capture_state["_search_discussions_func"] = search_discussions
    _capture_state["_telegram_send_file"] = _telegram_send_file
    _capture_state["_generate_xlsx"] = _generate_xlsx
    _capture_state["_generate_pptx"] = _generate_pptx
    init_feldwebel(BridgeContext(
        telegram_push=_telegram_push,
        get_inbox_summary=_get_inbox_summary,
        get_db=get_db,
        capture_state=_capture_state,
        siliconflow_api_key=SILICONFLOW_API_KEY,
        siliconflow_base_url=SILICONFLOW_BASE_URL,
        siliconflow_timeout=SILICONFLOW_TIMEOUT,
        siliconflow_models=SILICONFLOW_MODELS,
        capture_send_email=capture_send_email,
        capture_calendar_poll=capture_calendar_poll,
        create_calendar_event=create_calendar_event,
        list_tasks_func=list_tasks,
        create_task_func=create_task,
        telegram_chat_id=TELEGRAM_CHAT_ID or "",
    ))

# Operation Zahnrad — Plugin Auto-Discovery
try:
    from plugins import discover_and_register
    _plugin_deps = {
        "get_db": get_db,
        "siliconflow_api_key": SILICONFLOW_API_KEY,
        "siliconflow_base_url": SILICONFLOW_BASE_URL,
        "siliconflow_timeout": SILICONFLOW_TIMEOUT,
        "siliconflow_models": SILICONFLOW_MODELS,
        "ai_task_func": ai_task,
        # Operation Kabare: prefetch helper eleri a gmail/calendar service-t
        "capture_state": _capture_state,
    }
    _loaded_plugins = discover_and_register(mcp, _plugin_deps)
    if _loaded_plugins:
        logger.info("Operation Zahnrad: %d plugin(s) loaded: %s", len(_loaded_plugins), ", ".join(_loaded_plugins))
    else:
        logger.info("Operation Zahnrad: No plugins found")
except Exception as e:
    logger.error("Operation Zahnrad plugin loading failed: %s", e)

# Operation Zahnrad — Cron Scheduler
def _cron_matches(schedule: str, dt: datetime) -> bool:
    """Simple crontab matcher: 'minute hour day month weekday'. Supports: number, *, ranges (1-5), lists (1,3,5)."""
    parts = schedule.strip().split()
    if len(parts) != 5:
        return False
    checks = [
        (parts[0], dt.minute),
        (parts[1], dt.hour),
        (parts[2], dt.day),
        (parts[3], dt.month),
        (parts[4], dt.isoweekday() % 7),  # 0=sunday in crontab
    ]
    for pattern, value in checks:
        if pattern == "*":
            continue
        if "-" in pattern:
            lo, hi = pattern.split("-", 1)
            if not (int(lo) <= value <= int(hi)):
                return False
        elif "," in pattern:
            if value not in [int(x) for x in pattern.split(",")]:
                return False
        else:
            if int(pattern) != value:
                return False
    return True


async def _cron_loop():
    """Runs every 60s, checks scheduled recipes and executes them. Also cleans old uploads.

    Cron expressions are interpreted in Europe/Budapest time (CET/CEST), NOT UTC.
    This matches the user's expectation: '0 7 * * *' means 7 AM Budapest time,
    not 7 AM UTC which would be 9 AM Budapest during DST.
    """
    await asyncio.sleep(10)  # Let server fully start
    logger.info("Cron scheduler started (Europe/Budapest interpretation)")
    try:
        from zoneinfo import ZoneInfo
        BP_TZ = ZoneInfo("Europe/Budapest")
    except ImportError:
        BP_TZ = timezone(timedelta(hours=2))  # CEST fallback
    _last_cleanup = ""
    while True:
        try:
            now_dt = datetime.now(timezone.utc)           # UTC for DB timestamps
            now_local = datetime.now(BP_TZ)               # Budapest for cron matching

            # Daily cleanup: delete uploads older than 7 days (runs once at ~04:00 UTC)
            today_str = now_dt.strftime("%Y-%m-%d")
            if now_dt.hour == 4 and _last_cleanup != today_str:
                _last_cleanup = today_str
                try:
                    cutoff = (now_dt - timedelta(days=7)).isoformat()
                    conn = get_db()
                    deleted = conn.execute("DELETE FROM uploads WHERE uploaded_at < ?", (cutoff,)).rowcount
                    conn.commit()
                    conn.close()
                    if deleted:
                        logger.info("Uploads cleanup: %d files older than 7 days deleted", deleted)
                except Exception as e:
                    logger.error("Uploads cleanup failed: %s", e)
            conn = get_db()
            recipes = conn.execute(
                "SELECT id, name, cron_schedule, cron_model, cron_delivery, cron_last_run "
                "FROM pyramid_recipes WHERE cron_enabled = 1 AND cron_schedule IS NOT NULL AND enabled = 1"
            ).fetchall()
            conn.close()

            for r in recipes:
                # Match cron against Budapest LOCAL time (user's timezone).
                if not _cron_matches(r["cron_schedule"], now_local):
                    continue
                # Dedup: skip if already ran this minute (compare in UTC)
                if r["cron_last_run"]:
                    try:
                        last = datetime.fromisoformat(r["cron_last_run"])
                        if (last.year == now_dt.year and last.month == now_dt.month and
                                last.day == now_dt.day and last.hour == now_dt.hour and
                                last.minute == now_dt.minute):
                            continue
                    except (ValueError, TypeError):
                        pass

                logger.info("Cron trigger: %s (schedule=%s, model=%s)", r["name"], r["cron_schedule"], r["cron_model"])

                # Mark last_run immediately to prevent double fire
                conn = get_db()
                conn.execute("UPDATE pyramid_recipes SET cron_last_run = ? WHERE id = ?", (now_dt.isoformat(), r["id"]))
                conn.commit()
                conn.close()

                # Execute via ai_task dispatch (same path as execute_recipe)
                try:
                    recipe_conn = get_db()
                    recipe_row = recipe_conn.execute(
                        "SELECT prompt_template, required_tools FROM pyramid_recipes WHERE id = ?", (r["id"],)
                    ).fetchone()
                    recipe_conn.close()

                    if not recipe_row:
                        continue

                    prompt = recipe_row["prompt_template"]
                    req_tools = json.loads(recipe_row["required_tools"]) if recipe_row["required_tools"] else []
                    if req_tools:
                        prompt += f"\n\nELERHETO TOOL-OK: {', '.join(req_tools)}"
                    today = now_dt.strftime("%Y-%m-%d")
                    prompt += f"\n\n[Mai datum: {today}. Az adatoknak FRISSNEK kell lenniuk!]"

                    # ── Operation Kabare: prefetch real data (same as execute_recipe) ──
                    # A cron path korabban kikerulte a plugins/recipes.py execute_recipe-t,
                    # es igy a FACTUAL CONTEXT prefetch sem futott. Ezert a cron-bol futo
                    # recipe-k halucinaltak arfolyamot. Itt inline injektaljuk.
                    try:
                        from plugins._recipe_prefetch import run_prefetch
                        prefetch_deps = {
                            "get_db": get_db,
                            "capture_state": _capture_state,
                        }
                        factual_context = await run_prefetch(r["name"], prefetch_deps)
                        if factual_context:
                            prompt += (
                                "\n\n=== FACTUAL CONTEXT (Python-ban lehuzott valos adatok) ===\n"
                                f"{factual_context}\n"
                                "=== END FACTUAL CONTEXT ===\n\n"
                                "SZIGORU SZABALY: MINDEN szamadatnak (arfolyam, ar, index, idopont, "
                                "nev, cim) a fenti FACTUAL CONTEXT blokkbol kell szarmaznia. "
                                "TILOS fejbol szamot, adatot, forrast irni. Ha valami nincs a "
                                "CONTEXT-ben, ird: 'adat nem elerheto'. SOHA ne talalj ki semmit."
                            )
                            logger.info("Cron prefetch injected for: %s (%d chars)",
                                        r["name"], len(factual_context))
                    except Exception as e:
                        logger.error("Cron prefetch injection failed for %s: %s", r["name"], e)

                    model = r["cron_model"] or "glm5"
                    if model == "all":
                        result_json = await ai_task(
                            title=f"Cron: {r['name']}",
                            description=prompt,
                            assigned_by="cron-scheduler",
                        )
                    else:
                        max_tokens = 16000 if model == "glm5" else 8000
                        agent_tasks = json.dumps({model: {"prompt": prompt, "max_tokens": max_tokens}})
                        result_json = await ai_task(
                            title=f"Cron: {r['name']}",
                            description=prompt,
                            assigned_by="cron-scheduler",
                            agent_tasks=agent_tasks,
                        )

                    result = json.loads(result_json)
                    task_id = result.get("task_id")

                    # Poll for result (delivery needs actual content)
                    if task_id and r["cron_delivery"] in ("telegram", "both"):
                        for _ in range(90):
                            await asyncio.sleep(2)
                            conn = get_db()
                            row2 = conn.execute(
                                "SELECT content FROM ai_task_results WHERE task_id = ? ORDER BY id DESC LIMIT 1",
                                (task_id,)
                            ).fetchone()
                            status = conn.execute("SELECT status FROM ai_tasks WHERE id = ?", (task_id,)).fetchone()
                            conn.close()
                            if row2:
                                content = row2["content"] or ""
                                msg = f"📋 <b>{r['name']}</b> (scheduled)\n\n{content[:3800]}"
                                await _telegram_push(msg)
                                logger.info("Cron delivered: %s → Telegram (task #%d)", r["name"], task_id)
                                break
                            if status and status["status"] in ("completed", "failed"):
                                break

                except Exception as e:
                    logger.error("Cron execution failed for %s: %s", r["name"], e)

        except Exception as e:
            logger.error("Cron loop error: %s", e)

        await asyncio.sleep(60)


def _start_cron_scheduler():
    """Start cron loop in a background thread with its own event loop."""
    def _run():
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_cron_loop())

    t = threading.Thread(target=_run, daemon=True)
    t.start()
    logger.info("Cron scheduler background thread started")

_start_cron_scheduler()


if __name__ == "__main__":
    mcp.run(
        transport="streamable-http",
        host="0.0.0.0",
        port=int(os.environ.get("PORT", 8003))
    )
