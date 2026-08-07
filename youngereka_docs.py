"""
Dokumentum-pipeline — feltöltésből vision-kész anyag.
======================================================

Ez a felület lényegi része. Ha csak a bodyszöveg-kinyerés készülne el,
a felület fele annyit érne: egy biológiai szakcikkben az információ fele
a Figure 3-ban van, nem a bekezdésekben.

    upload → MIME-sniff
    ├── PDF   → szövegréteg + (szkennelt? raszterizálás) + ÁBRA-kinyerés
    ├── kép   → normalizálás (HEIC/MPO/EXIF-forgatás/GPS-strip)
    ├── office→ docx/xlsx/pptx → markdown
    └── text  → közvetlen

A KÉP-NORMALIZÁLÓ A KÖZÖS (task #20)
------------------------------------
A #20-as task (Agora iPhone fotófeltöltés) ugyanezt a pipeline-t kéri:
`register_heif_opener` · MPO → `seek(0)` · `ImageOps.exif_transpose` ·
`convert("RGB")` · hosszabb él max 2048px · JPEG q85 · TELJES EXIF strip.
A #20 még pending, tehát ez itt AZ a közös normalizáló — ott is ezt kell
majd hívni, nem újat írni.

Az EXIF-strip nem kozmetika: egy telefonos fotó GPS-koordinátát hordoz.
Réka labor- vagy otthoni fotóiról a helyadat nem megy fel egy külső
API-ra. A `convert("RGB")` + újrakódolás mellékhatásként mindent eldob,
de expliciten is kimondjuk, mert erre nem szabad mellékhatásként hagyatkozni.
"""
from __future__ import annotations

import base64
import io
import logging
import re

logger = logging.getLogger("bridge.yr_docs")

#: Egy üzenetben ennyi kép mehet fel. A SiliconFlow több modellje 4 kép
#: fölött HTTP 20015-öt dob (mérés: 2026-06, gemma-4 / Qwen3.6). A
#: munkaparancs 8-at ír; élesben a kisebbik a biztonságos, ezért a
#: tényleges korlátot mérés dönti el, és env-ből állítható.
VISION_MAX_IMAGES = 8

MAX_EDGE = 2048
JPEG_QUALITY = 85

#: 30k token ~ 120k karakter (magyar/angol vegyes szövegre ~4 kar/token).
MAX_TEXT_CHARS = 120_000

_HEIF_READY = False


def _ensure_heif() -> None:
    """A HEIC-olvasót egyszer regisztráljuk. Hiánya nem végzetes:
    a nem-HEIC képek ugyanúgy mennek tovább."""
    global _HEIF_READY
    if _HEIF_READY:
        return
    try:
        from pillow_heif import register_heif_opener
        register_heif_opener()
        _HEIF_READY = True
    except Exception as e:  # noqa: BLE001
        logger.warning("pillow_heif nem elérhető — a HEIC feltöltés el fog "
                       "bukni, minden más megy: %s", e)


def normalize_image(raw: bytes) -> bytes:
    """Tetszőleges fotó → tiszta, álló helyzetű, EXIF-mentes JPEG.

    Ez a közös normalizáló (#20 is ezt hívja). Dob, ha a Pillow nem
    tudja megnyitni — a hívó ebből tudja, hogy a fájl nem kép.

    A validáció elve „a Pillow megnyitja és újrakódolja", NEM
    format-whitelist: az iPhone MPO-t (multi-frame JPEG) a whitelist
    elutasítaná, pedig az első frame maga a fotó.
    """
    _ensure_heif()
    from PIL import Image, ImageOps

    im = Image.open(io.BytesIO(raw))

    # MPO (iPhone): több frame, az első a fotó
    if getattr(im, "format", "") == "MPO":
        try:
            im.seek(0)
        except Exception:  # noqa: BLE001
            pass

    # EXIF orientation → tényleges pixelforgatás. Enélkül a telefonos
    # fotó fekve megy a modellhez, és a leírás is fekve lesz.
    im = ImageOps.exif_transpose(im)

    if im.mode != "RGB":
        im = im.convert("RGB")

    if max(im.size) > MAX_EDGE:
        im.thumbnail((MAX_EDGE, MAX_EDGE), Image.LANCZOS)

    out = io.BytesIO()
    # Se exif=, se icc_profile= — így semmi metaadat (GPS!) nem megy tovább.
    im.save(out, format="JPEG", quality=JPEG_QUALITY, optimize=True)
    return out.getvalue()


def _b64(raw: bytes) -> str:
    return base64.b64encode(raw).decode("ascii")


def _data_url(jpeg: bytes) -> str:
    return f"data:image/jpeg;base64,{_b64(jpeg)}"


# ============================================================
# PDF
# ============================================================

def _pdf(raw: bytes, max_images: int) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=raw, filetype="pdf")
    pages = doc.page_count

    parts, figures = [], []
    total_chars = 0
    for pno in range(pages):
        page = doc.load_page(pno)
        t = page.get_text() or ""
        total_chars += len(t)
        if t.strip():
            parts.append(t)

    # Szkennelt-e? Kevés szöveg oldalanként → nincs valódi szövegréteg.
    per_page = (total_chars / pages) if pages else 0
    scanned = per_page < 100

    notes = []
    if scanned:
        # Raszterizálás: az oldal MAGA a kép, azt olvassa a vision.
        for pno in range(min(pages, max_images)):
            pix = doc.load_page(pno).get_pixmap(dpi=150)
            figures.append(normalize_image(pix.tobytes("png")))
        notes.append(
            f"Ebben a PDF-ben nincs szövegréteg (szkennelt vagy fotózott), "
            f"ezért képként olvasom: {len(figures)} oldal ment fel"
            + (f" a {pages}-ből." if pages > len(figures) else "."))
    else:
        # Beágyazott ÁBRÁK kiemelése. Ezek a Figure-ök.
        seen = set()
        for pno in range(pages):
            page = doc.load_page(pno)
            page_area = abs(page.rect.width * page.rect.height) or 1
            for info in page.get_images(full=True):
                xref = info[0]
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    img = doc.extract_image(xref)
                except Exception:  # noqa: BLE001
                    continue
                w, h = img.get("width", 0), img.get("height", 0)
                if w <= 200 or h <= 200:
                    continue  # ikon, logó, elválasztó
                # Teljes oldalt lefedő kép = oldalszkennelés, nem ábra
                try:
                    rects = page.get_image_rects(xref)
                    if rects and abs(rects[0].width * rects[0].height) > 0.9 * page_area:
                        continue
                except Exception:  # noqa: BLE001
                    pass
                try:
                    figures.append(normalize_image(img["image"]))
                except Exception:  # noqa: BLE001
                    continue

    doc.close()

    found = len(figures)
    if not scanned and found > max_images:
        notes.append(
            f"{found} ábrát találtam, az első {max_images} ment fel. "
            f"Ha egy másik kell, írd meg a számát — pl. „a 9. ábra”.")
        figures = figures[:max_images]

    text = "\n\n".join(parts)
    truncated = False
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS]
        truncated = True
        notes.append("A dokumentum hosszabb, mint amit egyben fel tudok "
                     "dolgozni — az elejét olvastam. Ha a vége kell, szólj, "
                     "és arra a részre koncentrálok.")

    label_bits = [f"{pages} oldal"]
    if scanned:
        label_bits.append(f"{len(figures)} oldalkép")
    elif found:
        label_bits.append(f"{min(found, max_images)} ábra kinyerve"
                          + (f" / {found}" if found > max_images else ""))
    if truncated:
        label_bits.append("szöveg vágva")

    return {
        "kind": "pdf",
        "label": "PDF · " + " · ".join(label_bits),
        "text": text,
        "images": figures,
        "images_found": found,
        "notes": notes,
    }


# ============================================================
# Office → markdown
# ============================================================

def _docx(raw: bytes) -> dict:
    import docx  # python-docx
    d = docx.Document(io.BytesIO(raw))
    lines = [p.text for p in d.paragraphs if p.text.strip()]
    for ti, table in enumerate(d.tables, 1):
        lines.append(f"\n**{ti}. táblázat**\n")
        for row in table.rows:
            lines.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return {"kind": "office", "label": f"DOCX · {len(d.paragraphs)} bekezdés"
            + (f" · {len(d.tables)} táblázat" if d.tables else ""),
            "text": "\n".join(lines), "images": [], "images_found": 0, "notes": []}


def _xlsx(raw: bytes) -> dict:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
    out, rows_total = [], 0
    for ws in wb.worksheets:
        out.append(f"\n## {ws.title}\n")
        for row in ws.iter_rows(values_only=True):
            if row is None or all(c is None for c in row):
                continue
            rows_total += 1
            out.append("| " + " | ".join("" if c is None else str(c) for c in row) + " |")
    return {"kind": "office",
            "label": f"XLSX · {len(wb.worksheets)} munkalap · {rows_total} sor",
            "text": "\n".join(out), "images": [], "images_found": 0, "notes": []}


def _pptx(raw: bytes) -> dict:
    from pptx import Presentation
    pres = Presentation(io.BytesIO(raw))
    out = []
    for i, slide in enumerate(pres.slides, 1):
        out.append(f"\n## {i}. dia\n")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return {"kind": "office", "label": f"PPTX · {len(pres.slides)} dia",
            "text": "\n".join(out), "images": [], "images_found": 0, "notes": []}


# ============================================================
# Belépési pont
# ============================================================

_OLVASHATATLAN = (
    "Ezt a fájlt nem tudom elolvasni — nem kép, nem PDF, nem Office és nem "
    "szöveg. PDF-ként vagy képként (akár telefonnal lefotózva) menni fog.")


def _szovegnek_tunik(text: str, minta: int = 4000) -> bool:
    """Igaz, ha a dekódolt tartalom tényleg olvasható szöveg.

    Az egybájtos kodekek minden bájtra adnak VALAMILYEN karaktert, tehát a
    sikeres `decode()` önmagában semmit nem bizonyít. A vezérlőkarakterek
    aránya viszont elárulja a bináris tartalmat: valódi szövegben a
    sortörésen és taboláson kívül alig van ilyen.
    """
    minta_szoveg = text[:minta]
    if not minta_szoveg.strip():
        return False
    vezerlo = sum(1 for c in minta_szoveg
                  if (ord(c) < 32 and c not in "\n\r\t") or ord(c) == 127)
    return vezerlo / len(minta_szoveg) < 0.02


_IMAGE_EXT = {"jpg", "jpeg", "png", "gif", "webp", "heic", "heif", "bmp", "tif", "tiff", "mpo"}
_TEXT_EXT = {"txt", "md", "csv", "tsv", "json", "log", "fasta", "fa", "gb", "r", "py"}


def _sniff(raw: bytes, filename: str) -> str:
    """Tartalom szerint, nem kiterjesztés szerint — a telefonok
    következetlenül neveznek (`image.jpg`, ami valójában HEIC)."""
    if raw[:4] == b"%PDF":
        return "pdf"
    if raw[:2] == b"PK":  # zip-alapú: docx/xlsx/pptx
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        return ext if ext in {"docx", "xlsx", "pptx"} else "zip"
    if raw[:3] == b"\xff\xd8\xff" or raw[:8] == b"\x89PNG\r\n\x1a\n" or raw[:6] in (b"GIF87a", b"GIF89a"):
        return "image"
    if raw[4:12] in (b"ftypheic", b"ftypheix", b"ftyphevc", b"ftypmif1", b"ftypmsf1"):
        return "image"
    if raw[:4] == b"RIFF" and raw[8:12] == b"WEBP":
        return "image"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in _IMAGE_EXT:
        return "image"
    if ext in {"docx", "xlsx", "pptx"}:
        return ext
    return "text"


def process_upload(raw: bytes, filename: str, max_images: int = VISION_MAX_IMAGES) -> dict:
    """Feltöltött fájl → {kind, label, text, images(JPEG bytes), notes}.

    A `label` a felület szignatúra-eleme: a preparátum-címke, ami
    pontosan megmondja Rékának, MIT lát a modell a fájlból. Ezért nem
    díszítés, hanem a legfontosabb információdarab — és ezért nem
    általánosít („feldolgozva"), hanem számol.
    """
    kind = _sniff(raw, filename)

    if kind == "pdf":
        return _pdf(raw, max_images)

    if kind == "image":
        jpeg = normalize_image(raw)
        from PIL import Image
        w, h = Image.open(io.BytesIO(jpeg)).size
        return {"kind": "image", "label": f"Kép · {w}×{h}",
                "text": "", "images": [jpeg], "images_found": 1, "notes": []}

    if kind == "docx":
        return _docx(raw)
    if kind == "xlsx":
        return _xlsx(raw)
    if kind == "pptx":
        return _pptx(raw)

    # text / csv / md / ismeretlen
    #
    # FIGYELEM: az iso-8859-2 és a cp1250 egybájtos kódolás, tehát
    # GYAKORLATILAG BÁRMIT dekódol — egy videó vagy egy zip is „sikeresen"
    # szöveggé válna, és halandzsa menne fel a modellhez. A csendes siker
    # itt rosszabb a hibánál, ezért a dekódolás után józansági próba van.
    if b"\x00" in raw[:8192]:
        raise ValueError(_OLVASHATATLAN)

    text = ""
    for enc in ("utf-8", "utf-8-sig", "iso-8859-2", "cp1250"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError(_OLVASHATATLAN)

    if not _szovegnek_tunik(text):
        raise ValueError(_OLVASHATATLAN)

    notes = []
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS]
        notes.append("Hosszú fájl — az elejét olvastam.")
    lines = text.count("\n") + 1
    return {"kind": "text", "label": f"Szöveg · {lines} sor · {len(text)} karakter",
            "text": text, "images": [], "images_found": 0, "notes": notes}


def vision_content(text: str, images: list[bytes]) -> list | str:
    """OpenAI-kompatibilis multimodális `content`. Kép nélkül sima string,
    mert néhány modell a fölösleges listát rosszul tűri."""
    if not images:
        return text
    out: list = [{"type": "text", "text": text}]
    for jpeg in images:
        out.append({"type": "image_url", "image_url": {"url": _data_url(jpeg)}})
    return out


def summarize_for_prompt(doc: dict, filename: str) -> str:
    """A kinyert anyag szöveges része, a modellnek szánt kerettel."""
    head = f"[Melléklet: {filename} — {doc['label']}]"
    body = (doc.get("text") or "").strip()
    if body:
        return f"{head}\n\n{body}"
    if doc.get("images"):
        return f"{head}\n(A tartalom képként megy fel.)"
    return head


_WS = re.compile(r"\s+")


def title_from(text: str, limit: int = 60) -> str:
    """Beszélgetés-cím az első üzenetből."""
    t = _WS.sub(" ", (text or "").strip())
    if not t:
        return "Új beszélgetés"
    return t[:limit] + ("…" if len(t) > limit else "")
